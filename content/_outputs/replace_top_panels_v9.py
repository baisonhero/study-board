#!/usr/bin/env python3
"""
V11 rev9 — in-place replacement of the 3 TOP panels only.
Preserves all user manual edits elsewhere on the slide.

Steps:
  1. Open existing V11 pptx
  2. Find every shape whose top is in band [6.45, 7.95] inches
     (= the current 3 TOP panels area). Remove those shapes only.
  3. Insert three new equal-width unified TOP5 panels in the same band:
       LEFT  : メンバー 6月コスト TOP5   ([千円])
       CENTER: メンバー 6月 生成行 TOP5  ([行])
       RIGHT : メンバー 6月 効率 TOP5    ([行/千円])
  4. Save in place.

Data integrity checks (asserted before save):
  田中: 2200 ÷ 78 = 28.2  (28.205…)
  高橋: 1750 ÷ 64 = 27.3  (27.344…)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V11.pptx"

# ---------- House palette (V11) ----------
DEEP_FOREST = RGBColor(0x1F, 0x3A, 0x2E)
SAGE        = RGBColor(0xB5, 0xC9, 0xA8)
SAGE_BG     = RGBColor(0xE8, 0xEF, 0xE3)
AMBER       = RGBColor(0xF4, 0xB9, 0x42)
AMBER_DEEP  = RGBColor(0xC8, 0x8E, 0x1F)
INK         = RGBColor(0x22, 0x28, 0x24)
MUTED       = RGBColor(0x6B, 0x6B, 0x6B)
CARD_BORDER = RGBColor(0xD7, 0xDC, 0xCF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

JP_FONT, EN_FONT = "メイリオ", "Segoe UI"

def set_run_font(run, latin=EN_FONT, ea=JP_FONT, size=12, color=INK, bold=False):
    run.font.name = latin
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea_elem = etree.SubElement(rPr, qn('a:ea'))
    ea_elem.set('typeface', ea)

def add_textbox(slide, text, l, t, w, h, size=12, color=INK, bold=False,
                align='left', valign='top', word_wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if valign == 'middle':
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == 'bottom':
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    set_run_font(r, size=size, color=color, bold=bold)
    return tb

def add_rect(slide, l, t, w, h, fill=WHITE, line=None, line_width=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        if line_width:
            shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_oval(slide, l, t, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

# ---------- Data integrity ----------
USERS = [
    # (name, cost千円, lines行)
    ("山田 太郎", 130, 4200),
    ("鈴木 花子", 108, 4500),
    ("佐藤 健",    92, 2800),
    ("田中 美咲",  78, 2200),    # mock added
    ("高橋 翔",    64, 1750),    # mock added
]
for name, c, ln in USERS:
    eff = round(ln / c, 1)
    print(f'  {name}: {ln} 行 / {c} 千円 = {eff} 行/千円')
# Assert mock math
assert round(2200/78, 1) == 28.2
assert round(1750/64, 1) == 27.3

top_cost = sorted(USERS, key=lambda u: -u[1])         # by cost desc
top_loc  = sorted(USERS, key=lambda u: -u[2])         # by lines desc
top_eff  = sorted(USERS, key=lambda u: -round(u[2]/u[1], 1))  # by efficiency desc

# ---------- Open and inspect ----------
prs = Presentation(PPTX)
slide = prs.slides[0]

# Identify shapes in the TOP panels Y band [6.45, 7.95]
BAND_TOP, BAND_BOT = 6.45, 7.95
to_remove = []
for sh in list(slide.shapes):
    try:
        t_in = sh.top / 914400
        if BAND_TOP <= t_in < BAND_BOT:
            to_remove.append(sh)
    except Exception:
        pass

print(f'\nRemoving {len(to_remove)} shapes in band {BAND_TOP}–{BAND_BOT}"...')
for sh in to_remove:
    sp = sh._element
    sp.getparent().remove(sp)

# ---------- Add new unified TOP5 panels ----------
TP_T   = 6.50
TP_H   = 1.40
PAGE_L = 0.30
PAGE_W = 7.67
n_panels = 3
gap     = 0.10
panel_w = (PAGE_W - (n_panels - 1) * gap) / n_panels   # = (7.67-0.20)/3 = 2.49

# Rank badge colors
RANK_BG = [AMBER, SAGE, RGBColor(0xE0, 0xE0, 0xE0), SAGE_BG, SAGE_BG]

panels = [
    # (title, unit_label, rankings, value_formatter)
    ("6月 コスト TOP5", "[千円]",      top_cost,
        lambda u: f"{u[1]}"),
    ("6月 生成行 TOP5", "[行]",        top_loc,
        lambda u: f"{u[2]:,}"),
    ("6月 効率 TOP5",   "[行/千円]",   top_eff,
        lambda u: f"{round(u[2]/u[1], 1)}"),
]

for col_i, (title, unit, ranking, value_fn) in enumerate(panels):
    px = PAGE_L + col_i * (panel_w + gap)

    # Outer card
    add_rect(slide, px, TP_T, panel_w, TP_H,
             fill=WHITE, line=CARD_BORDER, line_width=0.5)
    # Top header strip
    add_rect(slide, px, TP_T, panel_w, 0.04, fill=DEEP_FOREST)

    # Title + unit pill (same row, two halves)
    add_textbox(slide, title,
                l=px+0.08, t=TP_T+0.07, w=panel_w-0.80, h=0.18,
                size=9.5, color=DEEP_FOREST, bold=True,
                valign='middle', word_wrap=False)
    add_textbox(slide, unit,
                l=px+panel_w-0.80, t=TP_T+0.07, w=0.70, h=0.18,
                size=8, color=MUTED, align='right',
                valign='middle', word_wrap=False)

    # 5 rows
    row_top = TP_T + 0.32
    row_h   = 0.20
    inner_pad = 0.08
    rank_w  = 0.18
    name_w  = 1.12
    value_w = 0.62
    for ri, member in enumerate(ranking):
        ry = row_top + ri * row_h
        rx = px + inner_pad
        # Rank circle
        add_oval(slide, rx, ry + 0.02, rank_w, rank_w,
                 fill=RANK_BG[ri], line=None)
        add_textbox(slide, str(ri+1),
                    l=rx, t=ry+0.02, w=rank_w, h=rank_w,
                    size=9, color=DEEP_FOREST, bold=True,
                    align='center', valign='middle')
        # Name
        add_textbox(slide, member[0],
                    l=rx + rank_w + 0.05, t=ry, w=name_w, h=row_h,
                    size=10, color=INK, bold=(ri == 0), valign='middle',
                    word_wrap=False)
        # Value (right-aligned)
        add_textbox(slide, value_fn(member),
                    l=px + panel_w - inner_pad - value_w, t=ry,
                    w=value_w, h=row_h,
                    size=11, color=AMBER_DEEP, bold=True,
                    align='right', valign='middle', word_wrap=False)

# Save in place
prs.save(PPTX)
print(f'\nWROTE: {PPTX}')
