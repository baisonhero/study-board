#!/usr/bin/env python3
"""
V13 rev4 — left ref panel rewrite per user mockup.

Target layout (DEEP_FOREST background, white/amber text):
  上段:
    '当チーム実績 (6月)'          ← small bold label (white)
    '690 [千円] → 20,000行'       ← LARGE highlighted (amber accent)
  ─── Amber divider line ───
  下段:
    '外注費用換算 (1,600 千円/人月)' ← small bold label (white)
    '約 46,300 行 / 人月'           ← LARGE highlighted (amber accent)

Calculation: 20,000 行 × 1,600 千円/人月 ÷ 690 千円 ≈ 46,377 → 約 46,300

Current shapes (from inspection):
  [83] AUTO_SHAPE background panel (kept)
  [84] '当チーム実績 (5月)'                     → '当チーム実績 (6月)'
  [85] '620 [千円] → 17,000行'                  → '690 [千円] → 20,000行'
  [86] LINE divider (kept)
  [87] '業界基準 (新規開発上位25% 1,000行/人月)' → '外注費用換算 (1,600 千円/人月)'
  [88] '17,000 行 ≒ 17 人月相当'                → '約 46,300 行 / 人月'
  [89] '(≒ 外注費用 27,200 [千円] 相当)'        → DELETE
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V13.pptx"

AMBER       = RGBColor(0xF4, 0xB9, 0x42)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
JP_FONT, EN_FONT = "メイリオ", "Segoe UI"

def set_ea_font(run):
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', JP_FONT)

def replace_text(shape, new_text, size=None, color=None, bold=None, align=None):
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        return
    p0.runs[0].text = new_text
    for r in list(p0.runs)[1:]:
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    r0 = p0.runs[0]
    if size is not None:
        r0.font.size = Pt(size)
    if color is not None:
        r0.font.color.rgb = color
    if bold is not None:
        r0.font.bold = bold
    if align is not None:
        p0.alignment = align
    r0.font.name = EN_FONT
    set_ea_font(r0)

prs = Presentation(PPTX)
slide = prs.slides[0]

# Build a map of label-shapes by their current text inside the left ref band.
TARGETS = {
    "当チーム実績 (5月)":                         ("当チーム実績 (6月)",        9,  WHITE, True,  PP_ALIGN.LEFT),
    "620 [千円] → 17,000行":                      ("690 [千円] → 20,000行",      18, AMBER, True,  PP_ALIGN.LEFT),
    "業界基準 (新規開発上位25% 1,000行/人月)":     ("外注費用換算 (1,600 千円/人月)", 9,  WHITE, True,  PP_ALIGN.LEFT),
    "17,000 行 ≒ 17 人月相当":                    ("約 46,300 行 / 人月",        18, AMBER, True,  PP_ALIGN.LEFT),
}
DELETE_TEXTS = {"(≒ 外注費用 27,200 [千円] 相当)"}

to_delete = []
hits = 0
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    try:
        t_in = sh.top/914400; l_in = sh.left/914400
    except:
        continue
    if not (10.30 <= t_in <= 11.80 and l_in < 4.10):
        continue
    txt = sh.text_frame.text.strip()
    if txt in TARGETS:
        new_text, size, color, bold, align = TARGETS[txt]
        replace_text(sh, new_text, size=size, color=color, bold=bold, align=align)
        # bump heights for large highlighted values
        if size >= 16:
            sh.height = Inches(0.34)
        print(f"  ✓ '{txt}' → '{new_text}' ({size}pt)")
        hits += 1
    elif txt in DELETE_TEXTS:
        to_delete.append((sh, txt))

for sh, txt in to_delete:
    sh._element.getparent().remove(sh._element)
    print(f"  ✗ Deleted: '{txt}'")
    hits += 1

assert hits >= 5, f"Expected 5+ edits in left ref panel, got {hits}"

prs.save(PPTX)
print(f"\nWROTE: {PPTX}")
