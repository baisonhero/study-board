#!/usr/bin/env python3
"""
V11 rev16 — in-place: update LEFT reference band (当チーム実績).

Calculation basis (per user choice — 新規開発上位25% = 1,000 行/人月):
  20,000 行 ÷ 1,000 行/人月 = 20 人月
  20 人月 × 1,600 千円/人月 = 32,000 千円

3-line update (preserves shape positions/styles):
  [137] sub label  (10.92, 9pt white bold)
          "約 1,260 [千円] 投資換算 (人月単純相当)"
        → "外注費用換算 (1,600 千円/人月)"
  [138] big amber  (11.10, 18pt amber bold)
          "約 36,400 行 / 人/月"
        → "約 32,000 [千円] 相当"
  [139] caption    (11.40, 8pt sage)
          "業界中央値 600 行/人/月 の 約61倍 (※単純比較)"
        → "20人月 × 1,600 千円/人月 (新規開発上位25% 1,000行/人月 基準)"
"""
from pptx import Presentation
PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V11.pptx"

OLD_NEW = {
    "約 1,260 [千円] 投資換算 (人月単純相当)":
        "外注費用換算 (1,600 千円/人月)",
    "約 36,400 行 / 人/月":
        "約 32,000 [千円] 相当",
    "業界中央値 600 行/人/月 の 約61倍 (※単純比較)":
        "20人月 × 1,600 千円/人月 (新規開発上位25% 1,000行/人月 基準)",
}

def replace_text_keep_props(shape, new_text):
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        return False
    p0.runs[0].text = new_text
    for run in list(p0.runs)[1:]:
        run._r.getparent().remove(run._r)
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    return True

prs = Presentation(PPTX)
slide = prs.slides[0]

n = 0
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    try:
        t_in = sh.top/914400; l_in = sh.left/914400
    except Exception:
        continue
    # Target only LEFT reference band (L < 4.05, Y >= 10.30)
    if not (t_in >= 10.30 and l_in < 4.05):
        continue
    text = sh.text_frame.text.strip()
    if text in OLD_NEW:
        new = OLD_NEW[text]
        if replace_text_keep_props(sh, new):
            n += 1
            print(f"  ✓ {text!r}\n      → {new!r}")

assert n == 3, f"Expected 3 replacements, got {n}"
prs.save(PPTX)
print(f"\nReplaced {n} text shapes.\nWROTE: {PPTX}")
