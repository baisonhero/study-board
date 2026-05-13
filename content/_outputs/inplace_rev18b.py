#!/usr/bin/env python3
"""
V11 rev18b — fix the diagonal separator line in the plan/actual table.

Previous rev18 incorrectly resized the vertical separator (LINE shape) to
W=0.91, turning it diagonal. Reset it to a vertical line at L=3.68 (between
6月 and 7月) with W=0 and the same height.
"""
from pptx import Presentation
from pptx.util import Inches

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V11.pptx"
prs = Presentation(PPTX)
slide = prs.slides[0]

n = 0
for sh in slide.shapes:
    # find LINE shape (type 9) in the table band Y ~ 4.40-4.95
    if sh.shape_type != 9:
        continue
    try:
        t_in = sh.top/914400
    except Exception:
        continue
    if not (4.30 <= t_in < 5.00):
        continue
    # Reset to vertical line at L=3.68, W=0
    # End point Y = T + H, currently we want H ≈ 0.48 (spans table header + 2 rows)
    sh.left = Inches(3.68)
    sh.width = Inches(0.0)
    # Keep height (it was 0.48 originally per inspection)
    # Some line shapes have negative-only-by-direction; set height if needed
    if sh.height == 0:
        sh.height = Inches(0.48)
    n += 1
    print(f"  ✓ Separator reset: vertical line at L=3.68, W=0, H={sh.height/914400:.2f}")

assert n >= 1, "Separator line not found"
prs.save(PPTX)
print(f"WROTE: {PPTX}")
