#!/usr/bin/env python3
"""Targeted patch: top-right LEFT/RIGHT cards.
LEFT card (5.00, 1.00, 1.49, 1.20): keep CORAL_DARK fill, change ALL text → WHITE
RIGHT card (6.48, 1.00, 1.49, 1.20): change fill CORAL → SAND, text already CORAL (correct)
"""
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
import os

SRC = '2026-07-01_Claude_monthly_report_V13_tmp.pptx'
SAND  = RGBColor(0xE8, 0xDC, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CORAL = RGBColor(0xCC, 0x78, 0x5C)
CORAL_DARK = RGBColor(0xA8, 0x5F, 0x4A)
ANTHRACITE = RGBColor(0x19, 0x19, 0x19)

def in_rect(sh, l_min, l_max, t_min, t_max):
    l = Emu(sh.left).inches if sh.left else 0
    t = Emu(sh.top).inches if sh.top else 0
    return l_min <= l <= l_max and t_min <= t <= t_max

def main():
    os.chdir(os.path.dirname(os.path.abspath(__file__)))
    prs = Presentation(SRC)
    patched = []
    for slide in prs.slides:
        for sh in slide.shapes:
            # LEFT card text shapes (5.0–5.05 wide, t between 1.0–2.2)
            if in_rect(sh, 4.99, 5.06, 1.0, 2.2) and sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = WHITE
                        patched.append(f'LEFT-text({sh.left/914400:.2f},{sh.top/914400:.2f}) → WHITE: {r.text!r}')
            # RIGHT card fill shape: AUTO_SHAPE at (6.48, 1.00, 1.49, 1.20) — area >= 1.0
            if in_rect(sh, 6.47, 6.50, 0.99, 1.01) and not sh.has_text_frame:
                # large rect fill
                try:
                    sh.fill.solid(); sh.fill.fore_color.rgb = SAND
                    patched.append('RIGHT-fill → SAND')
                except Exception as e:
                    print('fill error:', e)
            # Also catch the rect if it has text_frame (auto_shape can have one)
            if in_rect(sh, 6.47, 6.50, 0.99, 1.01):
                w = Emu(sh.width).inches if sh.width else 0
                h = Emu(sh.height).inches if sh.height else 0
                if w > 1.0 and h > 1.0:
                    try:
                        sh.fill.solid(); sh.fill.fore_color.rgb = SAND
                        patched.append(f'RIGHT-bigrect ({w:.2f}x{h:.2f}) → SAND fill')
                    except Exception as e:
                        print('fill2 error:', e)
    # RIGHT card subtext (外注費用 1.2人月, (1,600 千円/月 換算)) → ANTHRACITE
    # Those text shapes sit at l~6.49, t in 1.82–2.20
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            l = Emu(sh.left).inches if sh.left else 0
            t = Emu(sh.top).inches if sh.top else 0
            if 6.45 <= l <= 6.55 and 1.80 <= t <= 2.20:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = ANTHRACITE
                        patched.append(f'RIGHT-subtext({l:.2f},{t:.2f}) → ANTHRACITE: {r.text!r}')
    prs.save(SRC)
    for p in patched:
        print(p)

if __name__ == '__main__':
    main()
