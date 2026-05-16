#!/usr/bin/env python3
"""
V13 rev7 — ポジショニング X軸の最大値「2000」を plot 右端に正確配置 + 軸タイトル
「行数 [行/人]」を「2000」の右側に重ならず配置.

Plot geometry (from inspection):
  X-axis baseline: L=1.20, W=5.00 → spans L=1.20 (=0) to L=6.20 (=2000)
  Midpoint = 3.70 (=1000), also has vertical divider at L=3.70

Current X tick positions (LEFT-aligned text in W=0.40 boxes):
  "0"    L=1.00 → visually near plot left ✓
  "1000" L=3.08 → MISALIGNED (should center at 3.70)
  "2000" L=5.17 → MISALIGNED (should center at 6.20)

Current axis title:
  '行数 [行/人]' at L=5.70, W=1.10 → overlaps with tick area, misplaced

Fix:
  Set CENTER alignment on tick text + reposition box so visual center aligns with axis points.
  Place axis title to the right of the new "2000" tick box (which ends at L=6.40).
"""
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V13.pptx"

prs = Presentation(PPTX)
slide = prs.slides[0]

# Target positions for X-axis tick labels (W=0.40, CENTER aligned)
# Plot: L_left=1.20, L_right=6.20, mid=3.70
TICK_POS = {
    "0":    (1.00, 9.71, 0.40, 0.14),  # center at L=1.20
    "1000": (3.50, 9.71, 0.40, 0.14),  # center at L=3.70
    "2000": (6.00, 9.71, 0.40, 0.14),  # center at L=6.20
}

# Axis title: place to right of "2000" tick (which ends at L=6.40)
AXIS_TITLE_TEXT = "行数 [行/人]"
AXIS_TITLE_POS  = (6.45, 9.70, 1.30, 0.14)  # L, T, W, H

hits_ticks = 0
hits_title = 0
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    try:
        t = sh.top/914400; l = sh.left/914400
    except:
        continue
    # X-axis tick area: T around 9.71, L between 0.9 and 5.6
    if not (9.65 <= t <= 9.80):
        continue
    txt = sh.text_frame.text.strip()
    if txt in TICK_POS:
        new_L, new_T, new_W, new_H = TICK_POS[txt]
        sh.left = Inches(new_L)
        sh.top = Inches(new_T)
        sh.width = Inches(new_W)
        sh.height = Inches(new_H)
        # Set CENTER alignment so text visually centers on tick mark
        for para in sh.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
        # Verify: visual center = new_L + new_W/2
        vis_center = new_L + new_W/2
        print(f"  ✓ X-tick '{txt}' → L={new_L:.2f} W={new_W:.2f} (vis center={vis_center:.2f}) CENTER")
        hits_ticks += 1
    elif txt == AXIS_TITLE_TEXT:
        new_L, new_T, new_W, new_H = AXIS_TITLE_POS
        sh.left = Inches(new_L)
        sh.top = Inches(new_T)
        sh.width = Inches(new_W)
        sh.height = Inches(new_H)
        for para in sh.text_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT
        print(f"  ✓ Axis title '{txt}' → L={new_L:.2f} T={new_T:.2f} (right of 2000-tick at L=6.40)")
        hits_title += 1

assert hits_ticks == 3, f"Expected 3 ticks updated, got {hits_ticks}"
assert hits_title == 1, f"Expected 1 axis title updated, got {hits_title}"

prs.save(PPTX)
print(f"\nWROTE: {PPTX}")
