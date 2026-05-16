#!/usr/bin/env python3
"""
V13 rev1 — in-place edit on V13.pptx (NOT V11/V12).

4 user requirements:
  1. Add team name below title (V12 didn't actually have one — adding new)
     Placeholder: "認証基盤プロジェクト" (from judge card "認証基盤PJ" reference)
  2. Q1合計 cell 1,890 → green text (Forest Green) — signal "全月計画内達成"
  3. KPI cards: separate big value (left) and delta (right) horizontally,
     same Y row, vertically centered
  4. Left ref panel — LOC-centric rewrite:
       ' 外注費用換算 (1,600 千円/人月)' → '業界基準 (新規開発上位25% 1,000行/人月)'
       '約 32,000 [千円] 相当'           → '20,000 行 ≒ 20 人月相当'
       '20人月 × 1,600 千円/人月 …'      → '(≒ 外注費用 32,000 [千円] 相当)'
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V13.pptx"

DEEP_FOREST = RGBColor(0x1F, 0x3A, 0x2E)
FOREST      = RGBColor(0x2E, 0x6F, 0x40)
SAGE_BG     = RGBColor(0xE8, 0xEF, 0xE3)
GREEN_OK    = RGBColor(0x2E, 0x6F, 0x40)
AMBER       = RGBColor(0xF4, 0xB9, 0x42)
AMBER_DEEP  = RGBColor(0xC8, 0x8E, 0x1F)
INK         = RGBColor(0x22, 0x28, 0x24)
MUTED       = RGBColor(0x6B, 0x6B, 0x6B)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
JP_FONT, EN_FONT = "メイリオ", "Segoe UI"

def set_run_font(run, size, color, bold=False):
    run.font.name = EN_FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea_elem = etree.SubElement(rPr, qn('a:ea'))
    ea_elem.set('typeface', JP_FONT)

def replace_text_keep(shape, new_text):
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        return
    p0.runs[0].text = new_text
    for r in list(p0.runs)[1:]:
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)

prs = Presentation(PPTX)
slide = prs.slides[0]

# ============================================================
# 1. Add team name below title (NEW shape)
# ============================================================
# Current header layout:
#   Title at T=0.20 H=0.40 (ends 0.60)
#   Subtitle at T=0.62 H=0.20
# Insert team name at T=0.55, h=0.18 — between title-baseline and subtitle.
# Shift subtitle down to T=0.78.
team_name_text = "認証基盤プロジェクト"

# Shift subtitle down
for sh in slide.shapes:
    if not sh.has_text_frame: continue
    try:
        t_in = sh.top/914400
    except: continue
    if 0.55 <= t_in < 0.75 and "2026年度Q1" in sh.text_frame.text:
        sh.top = Inches(0.80)
        print("  ✓ Subtitle shifted down to T=0.80")

# Add team name shape
tb = slide.shapes.add_textbox(Inches(0.42), Inches(0.58), Inches(6.30), Inches(0.20))
tf = tb.text_frame
tf.word_wrap = False
tf.margin_left = tf.margin_right = Emu(0)
tf.margin_top = tf.margin_bottom = Emu(0)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = team_name_text
set_run_font(r, size=12, color=DEEP_FOREST, bold=True)
print(f"  ✓ Team name added: '{team_name_text}'")

# ============================================================
# 2. Q1合計 cell 1,890 → green text
# ============================================================
for sh in slide.shapes:
    if not sh.has_text_frame: continue
    try:
        t_in = round(sh.top/914400, 2); l_in = round(sh.left/914400, 2)
    except: continue
    # Q1合計 cell at T=4.74 L=3.85
    if t_in == 4.74 and l_in == 3.85 and sh.text_frame.text.strip() == "1,890":
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = GREEN_OK
                run.font.bold = True
        print("  ✓ Q1合計 1,890 → Forest Green text")

# ============================================================
# 3. KPI cards: delta moved to right side, same Y as big value
# ============================================================
# Cards at L=0.30, 2.92, 5.54 (width 2.42 each)
# Current inner shapes:
#   label   T=5.65 (size 8, top of card)
#   big     T=5.78 W=2.26 (centered/left-aligned at L=card+0.08)
#   delta   T=6.06 W=2.02 (small)
#   sub     T=6.16 W=2.26
# New target:
#   label   T=5.65 (unchanged)
#   big     T=5.78 W=1.40 align=left
#   delta   T=5.78 W=0.80 align=right  ← MOVED UP, right-aligned, valign middle
#   sub     T=6.16 (unchanged)
KPI_CARDS_L = [0.30, 2.92, 5.54]
CARD_W = 2.42

for sh in slide.shapes:
    if not sh.has_text_frame: continue
    try:
        t_in = round(sh.top/914400, 2); l_in = round(sh.left/914400, 2)
    except: continue
    # detect KPI tile inner: T=5.78 (big) or T=6.06 (delta)
    if t_in not in {5.78, 6.06}:
        continue
    # Identify which card
    for card_l in KPI_CARDS_L:
        inner_l = card_l + 0.08
        if abs(l_in - inner_l) < 0.05:
            txt = sh.text_frame.text.strip()
            if t_in == 5.78:
                # Big value — shrink width, left-align
                sh.left = Inches(card_l + 0.08)
                sh.width = Inches(1.40)
                for para in sh.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.LEFT
                print(f"  ✓ Big value '{txt}' → L={card_l+0.08:.2f} W=1.40 left-aligned")
            elif t_in == 6.06:
                # Delta — move to T=5.78 (same Y as big), right-aligned, vertically middle
                sh.top = Inches(5.78)
                sh.left = Inches(card_l + 1.52)
                sh.width = Inches(0.82)
                sh.height = Inches(0.28)
                tf = sh.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                for para in tf.paragraphs:
                    para.alignment = PP_ALIGN.RIGHT
                # Slightly bigger delta font for visual balance with big number
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(10)
                        run.font.bold = True
                print(f"  ✓ Delta '{txt}' → T=5.78 L={card_l+1.52:.2f} W=0.82 right-aligned 10pt")
            break

# ============================================================
# 4. Left ref panel — LOC-centric rewrite
# ============================================================
left_swaps = {
    "外注費用換算 (1,600 千円/人月)":           "業界基準 (新規開発上位25% 1,000行/人月)",
    "約 32,000 [千円] 相当":                    "20,000 行 ≒ 20 人月相当",
    "20人月 × 1,600 千円/人月 (新規開発上位25% 1,000行/人月 基準)":
                                                 "(≒ 外注費用 32,000 [千円] 相当)",
}
for sh in slide.shapes:
    if not sh.has_text_frame: continue
    try:
        t_in = sh.top/914400; l_in = sh.left/914400
    except: continue
    if not (10.30 <= t_in < 11.65 and l_in < 4.05):
        continue
    txt = sh.text_frame.text.strip()
    if txt in left_swaps:
        replace_text_keep(sh, left_swaps[txt])
        print(f"  ✓ Left ref: '{txt}' → '{left_swaps[txt]}'")

prs.save(PPTX)
print(f"\nWROTE: {PPTX}")
