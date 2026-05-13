#!/usr/bin/env python3
"""
V11 rev18 — 6-month flexible reporting period (4-9月).

Changes:
  1. Re-render monthly trend chart to show 6 months (4-9月) only
     - 4-6月 actual (forest solid bar + burnt orange solid line)
     - 7-9月 plan (forest hatched bar, no LOC line)
  2. Forecast card RIGHT:
       Label : "年度末着地予測 (予測)" → "9月までの予測コスト"
       Value : 8,640 → 4,140 = 1,890 + 750×3
       外注  : 5.4人月 → 2.6人月 = 4,140 ÷ 1,600
  3. Judge card 投資規模:
       "来期の投資規模: 外注費用 5.4人月相当 (累積+残月計画ベース)"
       → "9月まで: 外注費用 2.6人月相当 (累積+残月計画ベース)"
  4. Plan/actual table → 6 months + 6ヶ月合計
       Delete 10-3月 columns
       Resize 4-9月 columns wider (0.53→0.91 each)
       Reposition 合計 column (text "年間合計"→"6ヶ月合計", value 9,000→4,500)
       Move vertical separator from L=2.53 to L=3.68
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import MaxNLocator

plt.rcParams["font.family"] = ["DejaVu Sans", "Droid Sans Fallback",
                                "Noto Sans CJK JP", "Hiragino Sans",
                                "Yu Gothic", "Meiryo"]
plt.rcParams["axes.unicode_minus"] = False

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V11.pptx"

forest_hex   = "#1F3A2E"
muted_hex    = "#6B6B6B"
burnt_orange = "#D97706"

MONTHS = ['4月','5月','6月','7月','8月','9月']
COST_ACT  = [580, 620, 690]
COST_PLAN = [750, 750, 750]
LOC_ACT   = [15, 17, 20]
ACTUAL_COUNT = 3

# =============================================================
# 1. Re-render chart (6 months only)
# =============================================================
fig, ax_left = plt.subplots(figsize=(7.67, 1.30), dpi=220)
fig.patch.set_facecolor("white")

x_pos = list(range(len(MONTHS)))
# Bars: 4-6 actual solid forest, 7-9 plan hatched
ax_left.bar(x_pos[:ACTUAL_COUNT], COST_ACT,
            color=forest_hex, edgecolor=forest_hex, width=0.55, linewidth=0.5,
            label="単月コスト (実績)", zorder=2)
ax_left.bar(x_pos[ACTUAL_COUNT:], COST_PLAN,
            facecolor="white", edgecolor=forest_hex, hatch="///",
            width=0.55, linewidth=0.6,
            label="単月コスト (計画)", zorder=2)

ax_left.set_xticks(x_pos)
ax_left.set_xticklabels(MONTHS, fontsize=6.5, color=muted_hex)
ax_left.set_ylabel("単月コスト [千円]", color=forest_hex, fontsize=7, labelpad=2)
ax_left.tick_params(axis="y", labelcolor=forest_hex, labelsize=6.5, length=2, pad=1)
ax_left.tick_params(axis="x", length=2, pad=1)
ax_left.grid(True, axis="y", linestyle="--", linewidth=0.4,
             color=muted_hex, alpha=0.35, zorder=1)
ax_left.set_axisbelow(True)
for spine in ("top", "right"):
    ax_left.spines[spine].set_visible(False)
ax_left.spines["left"].set_color(forest_hex)
ax_left.spines["left"].set_linewidth(0.6)
ax_left.spines["bottom"].set_color(muted_hex)
ax_left.spines["bottom"].set_linewidth(0.6)
ax_left.yaxis.set_major_locator(MaxNLocator(nbins=5, integer=True))
ax_left.set_ylim(0, max(max(COST_ACT), max(COST_PLAN)) * 1.25)

# LOC line — actual only (4-6月)
ax_right = ax_left.twinx()
ax_right.plot(x_pos[:ACTUAL_COUNT], LOC_ACT,
              color=burnt_orange, marker="o", markersize=4,
              markerfacecolor=burnt_orange, markeredgecolor=burnt_orange,
              linewidth=2.0, label="単月行数 (実績)", zorder=4)
ax_right.set_ylabel("単月行数 [千行]", color=burnt_orange, fontsize=7, labelpad=2)
ax_right.tick_params(axis="y", labelcolor=burnt_orange, labelsize=6.5,
                     length=2, pad=1)
for spine in ("top", "left"):
    ax_right.spines[spine].set_visible(False)
ax_right.spines["right"].set_color(burnt_orange)
ax_right.spines["right"].set_linewidth(0.8)
ax_right.set_ylim(0, max(LOC_ACT) * 1.40)

# Boundary 6月/7月
ax_left.axvline(x=ACTUAL_COUNT - 0.5, color=muted_hex,
                linestyle=":", linewidth=0.6, alpha=0.7, zorder=1)
boundary_y = max(max(COST_ACT), max(COST_PLAN)) * 1.13
ax_left.text(ACTUAL_COUNT - 0.6, boundary_y, "実績→", color=forest_hex,
             fontsize=6, ha="right", va="center")
ax_left.text(ACTUAL_COUNT - 0.4, boundary_y, "←計画", color=muted_hex,
             fontsize=6, ha="left", va="center")

# Legend external
h1, l1 = ax_left.get_legend_handles_labels()
h2, l2 = ax_right.get_legend_handles_labels()
ax_left.legend(h1 + h2, l1 + l2,
               loc="lower center", bbox_to_anchor=(0.5, 1.02), ncol=3,
               fontsize=6.5, frameon=False,
               handlelength=1.4, handletextpad=0.35,
               columnspacing=1.0, borderpad=0.0,
               labelcolor=[forest_hex, forest_hex, burnt_orange])

plt.subplots_adjust(left=0.06, right=0.95, top=0.85, bottom=0.14)
NEW_CHART_PNG = "/tmp/v11r18_trend.png"
fig.savefig(NEW_CHART_PNG, dpi=220, facecolor="white", edgecolor="none")
plt.close(fig)

# =============================================================
# Apply edits to pptx
# =============================================================
prs = Presentation(PPTX)
slide = prs.slides[0]

def replace_text_keep_props(shape, new_text):
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        return False
    p0.runs[0].text = new_text
    for r in list(p0.runs)[1:]:
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    return True

# --- Replace chart picture ---
old_pic = None
for sh in slide.shapes:
    if sh.shape_type == 13:
        old_pic = sh
        break
assert old_pic is not None
pic_l, pic_t = old_pic.left, old_pic.top
pic_w, pic_h = old_pic.width, old_pic.height
old_pic._element.getparent().remove(old_pic._element)
slide.shapes.add_picture(NEW_CHART_PNG, pic_l, pic_t, width=pic_w, height=pic_h)
print("✓ Chart picture re-rendered (4-9月 only)")

# --- Text swaps ---
TEXT_SWAPS = {
    "年度末着地予測 (予測)": "9月までの予測コスト",
    "8,640": "4,140",
    "外注費用 5.4人月": "外注費用 2.6人月",
    "来期の投資規模: 外注費用 5.4人月相当 (累積+残月計画ベース)":
        "9月まで: 外注費用 2.6人月相当 (累積+残月計画ベース)",
    # Table totals
    "年間合計": "6ヶ月合計",
    "9,000": "4,500",
}
n_text = 0
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    raw = sh.text_frame.text.strip()
    if raw in TEXT_SWAPS:
        replace_text_keep_props(sh, TEXT_SWAPS[raw])
        n_text += 1
        print(f"  ✓ '{raw}' → '{TEXT_SWAPS[raw]}'")
print(f"Text swaps: {n_text}")

# --- Table restructure: delete 10-3月 columns, resize 4-9月, move 合計 ---
# Identify shapes by current positions:
#   10-3月 header cells: L in {4.11, 4.64, 5.16, 5.69, 6.22, 6.74}, T=4.42
#   10-3月 計画 cells:    same L,  T=4.58
#   10-3月 実績 cells:    same L,  T=4.74
#   (年間合計 cells at L=7.27 are already retextd by the swap above)

DELETE_L = {round(x, 2) for x in (4.11, 4.64, 5.16, 5.69, 6.22, 6.74)}
DELETE_T = {round(x, 2) for x in (4.42, 4.58, 4.74)}

to_delete = []
for sh in slide.shapes:
    try:
        t_in = round(sh.top/914400, 2)
        l_in = round(sh.left/914400, 2)
    except Exception:
        continue
    if t_in in DELETE_T and l_in in DELETE_L:
        to_delete.append(sh)
for sh in to_delete:
    sh._element.getparent().remove(sh._element)
print(f"Deleted {len(to_delete)} table cells (10-3月)")

# Resize 4-9月 columns (header, 計画, 実績) — keep T but widen
# Old: L = 0.95, 1.48, 2.00, 2.53, 3.06, 3.58 each W=0.53
# New: L = 0.95, 1.86, 2.77, 3.68, 4.59, 5.50 each W=0.91
OLD_L_NEW = {
    0.95: (0.95, 0.91),
    1.48: (1.86, 0.91),
    2.00: (2.77, 0.91),
    2.53: (3.68, 0.91),
    3.06: (4.59, 0.91),
    3.58: (5.50, 0.91),
}
TABLE_TS = {4.42, 4.58, 4.74}
n_moved = 0
for sh in slide.shapes:
    try:
        t_in = round(sh.top/914400, 2)
        l_in = round(sh.left/914400, 2)
    except Exception:
        continue
    if t_in in TABLE_TS and l_in in OLD_L_NEW:
        new_l, new_w = OLD_L_NEW[l_in]
        sh.left = Inches(new_l)
        sh.width = Inches(new_w)
        n_moved += 1
print(f"Repositioned {n_moved} table cells (4-9月)")

# Move 合計 column (originally L=7.27, W=0.70) to L=6.41, W=1.56
for sh in slide.shapes:
    try:
        t_in = round(sh.top/914400, 2)
        l_in = round(sh.left/914400, 2)
    except Exception:
        continue
    if t_in in TABLE_TS and l_in == 7.27:
        sh.left = Inches(6.41)
        sh.width = Inches(1.56)
        print(f"  ✓ 合計 cell moved: L=7.27→6.41, W=0.70→1.56 (T={t_in})")

# Move separator line — originally L=2.53 between 6月/7月; new position L=3.68
for sh in slide.shapes:
    try:
        t_in = round(sh.top/914400, 2)
        l_in = round(sh.left/914400, 2)
        w_in = round(sh.width/914400, 2)
    except Exception:
        continue
    # Look for vertical line (LINE shape, height ≈0.48, width ≈0)
    if sh.shape_type == 9 and abs(l_in - 2.53) < 0.01 and t_in == 4.42:
        sh.left = Inches(3.68)
        print(f"  ✓ Separator line moved: L=2.53 → 3.68")

prs.save(PPTX)
print(f"\nWROTE: {PPTX}")
