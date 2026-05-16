#!/usr/bin/env python3
"""
V13 rev8 — KPI 内縦間隔拡張 + TOP5 重なり解消 + KPI を 5月ベースに統一.

Layout changes:
  KPI card height: 0.90 → 1.10 (T=5.60, bottom=6.70)
  KPI internal stack with padding:
    Label  T=5.66 H=0.16 size=10
    gap 0.06
    Value  T=5.88 H=0.40 size=28
    gap 0.06
    Delta  T=6.34 H=0.18 size=13
    gap 0.06
    Sub    T=6.58 H=0.10 size=8

  TOP5 cards:
    bg T=6.50→6.85, H=1.40→1.20 (bottom 8.05)
    Gap KPI bottom 6.70 → TOP5 top 6.85 = 0.15 ✓
    Internal rows compressed: spacing 0.20 → 0.18

  Title 全社14: T=8.10 unchanged. Gap to TOP5 = 8.10-8.05 = 0.05

KPI data update (5月 統一):
  Card1: '6月 コスト [千円]' / 690 / '+11.3%' / 'vs 5月 620 [千円]'
      →  '5月 コスト [千円]' / 620 / '+6.9%'  / 'vs 4月 580 [千円]'  [+ → GREEN]
  Card2: 'ユーザー数 (MAU)' / '12名' / '+2名' / 'vs 5月 10名'
      →  'ユーザー数 (MAU)' / '10名' / '+2名' / 'vs 4月 8名'         [+ → GREEN]
  Card3: '6月 1人あたり月額 [千円]' / '57.5' / '-4.5%' / 'vs 5月 62.0 [千円]'
      →  '5月 1人あたり月額 [千円]' / '62.0' / '-14.5%' / 'vs 4月 72.5 [千円]' [- → RED]
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V13.pptx"

DEEP_FOREST = RGBColor(0x1F, 0x3A, 0x2E)
FOREST      = RGBColor(0x2E, 0x6F, 0x40)
RED_CRIMSON = RGBColor(0xC0, 0x2A, 0x2A)
MUTED       = RGBColor(0x6B, 0x6B, 0x6B)
JP_FONT, EN_FONT = "メイリオ", "Segoe UI"

def set_ea(run):
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', JP_FONT)

def restyle(shape, new_text, *, size, color, bold, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p0 = tf.paragraphs[0]
    if not p0.runs:
        r = p0.add_run()
    else:
        p0.runs[0].text = new_text
        for r in list(p0.runs)[1:]:
            r._r.getparent().remove(r._r)
    if p0.runs:
        p0.runs[0].text = new_text
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    r0 = p0.runs[0]
    r0.font.name = EN_FONT
    r0.font.size = Pt(size)
    r0.font.bold = bold
    r0.font.color.rgb = color
    set_ea(r0)
    p0.alignment = align

prs = Presentation(PPTX)
slide = prs.slides[0]

# ============================================================
# 1) KPI cards: change data, restyle layout with padding
# ============================================================
CARD_W = 2.42
CARD_H_NEW = 1.10  # 5.60 → 6.70
LABEL_T, LABEL_H = 5.66, 0.16
VALUE_T, VALUE_H = 5.88, 0.40
DELTA_T, DELTA_H = 6.34, 0.18
SUB_T,   SUB_H   = 6.58, 0.10

CARDS = [
    {  # Card 1: 5月 コスト
        "L": 0.30,
        "label":  ("5月 コスト [千円]",   10, FOREST,      True),
        "value":  ("620",                 28, DEEP_FOREST, True),
        "delta":  ("+6.9%",               13, FOREST,      True),
        "sub":    ("vs 4月 580 [千円]",    8, MUTED,       False),
        "old_label": "6月 コスト [千円]",
        "old_value": "690",
        "old_delta": "+11.3%",
        "old_sub":   "vs 5月 620 [千円]",
    },
    {  # Card 2: MAU
        "L": 2.92,
        "label":  ("ユーザー数 (MAU)",    10, FOREST,      True),
        "value":  ("10名",                28, DEEP_FOREST, True),
        "delta":  ("+2名",                13, FOREST,      True),
        "sub":    ("vs 4月 8名",           8, MUTED,       False),
        "old_label": "ユーザー数 (MAU)",
        "old_value": "12名",
        "old_delta": "+2名",
        "old_sub":   "vs 5月 10名",
    },
    {  # Card 3: 5月 1人月額
        "L": 5.54,
        "label":  ("5月 1人あたり月額 [千円]", 10, FOREST,      True),
        "value":  ("62.0",                    28, DEEP_FOREST, True),
        "delta":  ("-14.5%",                  13, RED_CRIMSON, True),
        "sub":    ("vs 4月 72.5 [千円]",       8, MUTED,       False),
        "old_label": "6月 1人あたり月額 [千円]",
        "old_value": "57.5",
        "old_delta": "-4.5%",
        "old_sub":   "vs 5月 62.0 [千円]",
    },
]

# 1a) Resize KPI card backgrounds (NOT accent stripes)
for sh in slide.shapes:
    try:
        t = sh.top/914400; l = sh.left/914400; w = sh.width/914400; h = sh.height/914400
    except: continue
    if not (5.55 <= t <= 5.65 and h > 0.50):
        continue
    for card in CARDS:
        if abs(l - card["L"]) < 0.05 and abs(w - CARD_W) < 0.05:
            sh.height = Inches(CARD_H_NEW)

# 1b) Restyle each KPI text element
for card in CARDS:
    card_l = card["L"]
    inner_l = card_l + 0.08
    inner_w = CARD_W - 0.16
    hits = 0
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        try:
            t = sh.top/914400; l = sh.left/914400
        except: continue
        if not (5.60 <= t <= 6.55): continue
        if not (card_l - 0.05 <= l <= card_l + CARD_W):
            continue
        txt = sh.text_frame.text.strip()
        which = None
        if txt == card["old_label"]: which = "label"
        elif txt == card["old_value"]: which = "value"
        elif txt == card["old_delta"]: which = "delta"
        elif txt == card["old_sub"]: which = "sub"
        if which is None: continue

        new_text, size, color, bold = card[which]
        restyle(sh, new_text, size=size, color=color, bold=bold, align=PP_ALIGN.CENTER)
        sh.left = Inches(inner_l)
        sh.width = Inches(inner_w)
        if which == "label":
            sh.top = Inches(LABEL_T); sh.height = Inches(LABEL_H)
        elif which == "value":
            sh.top = Inches(VALUE_T); sh.height = Inches(VALUE_H)
        elif which == "delta":
            sh.top = Inches(DELTA_T); sh.height = Inches(DELTA_H)
        elif which == "sub":
            sh.top = Inches(SUB_T); sh.height = Inches(SUB_H)
        hits += 1
        print(f"  ✓ KPI L={card_l:.2f} {which:6s} → '{new_text}' ({size}pt)")
    assert hits == 4, f"KPI Card L={card_l} expected 4 hits, got {hits}"

# ============================================================
# 2) TOP5 cards: shift down + shrink + compress rows
# ============================================================
# bg: T=6.50 → 6.85, H=1.40 → 1.20 (bottom 8.05)
# accent stripe: T=6.50 → 6.85
# header: T=6.57 → 6.91
# row T mapping (old → new): 6.84 → 7.15, 7.04 → 7.33, 7.24 → 7.51, 7.44 → 7.69, 7.64 → 7.87
#                           (alt 6.82) → (alt 7.13), 7.02→7.31, 7.22→7.49, 7.42→7.67, 7.62→7.85

ROW_MAP_RANK = {6.84: 7.15, 7.04: 7.33, 7.24: 7.51, 7.44: 7.69, 7.64: 7.87}
ROW_MAP_NAME = {6.82: 7.13, 7.02: 7.31, 7.22: 7.49, 7.42: 7.67, 7.62: 7.85}

TOP5_BG_NEW_T = 6.85
TOP5_BG_NEW_H = 1.20
TOP5_HEADER_NEW_T = 6.91

# Gather TOP5 shapes (T 6.50-7.90)
top5_shifts = 0
for sh in slide.shapes:
    try:
        t = sh.top/914400; l = sh.left/914400; w = sh.width/914400; h = sh.height/914400
    except: continue
    if not (6.50 <= t <= 7.95): continue
    # Identify type
    if abs(t - 6.50) < 0.02 and h > 1.30:
        # Background card
        sh.top = Inches(TOP5_BG_NEW_T)
        sh.height = Inches(TOP5_BG_NEW_H)
        top5_shifts += 1
    elif abs(t - 6.50) < 0.02 and h < 0.10:
        # Accent stripe
        sh.top = Inches(TOP5_BG_NEW_T)
        top5_shifts += 1
    elif abs(t - 6.57) < 0.02:
        # Header text or unit text
        sh.top = Inches(TOP5_HEADER_NEW_T)
        top5_shifts += 1
    else:
        # Row item (rank circle, rank number, name, value)
        # Match by closest old T
        new_t = None
        for old_t, target_t in ROW_MAP_RANK.items():
            if abs(t - old_t) < 0.02:
                new_t = target_t; break
        if new_t is None:
            for old_t, target_t in ROW_MAP_NAME.items():
                if abs(t - old_t) < 0.02:
                    new_t = target_t; break
        if new_t is not None:
            sh.top = Inches(new_t)
            top5_shifts += 1

print(f"  ✓ TOP5 shifts applied: {top5_shifts}")

prs.save(PPTX)
print(f"\nWROTE: {PPTX}")
