#!/usr/bin/env python3
"""
V11 rev12 — in-place: unify the X-axis label「使用率」with the Y-axis style.

Y-axis: "↑ 効率\n(行/コスト)"  — 7pt Bold, Deep Forest
X-axis (rev11): "使用率"        — 10pt Bold, Forest Green
       (rev12): "→ 使用率"      — 7pt Bold, Deep Forest  ← align to Y style

Find the existing X-axis textbox by:
  - text containing "使用率"
  - Y position in 9.50–9.90 (below plot, above ▲ row)
Replace the run text and font size/color. Tighten the box width so the
shorter 7pt text sits cleanly centered under the plot.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V11.pptx"

DEEP_FOREST = RGBColor(0x1F, 0x3A, 0x2E)
JP_FONT = "メイリオ"
EN_FONT = "Segoe UI"

prs = Presentation(PPTX)
slide = prs.slides[0]

n_changed = 0
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    try:
        t_in = sh.top / 914400
    except Exception:
        continue
    # Target: X-axis label area, Y in 9.50–9.90
    if not (9.50 <= t_in < 9.90):
        continue
    tf = sh.text_frame
    full = tf.text
    if "使用率" not in full:
        continue
    # Modify run(s) text + font props
    for para in tf.paragraphs:
        for run in para.runs:
            if "使用率" in run.text:
                # Replace with arrow-prefixed form to mirror Y-axis "↑ 効率"
                run.text = "→ 使用率"
                # Style → Y-axis spec: 7pt Bold Deep Forest, EN font Segoe UI
                run.font.name = EN_FONT
                run.font.size = Pt(7)
                run.font.color.rgb = DEEP_FOREST
                run.font.bold = True
                # East-Asia font (Meiryo) — preserve Japanese rendering
                rPr = run._r.get_or_add_rPr()
                for old in rPr.findall(qn('a:ea')):
                    rPr.remove(old)
                ea_elem = etree.SubElement(rPr, qn('a:ea'))
                ea_elem.set('typeface', JP_FONT)
                n_changed += 1
    # Tighten the box width so the now-smaller text sits compact and
    # remains center-aligned under the plot.
    # Current box: L=3.30 W=0.80 H=0.16; for 7pt label width≈0.55 is enough.
    new_left = Inches(3.45)   # was 3.30
    new_width = Inches(0.50)  # was 0.80
    sh.left = new_left
    sh.width = new_width

print(f"Changed {n_changed} run(s).")
assert n_changed >= 1, "X-axis label '使用率' textbox not found."

prs.save(PPTX)
print(f"WROTE: {PPTX}")
