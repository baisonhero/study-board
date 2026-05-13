#!/usr/bin/env python3
"""
V11 rev11 — in-place: add X-axis label「使用率」to the team mapping plot.

Y-axis label「↑ 効率 (行/コスト)」 already exists at (L=0.35, T=8.88) — keep as is.
X-axis label NOT present → add at the gap between plot bottom edge (T=9.70) and
the ▲ bottom-row quadrant labels (T=9.86). Centered at plot horizontal middle
(x ≈ 3.70).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V11.pptx"

FOREST      = RGBColor(0x2E, 0x6F, 0x40)
DEEP_FOREST = RGBColor(0x1F, 0x3A, 0x2E)
JP_FONT = "メイリオ"
EN_FONT = "Segoe UI"

def set_run_font(run, size, color, bold=True):
    run.font.name = EN_FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', JP_FONT)

prs = Presentation(PPTX)
slide = prs.slides[0]

# Idempotency guard: don't add twice if rev11 already ran
for sh in slide.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip() == "使用率":
        # Check it's in team-mapping area (Y near 9.7)
        try:
            t_in = sh.top / 914400
            if 9.50 < t_in < 9.90:
                print("X-axis label already present — skipping.")
                raise SystemExit(0)
        except Exception:
            pass

# Plot geometry (from rev6 onwards):
#   PX = 1.20, PY = 8.65 (= TM_T 8.35 + 0.30), PW = 5.00, PH = 1.05
# Bottom quadrant labels at T = 9.86.
# Plot bottom edge at PY+PH = 9.70.
# Available gap: 9.70 - 9.86 = 0.16 inch for the X axis label.
PX, PW = 1.20, 5.00
plot_center = PX + PW / 2          # = 3.70

# Place "使用率" centered horizontally at plot center, between bottom edge
# and ▲ row. Narrow box (0.80 in) so it never overlaps the quadrant labels
# (left label ends at x≈3.60, right label starts at x≈3.80).
LABEL_X = plot_center - 0.40       # = 3.30
LABEL_W = 0.80
LABEL_Y = 9.70
LABEL_H = 0.16

tb = slide.shapes.add_textbox(Inches(LABEL_X), Inches(LABEL_Y),
                              Inches(LABEL_W), Inches(LABEL_H))
tf = tb.text_frame
tf.word_wrap = False
tf.margin_left = tf.margin_right = Emu(0)
tf.margin_top = tf.margin_bottom = Emu(0)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "使用率"
set_run_font(r, size=10, color=FOREST, bold=True)
print(f"Added X-axis label '使用率' at L={LABEL_X:.2f} T={LABEL_Y:.2f} "
      f"W={LABEL_W:.2f} H={LABEL_H:.2f}")

prs.save(PPTX)
print(f"WROTE: {PPTX}")
