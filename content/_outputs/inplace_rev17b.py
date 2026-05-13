#!/usr/bin/env python3
"""
V11 rev17b — fix LEFT title wrap.
"前月までの累積コスト (実績)" → "前月までの累積 (実績)" (drop "コスト" to
parallel the RIGHT side "年度末着地予測 (予測)" — both 8 char + qualifier).
"""
from pptx import Presentation

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V11.pptx"

prs = Presentation(PPTX)
slide = prs.slides[0]

n = 0
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    try:
        t_in = sh.top/914400; l_in = sh.left/914400
    except Exception:
        continue
    if not (0.95 <= t_in < 1.30 and 4.95 <= l_in < 6.50):
        continue
    if "前月までの累積コスト" in sh.text_frame.text:
        p0 = sh.text_frame.paragraphs[0]
        if p0.runs:
            p0.runs[0].text = "前月までの累積 (実績)"
            for r in list(p0.runs)[1:]:
                r._r.getparent().remove(r._r)
        n += 1
        print(f"  ✓ Title shortened: '前月までの累積コスト (実績)' → '前月までの累積 (実績)'")

assert n == 1, f"Expected 1 update, got {n}"
prs.save(PPTX)
print(f"WROTE: {PPTX}")
