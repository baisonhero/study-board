#!/usr/bin/env python3
"""
V13 rev2 — Apply remaining 3 user requirements (5, 6, 7) consolidated into
a single in-place edit, plus all numerical cascade from "5月末カットオフ" model.

CHANGES:
  (5) Monthly chart: 4-5月 actual, 6月 plan
      → cascading: Q1累積実績 1,890 → 1,200 (4+5月 only)
      → 6月実績欄 → "-" / Q1合計消費率 → 53.3%
      → "6月の実績" divider → "5月の実績 (前月)"
      → KPI/TOP5/ref panel labels & numbers re-grounded to 5月
      → Judge / forecast cards updated
      Mock data for 5月: MAU=10, 1人月額=62.0, 1人行数=1700,
                          team total LOC=17,000

  (6) Plan/actual table layout fixed: uniform 4 columns × 1.78" wide

  (7) Team mapping:
      (a) Add 5-tick axis labels on both axes
      (b) Move title to top-right (smaller, less prominent)
      (c) Y axis label simplified: "↑ 1人当たり / コスト (千円/人)"
                                    → "コスト (千円/人)"
          X axis label aligned style: "行数 (行/人)"
      → Bubble for own team updated to 5月 data: X=1700, Y=62, MAU=10
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import MaxNLocator

plt.rcParams["font.family"] = ["DejaVu Sans", "Droid Sans Fallback",
                                "Noto Sans CJK JP", "Hiragino Sans",
                                "Yu Gothic", "Meiryo"]
plt.rcParams["axes.unicode_minus"] = False

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V13.pptx"

DEEP_FOREST = RGBColor(0x1F, 0x3A, 0x2E)
FOREST      = RGBColor(0x2E, 0x6F, 0x40)
SAGE        = RGBColor(0xB5, 0xC9, 0xA8)
SAGE_BG     = RGBColor(0xE8, 0xEF, 0xE3)
AMBER       = RGBColor(0xF4, 0xB9, 0x42)
AMBER_DEEP  = RGBColor(0xC8, 0x8E, 0x1F)
INK         = RGBColor(0x22, 0x28, 0x24)
MUTED       = RGBColor(0x6B, 0x6B, 0x6B)
CARD_BORDER = RGBColor(0xD7, 0xDC, 0xCF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_OK    = RGBColor(0x2E, 0x6F, 0x40)
COLOR_OWN   = AMBER
COLOR_OTHER = RGBColor(0xA8, 0xB8, 0xA0)
JP_FONT, EN_FONT = "メイリオ", "Segoe UI"

def set_run_font(run, size=10, color=INK, bold=False, latin=EN_FONT, ea=JP_FONT):
    run.font.name = latin
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea_elem = etree.SubElement(rPr, qn('a:ea'))
    ea_elem.set('typeface', ea)

def replace_text_keep(shape, new_text, *, size=None, color=None, bold=None):
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs: return False
    r0 = p0.runs[0]
    r0.text = new_text
    if size is not None: r0.font.size = Pt(size)
    if color is not None: r0.font.color.rgb = color
    if bold is not None: r0.font.bold = bold
    for r in list(p0.runs)[1:]:
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    return True

def add_textbox(slide, text, l, t, w, h, *, size=12, color=INK, bold=False,
                align='left', valign='top', word_wrap=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if valign == 'middle': tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if align == 'center': p.alignment = PP_ALIGN.CENTER
    elif align == 'right': p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    set_run_font(r, size=size, color=color, bold=bold)
    return tb

def add_rect(slide, l, t, w, h, *, fill=WHITE, line=None, line_width=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                                 Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        if line_width: shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_oval(slide, l, t, w, h, *, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t),
                                 Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_line(slide, x1, y1, x2, y2, *, color=CARD_BORDER, width=0.5):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln

# =========================================================
# (5a) Re-render chart with 5月 cutoff
# =========================================================
forest_hex   = "#1F3A2E"
muted_hex    = "#6B6B6B"
burnt_orange = "#D97706"

MONTHS = ['4月','5月','6月']
COST_ACT  = [580, 620]              # 4-5月 actual
COST_PLAN = [750]                    # 6月 plan
LOC_ACT   = [15, 17]                 # 4-5月 actual
ACTUAL_COUNT = 2                     # ← 5月末 cutoff

fig, ax_left = plt.subplots(figsize=(7.67, 1.30), dpi=220)
fig.patch.set_facecolor("white")
x_pos = list(range(len(MONTHS)))
ax_left.bar(x_pos[:ACTUAL_COUNT], COST_ACT,
            color=forest_hex, edgecolor=forest_hex, width=0.55, linewidth=0.5,
            label="単月コスト (実績)", zorder=2)
ax_left.bar(x_pos[ACTUAL_COUNT:], COST_PLAN,
            facecolor="white", edgecolor=forest_hex, hatch="///",
            width=0.55, linewidth=0.6,
            label="単月コスト (計画)", zorder=2)
ax_left.set_xticks(x_pos)
ax_left.set_xticklabels(MONTHS, fontsize=7, color=muted_hex)
ax_left.set_ylabel("単月コスト [千円]", color=forest_hex, fontsize=7, labelpad=2)
ax_left.tick_params(axis="y", labelcolor=forest_hex, labelsize=6.5, length=2, pad=1)
ax_left.tick_params(axis="x", length=2, pad=1)
ax_left.grid(True, axis="y", linestyle="--", linewidth=0.4,
             color=muted_hex, alpha=0.35, zorder=1)
ax_left.set_axisbelow(True)
for spine in ("top", "right"):
    ax_left.spines[spine].set_visible(False)
ax_left.spines["left"].set_color(forest_hex)
ax_left.spines["left"].set_linewidth(0.6)
ax_left.spines["bottom"].set_color(muted_hex)
ax_left.spines["bottom"].set_linewidth(0.6)
ax_left.yaxis.set_major_locator(MaxNLocator(nbins=5, integer=True))
ax_left.set_ylim(0, max(max(COST_ACT), max(COST_PLAN)) * 1.25)

ax_right = ax_left.twinx()
ax_right.plot(x_pos[:ACTUAL_COUNT], LOC_ACT,
              color=burnt_orange, marker="o", markersize=4,
              markerfacecolor=burnt_orange, markeredgecolor=burnt_orange,
              linewidth=2.0, label="単月行数 (実績)", zorder=4)
ax_right.set_ylabel("単月行数 [千行]", color=burnt_orange, fontsize=7, labelpad=2)
ax_right.tick_params(axis="y", labelcolor=burnt_orange, labelsize=6.5, length=2, pad=1)
for spine in ("top", "left"):
    ax_right.spines[spine].set_visible(False)
ax_right.spines["right"].set_color(burnt_orange)
ax_right.spines["right"].set_linewidth(0.8)
ax_right.set_ylim(0, max(LOC_ACT) * 1.40)

# Boundary at 5/6月
ax_left.axvline(x=ACTUAL_COUNT - 0.5, color=muted_hex,
                linestyle=":", linewidth=0.6, alpha=0.7, zorder=1)
boundary_y = max(max(COST_ACT), max(COST_PLAN)) * 1.13
ax_left.text(ACTUAL_COUNT - 0.55, boundary_y, "実績→", color=forest_hex,
             fontsize=6, ha="right", va="center")
ax_left.text(ACTUAL_COUNT - 0.45, boundary_y, "←計画", color=muted_hex,
             fontsize=6, ha="left", va="center")

h1, l1 = ax_left.get_legend_handles_labels()
h2, l2 = ax_right.get_legend_handles_labels()
ax_left.legend(h1 + h2, l1 + l2,
               loc="lower center", bbox_to_anchor=(0.5, 1.02), ncol=3,
               fontsize=6.5, frameon=False,
               handlelength=1.4, handletextpad=0.35,
               columnspacing=1.0, borderpad=0.0,
               labelcolor=[forest_hex, forest_hex, burnt_orange])
plt.subplots_adjust(left=0.06, right=0.95, top=0.85, bottom=0.16)
CHART_PNG = "/tmp/v13r2_trend.png"
fig.savefig(CHART_PNG, dpi=220, facecolor="white", edgecolor="none")
plt.close(fig)
print("✓ Chart rendered (4-5月 実績 / 6月 計画)")

# =========================================================
# Open + edit
# =========================================================
prs = Presentation(PPTX)
slide = prs.slides[0]

# --- Replace chart picture ---
for sh in slide.shapes:
    if sh.shape_type == 13:
        l, t, w, h = sh.left, sh.top, sh.width, sh.height
        sh._element.getparent().remove(sh._element)
        slide.shapes.add_picture(CHART_PNG, l, t, width=w, height=h)
        print("✓ Chart picture swapped")
        break

# =========================================================
# (5b) Numerical cascade — all 5月-cutoff updates via text swaps
# =========================================================
# Q1累積実績: 1,890 → 1,200 (= 580+620)
# Q1合計消費率: 84% → 53% (= 1200/2250)
SWAPS = {
    # Forecast cards
    "Q1実績累積":                                   "Q1進捗実績 (4-5月)",
    # NOTE: keep "1,890" → "1,200" only for the BIG forecast number — there are
    # also "1,890" cells in table; handle those by position separately
    # Forecast RIGHT card content
    "計画消費率 84%":                                "計画進捗率 53%",
    "(実績 1,890 / 計画 2,250)":                     "(実績 1,200 / 計画 2,250)",
    # Judge card
    "Q1消費率 84% (実績 1,890 / 計画 2,250 [千円])":
        "Q1進捗 53% (5月末 / 実績 1,200 / 計画 2,250 [千円])",
    # Effizienz line: "過去3ヶ月で..." → "過去2ヶ月で..."
    "過去3ヶ月でコスト 1.19倍 vs 採用行数 1.33倍":
        "過去2ヶ月でコスト 1.07倍 vs 採用行数 1.13倍",
    # "6月の実績" divider title
    "6月の実績":                                    "5月の実績 (前月)",
    "↓ 以下は 2026 年 6 月 単月のスナップショット":
        "↓ 以下は 2026 年 5 月 単月のスナップショット",
    # KPI labels
    "6月 コスト [千円]":                            "5月 コスト [千円]",
    "MAU [6月]":                                    "MAU [5月]",
    "6月 1人あたり月額 [千円]":                      "5月 1人あたり月額 [千円]",
    # KPI big numbers (replace by exact value)
    # NOTE: handle 690→620, 57.5→62.0, 12名→10名 by position to avoid touching
    # other cells with same strings
    # KPI delta lines
    "+11.3%":      "+6.9%",
    "−4.5%":       "−14.5%",
    "+2名":        "+2名",        # keep — assume 5月 MAU 10 vs 4月 MAU 8
    "vs 5月 620 [千円]": "vs 4月 580 [千円]",
    "vs 5月 10名":      "vs 4月 8名",
    "vs 5月 62.0 [千円]": "vs 4月 72.5 [千円]",
    # TOP5/TOP3 panel headers
    "6月 コスト TOP5":  "5月 コスト TOP5",
    "6月 生成行 TOP5":  "5月 生成行 TOP5",
    "6月 効率 TOP5":    "5月 効率 TOP5",
    # Left ref panel (re-ground to 5月 mock data)
    "当チーム実績 (6月)":         "当チーム実績 (5月)",
    "690 [千円] → 20,000行":     "620 [千円] → 17,000行",
    "業界基準 (新規開発上位25% 1,000行/人月)":   "業界基準 (新規開発上位25% 1,000行/人月)",  # no change
    "20,000 行 ≒ 20 人月相当":   "17,000 行 ≒ 17 人月相当",
    "(≒ 外注費用 32,000 [千円] 相当)":  "(≒ 外注費用 27,200 [千円] 相当)",
}

# Apply position-targeted KPI big value swaps
KPI_BIG_SWAPS = {
    # (T, L, old) → new
    (5.78, 0.38, "690"):  "620",
    (5.78, 3.00, "12名"): "10名",
    (5.78, 5.62, "57.5"): "62.0",
}
# Forecast card LEFT big number 1,890 (T=1.24 L=5.00)
FORECAST_LEFT_BIG = (1.24, 5.00, "1,890")

# Apply position-targeted Q1合計 cell swaps (table row 実績)
TABLE_Q1_ACTUAL = (4.74, 3.85, "1,890")
# 消費率 cell of 6月 — change to "-"
TABLE_JUN_RATE  = (4.90, 2.77, "92.0%")
# Q1合計 消費率 → 53.3%
TABLE_Q1_RATE   = (4.90, 3.85, "84.0%")
# 6月 実績 cell → "-"
TABLE_JUN_ACTUAL = (4.74, 2.77, "690")

n_text = 0
for sh in slide.shapes:
    if not sh.has_text_frame: continue
    try:
        t_in = round(sh.top/914400, 2); l_in = round(sh.left/914400, 2)
    except: continue
    txt = sh.text_frame.text.strip()

    # KPI big values (position-targeted)
    if (t_in, l_in, txt) in KPI_BIG_SWAPS:
        replace_text_keep(sh, KPI_BIG_SWAPS[(t_in, l_in, txt)])
        n_text += 1
        print(f"  ✓ KPI big: '{txt}' → '{KPI_BIG_SWAPS[(t_in, l_in, txt)]}'")
        continue

    # Forecast LEFT big 1,890 → 1,200
    if (t_in, l_in, txt) == FORECAST_LEFT_BIG:
        replace_text_keep(sh, "1,200")
        n_text += 1
        print(f"  ✓ Forecast LEFT big: 1,890 → 1,200")
        continue

    # Table Q1合計 実績 cell 1,890 → 1,200 (also color green per rev1)
    if (t_in, l_in, txt) == TABLE_Q1_ACTUAL:
        replace_text_keep(sh, "1,200")
        # Keep green color (set in rev1)
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = GREEN_OK
                r.font.bold = True
        n_text += 1
        print(f"  ✓ Table Q1合計 実績: 1,890 → 1,200 (green)")
        continue

    # Table 6月 実績 → "-"
    if (t_in, l_in, txt) == TABLE_JUN_ACTUAL:
        replace_text_keep(sh, "-")
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = MUTED
        n_text += 1
        print(f"  ✓ Table 6月 実績: 690 → '-'")
        continue

    # Table 6月 消費率 → "-"
    if (t_in, l_in, txt) == TABLE_JUN_RATE:
        replace_text_keep(sh, "-")
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = MUTED
        n_text += 1
        print(f"  ✓ Table 6月 消費率: 92.0% → '-'")
        continue

    # Table Q1合計 消費率 → 53.3%
    if (t_in, l_in, txt) == TABLE_Q1_RATE:
        replace_text_keep(sh, "53.3%")
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = GREEN_OK
                r.font.bold = True
        n_text += 1
        print(f"  ✓ Table Q1 消費率: 84.0% → 53.3%")
        continue

    # General text swaps
    if txt in SWAPS and txt != SWAPS[txt]:
        replace_text_keep(sh, SWAPS[txt])
        n_text += 1
        print(f"  ✓ '{txt}' → '{SWAPS[txt]}'")

print(f"Total text/positioned swaps: {n_text}")

# =========================================================
# (6) Plan/actual table layout — uniform 4 columns × 1.78"
# =========================================================
# New positions:
#   Label col: L=0.30 W=0.55
#   4月:  L=0.85 W=1.78  (ends 2.63)
#   5月:  L=2.63 W=1.78  (ends 4.41)
#   6月:  L=4.41 W=1.78  (ends 6.19)
#   Q1合計: L=6.19 W=1.78  (ends 7.97)
NEW_COLS = {
    "4月_label": (0.85, 1.78),
    "5月_label": (2.63, 1.78),
    "6月_label": (4.41, 1.78),
    "Q1合計_label": (6.19, 1.78),
}
# Map current L (after V13 build_v13) → new L
# In current V13 the columns are at L = 0.95, 1.86, 2.77, 3.85 (with widths 0.91, 0.91, 0.91, 4.12)
OLD_L_TO_NEW = {
    0.95: (0.85, 1.78),
    1.86: (2.63, 1.78),
    2.77: (4.41, 1.78),
    3.85: (6.19, 1.78),
}
TABLE_TS = {4.42, 4.58, 4.74, 4.90}

for sh in slide.shapes:
    try:
        t_in = round(sh.top/914400, 2); l_in = round(sh.left/914400, 2)
    except: continue
    if t_in not in TABLE_TS: continue
    if l_in in OLD_L_TO_NEW:
        new_l, new_w = OLD_L_TO_NEW[l_in]
        sh.left = Inches(new_l)
        sh.width = Inches(new_w)

# Also resize/reposition table separator line (currently at L=3.68 from rev18b)
# New: between 5月 and 6月 → L=4.41
for sh in slide.shapes:
    if sh.shape_type == 9:   # LINE
        try:
            t_in = round(sh.top/914400, 2); l_in = round(sh.left/914400, 2)
        except: continue
        if 4.30 <= t_in < 5.10 and 3.50 < l_in < 4.50:
            sh.left = Inches(4.41)
            sh.width = Inches(0)
            print("  ✓ Table separator → L=4.41 (between 5月 and 6月)")

print("✓ Table layout updated (4 uniform 1.78\" columns)")

# =========================================================
# (7) Team mapping — axis ticks, title repositioning, simplified labels
# =========================================================
# Plot geometry (from V13 build): PX=1.20 PY=8.65 PW=5.00 PH=1.05
PX, PY, PW, PH = 1.20, 8.65, 5.00, 1.05
X_MIN, X_MAX = 0, 2400        # 行/人
Y_MIN, Y_MAX = 0, 120         # 千円/人

# (7c) Y axis label simplify
for sh in slide.shapes:
    if sh.has_text_frame and "1人当たり" in sh.text_frame.text and "コスト" in sh.text_frame.text:
        tf = sh.text_frame
        p0 = tf.paragraphs[0]
        if p0.runs:
            p0.runs[0].text = "コスト\n(千円/人)"
            for r in list(p0.runs)[1:]:
                r._r.getparent().remove(r._r)
        print("  ✓ Y axis label simplified")

# X axis label simplified consistently: "→ 行数/人" → "行数 (行/人)"
for sh in slide.shapes:
    if not sh.has_text_frame: continue
    txt = sh.text_frame.text.strip()
    if txt == "→ 行数/人":
        replace_text_keep(sh, "行数 (行/人)")
        # Also reposition slightly: smaller right-aligned label
        print("  ✓ X axis label simplified")

# (7b) Move team-map title to top-right of plot area
for sh in slide.shapes:
    if not sh.has_text_frame: continue
    txt = sh.text_frame.text.strip()
    if txt == "全社14チームのポジショニング":
        # Original: L=0.30 T=8.05 W=7.67 H=0.22 size 14pt bold
        # New: top-right outside plot, smaller
        sh.left = Inches(3.50)
        sh.top = Inches(8.10)
        sh.width = Inches(4.45)
        sh.height = Inches(0.18)
        for p in sh.text_frame.paragraphs:
            p.alignment = PP_ALIGN.RIGHT
            for r in p.runs:
                r.font.size = Pt(11)
        print("  ✓ Team-map title repositioned to top-right (11pt right-aligned)")

# (7a) Add 5-tick axis numbers
# X axis ticks (below plot bottom edge at PY+PH=9.70): 0, 500, 1000, 1500, 2000
X_TICKS = [0, 500, 1000, 1500, 2000]
TICK_T = 9.71  # just below plot bottom
TICK_H = 0.12
for v in X_TICKS:
    nx = (v - X_MIN) / (X_MAX - X_MIN)
    cx = PX + nx * PW
    add_textbox(slide, str(v), l=cx - 0.20, t=TICK_T, w=0.40, h=TICK_H,
                size=6, color=MUTED, align='center', valign='top')

# Y axis ticks (left of plot edge at PX=1.20): 0, 30, 60, 90, 120
Y_TICKS = [0, 30, 60, 90, 120]
TICK_L = 0.90
TICK_W = 0.28
for v in Y_TICKS:
    ny = (v - Y_MIN) / (Y_MAX - Y_MIN)
    cy = PY + (1 - ny) * PH
    add_textbox(slide, str(v), l=TICK_L, t=cy - 0.06, w=TICK_W, h=TICK_H,
                size=6, color=MUTED, align='right', valign='middle')
print("✓ Axis ticks added (X: 0/500/1000/1500/2000, Y: 0/30/60/90/120)")

# Update own-team bubble position to new 5月 data: X=1700, Y=62, MAU=10
# Existing 14 bubbles in shapes (small ovals in 8.55-9.70 band)
# Identify own bubble (amber) — keep but reposition. Or simpler: re-position any
# bubble that's near the old own location (cx≈4.72, cy≈9.50 from 1667,57.5).
def map_xy(xv, yv):
    nx = (xv - X_MIN) / (X_MAX - X_MIN)
    ny = (yv - Y_MIN) / (Y_MAX - Y_MIN)
    return PX + nx*PW, PY + (1-ny)*PH

# Find own-team bubble (amber color, in the lower-right area)
own_bubble = None
for sh in slide.shapes:
    if sh.shape_type != 1: continue   # AUTO_SHAPE
    try:
        t_in = sh.top/914400; l_in = sh.left/914400
        w_in = sh.width/914400; h_in = sh.height/914400
    except: continue
    if not (8.55 <= t_in < 9.70 and 1.20 < l_in < 6.30):
        continue
    if not (0.08 <= w_in <= 0.30 and abs(w_in - h_in) < 0.01):
        continue
    # Check fill color
    try:
        if sh.fill.fore_color.rgb == AMBER:
            own_bubble = sh
            break
    except:
        pass

if own_bubble is not None:
    new_cx, new_cy = map_xy(1700, 62)
    new_dia = 0.10 + (10 - 7) / (25 - 7) * 0.12  # MAU=10 → dia ~0.12
    own_bubble.left = Inches(new_cx - new_dia/2)
    own_bubble.top = Inches(new_cy - new_dia/2)
    own_bubble.width = Inches(new_dia)
    own_bubble.height = Inches(new_dia)
    print(f"  ✓ Own bubble repositioned to (X=1700, Y=62): center ({new_cx:.2f}, {new_cy:.2f})")
    # Update "当チーム" label position to follow
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "当チーム":
            try:
                t_in = sh.top/914400
                if 8.50 <= t_in < 9.80:
                    label_x = new_cx + new_dia/2 + 0.02
                    label_y = new_cy - new_dia/2 - 0.04
                    if label_x + 0.65 > PX + PW + 0.10:
                        label_x = new_cx - new_dia/2 - 0.65
                    sh.left = Inches(label_x)
                    sh.top = Inches(label_y)
                    print(f"  ✓ 当チーム label repositioned")
            except: pass

prs.save(PPTX)
print(f"\n✓ WROTE: {PPTX}")
