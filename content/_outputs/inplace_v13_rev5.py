#!/usr/bin/env python3
"""
V13 rev5 — right top card numbers aligned with 予実表.

Q1期間内なので「9月までの予測」→「6月までの予測 (Q1着地)」に変更。

Target:
  LEFT (DEEP_FOREST):
    '前月までの累積 (実績)'
    '1,200'                    ← was 1,890 (Q1全体実績) → 5月末まで実績
    '[千円]'

  RIGHT (AMBER):
    '6月までの予測コスト'      ← was '9月までの予測コスト'
    '1,950'                    ← was 4,140
    '[千円]'
    '外注費用 1.2人月'         ← was '2.6人月'
    '(1,600 千円/月 換算)'

数値整合：
  累積 = 4月580 + 5月620 = 1,200
  予測 = 1,200 + 6月計画750 = 1,950
  換算 = 1,950 / 1,600 = 1.22 → 1.2人月
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V13.pptx"
JP_FONT, EN_FONT = "メイリオ", "Segoe UI"

def set_ea(run):
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', JP_FONT)

def replace_text_only(shape, new_text):
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        return
    p0.runs[0].text = new_text
    for r in list(p0.runs)[1:]:
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    r0 = p0.runs[0]
    r0.font.name = EN_FONT
    set_ea(r0)

prs = Presentation(PPTX)
slide = prs.slides[0]

REPLACEMENTS = {
    "1,890":                       "1,200",
    "9月までの予測コスト":          "6月までの予測コスト",
    "4,140":                       "1,950",
    "外注費用 2.6人月":             "外注費用 1.2人月",
}

hits = 0
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    try:
        t = sh.top/914400; l = sh.left/914400
    except:
        continue
    # Confine to right top card (T 1.00-2.20, L > 4.95)
    if not (1.00 <= t <= 2.20 and l > 4.95):
        continue
    txt = sh.text_frame.text.strip()
    if txt in REPLACEMENTS:
        new_text = REPLACEMENTS[txt]
        replace_text_only(sh, new_text)
        print(f"  ✓ '{txt}' → '{new_text}'  (T={t:.2f} L={l:.2f})")
        hits += 1

assert hits == 4, f"Expected 4 edits in right top card, got {hits}"

prs.save(PPTX)
print(f"\nWROTE: {PPTX}")
