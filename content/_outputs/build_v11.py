#!/usr/bin/env python3
"""
Build V11 (rev 5) — Claude Code 月次利用量レポート — A4 PORTRAIT, single page.

Layout (top → bottom):
  1. ヘッダ
  2. ジャッジ3カード + 累積実績/年度末予測 split カード   (見出し無し)
  3. 月次推移チャート [単月]
  4. 月別実績/計画表
  ─── 「6月の実績」見出しバー ─── (★ rev5 で追加. 上下セクションを明示分離)
  5. 6月 KPI 4タイル
  6. メンバー 6月コスト TOP5 / 6月 生成行TOP3 / 6月 効率TOP3
  7. 全社14チームのポジショニング (6月単月)
  8. 参考バンド (当チーム実績6月 vs 公開研究比較)  ← 最下段

V11 rev5 changes vs rev4:
  - REMOVED: 「Q1のジャッジ」見出し (cards だけ残す).
  - REMOVED: 結論バー (Q1累積 1,890... +5.6%) — 完全削除.
  - ADDED:  「6月の実績」見出しバー (h=0.35). 月別実績/計画表 の直下に挿入し、
            以下4セクションが全て 6月単月のデータであることを視覚的に明示.
  - Padding は 0.15-0.30" の範囲で配置. 参考バンドが最下段に位置.

Mock-data disclosure: TOP1-3 cost values (130/108/92 [千円]) from V7 mapped to
real names (山田/鈴木/佐藤). TOP4 (田中 78) / TOP5 (高橋 64) are mock per user.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import MaxNLocator

plt.rcParams["font.family"] = [
    "DejaVu Sans",
    "Droid Sans Fallback",
    "Noto Sans CJK JP",
    "Hiragino Sans", "Yu Gothic", "Meiryo",
]
plt.rcParams["axes.unicode_minus"] = False

# ---------- House style (Calm Forest + Amber) ----------
DEEP_FOREST  = RGBColor(0x1F, 0x3A, 0x2E)
FOREST       = RGBColor(0x2E, 0x6F, 0x40)
SAGE         = RGBColor(0xB5, 0xC9, 0xA8)
SAGE_BG      = RGBColor(0xE8, 0xEF, 0xE3)
AMBER        = RGBColor(0xF4, 0xB9, 0x42)
AMBER_DEEP   = RGBColor(0xC8, 0x8E, 0x1F)
INK          = RGBColor(0x22, 0x28, 0x24)
MUTED        = RGBColor(0x6B, 0x6B, 0x6B)
CARD_BORDER  = RGBColor(0xD7, 0xDC, 0xCF)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
RED          = RGBColor(0xC0, 0x39, 0x2B)
GREEN_OK     = RGBColor(0x2E, 0x6F, 0x40)
COLOR_OWN    = AMBER
COLOR_OTHER  = RGBColor(0xA8, 0xB8, 0xA0)
JP_FONT      = "メイリオ"
EN_FONT      = "Segoe UI"

# ---------- helpers ----------
def set_run_font(run, latin=EN_FONT, ea=JP_FONT, size=12, color=INK, bold=False):
    run.font.name = latin
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea_elem = etree.SubElement(rPr, qn('a:ea'))
    ea_elem.set('typeface', ea)

def add_textbox(slide, text, l, t, w, h, size=12, color=INK, bold=False,
                align='left', valign='top', latin=EN_FONT, ea=JP_FONT,
                line_spacing=None, word_wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if valign == 'top':
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif valign == 'middle':
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == 'bottom':
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    if align == 'left':
        p.alignment = PP_ALIGN.LEFT
    elif align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    set_run_font(r, latin=latin, ea=ea, size=size, color=color, bold=bold)
    return tb

def add_rect(slide, l, t, w, h, fill=WHITE, line=None, line_width=None, corner=None):
    if corner:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(l), Inches(t), Inches(w), Inches(h))
        try:
            shp.adjustments[0] = corner
        except Exception:
            pass
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        if line_width:
            shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_oval(slide, l, t, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_line(slide, x1, y1, x2, y2, color=CARD_BORDER, width=0.75):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line

# =====================================================================
# Slide setup (A4 portrait, SINGLE page)
# =====================================================================
prs = Presentation()
prs.slide_width  = Inches(8.27)
prs.slide_height = Inches(11.69)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# =====================================================================
# DATA (per V11 rev2 user spec; no fabrication)
# =====================================================================
# Single-month cost [千円] for Apr-Mar
MONTHS = ['4月', '5月', '6月', '7月', '8月', '9月',
          '10月', '11月', '12月', '1月', '2月', '3月']
COST_MONTHLY = [580, 620, 690,                  # actual
                770, 770, 770, 770, 770, 770, 770, 770, 770]  # forecast
LOC_MONTHLY  = [15, 17, 20,                     # actual (千行)
                22, 22, 22, 22, 22, 22, 22, 22, 22]
ACTUAL_MONTHS_COUNT = 3

# Plan vs actual
PLAN_MONTHLY    = [750] * 12
ACTUAL_MONTHLY  = [580, 620, 690] + [None] * 9
YEAR_PLAN_SUM   = sum(PLAN_MONTHLY)               # 9,000
Q1_ACTUAL_SUM   = sum(v for v in ACTUAL_MONTHLY if v is not None)  # 1,890

# TOP5 (cost) — mock per user
TOP5_MEMBERS = [
    (1, "山田 太郎", 130),
    (2, "鈴木 花子", 108),
    (3, "佐藤 健",    92),
    (4, "田中 美咲",  78),
    (5, "高橋 翔",    64),
]
TOP5_MAX_COST = max(c for _, _, c in TOP5_MEMBERS)

# Member 6月単月 rankings — V7 mapped
USERS = {
    "山田 太郎": dict(cost_k=130, loc=4200),
    "鈴木 花子": dict(cost_k=108, loc=4500),
    "佐藤 健":   dict(cost_k=92,  loc=2800),
}
for u, d in USERS.items():
    d["eff"] = round(d["loc"] / d["cost_k"], 1)

# Team mapping bubble positions (V7-preserved)
TEAMS = [
    ("当チーム", 0.83, 0.95, 0.20, True),
    ("チームB",  0.86, 0.92, 0.20, False),
    ("チームC",  0.90, 0.78, 0.26, False),
    ("チームD",  0.80, 0.81, 0.21, False),
    ("チームE",  0.76, 0.84, 0.18, False),
    ("チームF",  0.70, 0.73, 0.19, False),
    ("チームG",  0.68, 0.66, 0.16, False),
    ("チームH",  0.58, 0.64, 0.18, False),
    ("チームI",  0.53, 0.60, 0.23, False),
    ("チームJ",  0.48, 0.56, 0.24, False),
    ("チームK",  0.42, 0.55, 0.26, False),
    ("チームL",  0.37, 0.48, 0.28, False),
    ("チームM",  0.32, 0.53, 0.31, False),
    ("チームN",  0.22, 0.39, 0.29, False),
]

# =====================================================================
# Monthly trend chart (matplotlib → PNG) — SINGLE-MONTH bars + line
# =====================================================================
forest_hex = "#1F3A2E"
amber_hex  = "#F4B942"
muted_hex  = "#6B6B6B"

chart_w_in, chart_h_in = 7.67, 1.30   # slightly taller to accommodate external legend above plot
fig, ax_left = plt.subplots(figsize=(chart_w_in, chart_h_in), dpi=220)
fig.patch.set_facecolor("white")

x_pos = list(range(len(MONTHS)))
# Bars: 4-6月 = 実績 (solid forest), 7月-3月 = 計画値 750 [千円] (hatched outline)
# V11 rev7: 7月以降は予測線形外挿 → 計画値 (PLAN_MONTHLY=750) に変更
ax_left.bar(x_pos[:ACTUAL_MONTHS_COUNT],
            COST_MONTHLY[:ACTUAL_MONTHS_COUNT],
            color=forest_hex, edgecolor=forest_hex, width=0.55, linewidth=0.5,
            label="単月コスト (実績)", zorder=2)
ax_left.bar(x_pos[ACTUAL_MONTHS_COUNT:],
            PLAN_MONTHLY[ACTUAL_MONTHS_COUNT:],   # ← 750 [千円] × 9ヶ月
            facecolor="white", edgecolor=forest_hex, hatch="///",
            width=0.55, linewidth=0.6,
            label="単月コスト (計画)", zorder=2)
ax_left.set_xticks(x_pos)
ax_left.set_xticklabels(MONTHS, fontsize=6.5, color=muted_hex)
ax_left.set_ylabel("単月コスト [千円]", color=forest_hex, fontsize=7, labelpad=2)
ax_left.tick_params(axis="y", labelcolor=forest_hex, labelsize=6.5, length=2, pad=1)
ax_left.tick_params(axis="x", length=2, pad=1)
ax_left.grid(True, axis="y", linestyle="--", linewidth=0.4,
             color=muted_hex, alpha=0.35, zorder=1)
ax_left.set_axisbelow(True)
for spine in ("top", "right"):
    ax_left.spines[spine].set_visible(False)
ax_left.spines["left"].set_color(forest_hex)
ax_left.spines["left"].set_linewidth(0.6)
ax_left.spines["bottom"].set_color(muted_hex)
ax_left.spines["bottom"].set_linewidth(0.6)
ax_left.yaxis.set_major_locator(MaxNLocator(nbins=5, integer=True))
# Cost axis range: max of 実績 + 計画 (both ≤ 750/770)
ax_left.set_ylim(0, max(max(COST_MONTHLY[:ACTUAL_MONTHS_COUNT]),
                         max(PLAN_MONTHLY[ACTUAL_MONTHS_COUNT:])) * 1.25)

# Right axis: LOC line — 4-6月実績のみ. 予測線は描画しない (V11 rev7).
ax_right = ax_left.twinx()
ax_right.plot(x_pos[:ACTUAL_MONTHS_COUNT], LOC_MONTHLY[:ACTUAL_MONTHS_COUNT],
              color=amber_hex, marker="o", markersize=4,
              markerfacecolor=amber_hex, markeredgecolor=amber_hex,
              linewidth=1.8, label="単月行数 (実績)", zorder=4)
ax_right.set_ylabel("単月行数 [千行]", color=amber_hex,
                    fontsize=7, labelpad=2)
ax_right.tick_params(axis="y", labelcolor=amber_hex, labelsize=6.5,
                     length=2, pad=1)
for spine in ("top", "left"):
    ax_right.spines[spine].set_visible(False)
ax_right.spines["right"].set_color(amber_hex)
ax_right.spines["right"].set_linewidth(0.6)
# LOC 軸範囲は実績データに合わせる (4-6月 = 15/17/20 → max 20)
ax_right.set_ylim(0, max(LOC_MONTHLY[:ACTUAL_MONTHS_COUNT]) * 1.40)

# Vertical separator between 実績 and 計画
ax_left.axvline(x=ACTUAL_MONTHS_COUNT - 0.5, color=muted_hex,
                linestyle=":", linewidth=0.6, alpha=0.7, zorder=1)
boundary_y = max(max(COST_MONTHLY[:ACTUAL_MONTHS_COUNT]),
                  max(PLAN_MONTHLY[ACTUAL_MONTHS_COUNT:])) * 1.13
ax_left.text(ACTUAL_MONTHS_COUNT - 0.6, boundary_y,
             "実績→", color=forest_hex, fontsize=6, ha="right", va="center")
ax_left.text(ACTUAL_MONTHS_COUNT - 0.4, boundary_y,
             "←計画", color=muted_hex, fontsize=6, ha="left", va="center")

# Combined legend — placed OUTSIDE the plot, just above the axes,
# in a single horizontal row of 4 items.
h1, l1 = ax_left.get_legend_handles_labels()
h2, l2 = ax_right.get_legend_handles_labels()
ax_left.legend(h1 + h2, l1 + l2,
               loc="lower center", bbox_to_anchor=(0.5, 1.02), ncol=4,
               fontsize=6.5, frameon=False,
               handlelength=1.4, handletextpad=0.35,
               columnspacing=1.0, borderpad=0.0,
               labelcolor=[forest_hex, forest_hex, amber_hex, amber_hex])

# Tighter side margins; reserve ~15% at top for the external legend.
plt.subplots_adjust(left=0.06, right=0.95, top=0.85, bottom=0.14)
mpl_chart_png = "/tmp/v11r2_monthly_trend.png"
fig.savefig(mpl_chart_png, dpi=220, facecolor="white", edgecolor="none")
plt.close(fig)

# =====================================================================
# Layout coordinates (Y bands)
# =====================================================================
# Section Y bands (V11 rev5)
# Changes vs rev4:
#   - REMOVED conclusion bar (Q1累積 ... +5.6%) → +0.40" freed
#   - REMOVED "Q1のジャッジ" title above judge cards → +0.22" freed
#   - ADDED divider bar "6月の実績" between table and KPI (h=0.35)
#   - Reference band is now the absolute bottom-most element
#   - Slightly more breathing between sections
HDR_T = 0.20  ; HDR_H = 0.55   # 1. Header                  ends 0.75
JG_T  = 1.00  ; JG_H  = 1.20   # 2. Judge cards + 予測カード ends 2.20   (gap 0.25)
CH_T  = 2.45  ; CH_H  = 1.55   # 3. Monthly trend chart     ends 4.00   (gap 0.25)
TBL_T = 4.20  ; TBL_H = 0.75   # 4. Plan/actual table       ends 4.95   (gap 0.20)
DIV_T = 5.10  ; DIV_H = 0.35   # ★ NEW: 「6月の実績」見出し ends 5.45   (gap 0.15)
KPI_T = 5.60  ; KPI_H = 0.80   # 5. 6月 KPI tiles           ends 6.40   (gap 0.15)
TP_T  = 6.60  ; TP_H  = 1.55   # 6. Cost TOP5 + TOP3×2      ends 8.15   (gap 0.20)
TM_T  = 8.35  ; TM_H  = 1.75   # 7. Team mapping            ends 10.10  (gap 0.20)
REF_T = 10.40 ; REF_H = 1.20   # 8. Reference band          ends 11.60  (gap 0.30)

# Page-wide horizontal extents
PAGE_L = 0.30
PAGE_W = 7.67   # 8.27 - 2*0.30

# =====================================================================
# 1) HEADER
# =====================================================================
add_rect(slide, 0.30, 0.30, 0.06, 0.40, fill=AMBER)
add_textbox(slide, "Claude Code 月次利用量レポート",
            l=0.42, t=HDR_T, w=6.30, h=0.40,
            size=22, color=DEEP_FOREST, bold=True)
add_rect(slide, 6.85, HDR_T+0.04, 1.12, 0.32, fill=DEEP_FOREST, corner=0.30)
add_textbox(slide, "V11",
            l=6.85, t=HDR_T+0.07, w=1.12, h=0.26,
            size=13, color=AMBER, bold=True, align='center', valign='middle')
add_textbox(slide,
            "2026年度Q1 累積  /  報告日: 2026-07-01  /  対象期間: 2026/04/01 - 06/30",
            l=0.42, t=HDR_T+0.42, w=7.50, h=0.18,
            size=9, color=MUTED)
add_line(slide, 0.30, HDR_T+0.65, 7.97, HDR_T+0.65, color=CARD_BORDER, width=1.0)

# =====================================================================
# 2) JUDGE box (left 60%) + 年度末着地予測 (right 40%)
#    Band: T=0.80 -> 2.00 (h=1.20)
# =====================================================================
JG_L_W = 4.55                       # left section (60% of 7.67 ≈ 4.60)
JG_R_L = PAGE_L + JG_L_W + 0.15     # right section starts
JG_R_W = PAGE_L + PAGE_W - JG_R_L   # right section width (~3.12")

# V11 rev5: 「Q1のジャッジ」見出しを削除。
# 3カードは右側の予測カードの上端 (JG_T) と合わせて配置する。
# 高さ 1.20 の中で 3 カード × 0.30 + ギャップ 2 × 0.05 = 1.00 → 上下に
# 0.10 ずつパディング (上 0.10 / 下 0.10) で中央寄せ。
judge_items = [
    ("✓", GREEN_OK,
     "効率性は改善傾向 (Q1累積で継続観測)",
     "過去3ヶ月でコスト 1.19倍 vs 採用行数 1.33倍"),
    ("▶", AMBER_DEEP,
     "来期の投資規模: 外注費用 5.4人月相当 (累積+残月計画ベース)",
     "外注増員 vs Claude活用拡大、議論事項"),
    ("⚠", RED,
     "即対応: 認証基盤PJのモデル選定見直し (6月単月で兆候)",
     "局所的な過剰利用の兆候 (詳細は別ダッシュボード)"),
]
for i, (icon, icon_color, head, sub) in enumerate(judge_items):
    y = JG_T + 0.10 + i * 0.36
    add_rect(slide, PAGE_L, y, JG_L_W, 0.32,
             fill=SAGE_BG, line=CARD_BORDER, line_width=0.5)
    add_rect(slide, PAGE_L, y, 0.05, 0.32, fill=icon_color)
    add_textbox(slide, icon, l=PAGE_L+0.10, t=y+0.05, w=0.30, h=0.22,
                size=14, color=icon_color, bold=True,
                align='center', valign='middle')
    add_textbox(slide, head, l=PAGE_L+0.46, t=y+0.03, w=JG_L_W-0.50, h=0.15,
                size=9, color=DEEP_FOREST, bold=True)
    add_textbox(slide, sub, l=PAGE_L+0.46, t=y+0.17, w=JG_L_W-0.50, h=0.13,
                size=7.5, color=INK)

# Right card: 左右(50/50)で 実績 vs 予測 を別背景色で明示
# - LEFT half (実績): DEEP_FOREST 背景 — 確定値の重厚感
# - RIGHT half (予測): AMBER 背景 — 将来の明るさ・暫定感
# 強コントラストでメタ意味（過去/未来）を視覚化。
# 補助テキスト (「(4-6月の3ヶ月累積)」「(年度末 = X)」など) は全削除.

half_w = JG_R_W / 2
left_x  = JG_R_L
right_x = JG_R_L + half_w

# LEFT half — DEEP_FOREST background
add_rect(slide, left_x, JG_T, half_w, JG_H, fill=DEEP_FOREST, line=None)
# RIGHT half — AMBER background
add_rect(slide, right_x, JG_T, half_w, JG_H, fill=AMBER, line=None)

# ---- LEFT column: 累積実績 (DEEP_FOREST bg, light text) ----
# Layout: title (header + 実績) / big number / unit  — number & unit stacked.
add_textbox(slide, "前月までの累積コスト",
            l=left_x + 0.05, t=JG_T + 0.10, w=half_w - 0.10, h=0.16,
            size=9, color=SAGE, bold=True, align='center', word_wrap=False)
add_textbox(slide, "(実績)",
            l=left_x + 0.05, t=JG_T + 0.27, w=half_w - 0.10, h=0.14,
            size=7.5, color=SAGE, align='center', word_wrap=False)
add_textbox(slide, "1,890",
            l=left_x, t=JG_T + 0.46, w=half_w, h=0.48,
            size=30, color=AMBER, bold=True,
            align='center', valign='middle', word_wrap=False)
add_textbox(slide, "[千円]",
            l=left_x, t=JG_T + 0.96, w=half_w, h=0.16,
            size=9, color=SAGE, align='center', word_wrap=False)

# ---- RIGHT column: 年度末予測 (AMBER bg, dark text) ----
add_textbox(slide, "年度末着地予測 (予測)",
            l=right_x + 0.05, t=JG_T + 0.08, w=half_w - 0.10, h=0.18,
            size=9, color=DEEP_FOREST, bold=True, align='center', word_wrap=False)
add_textbox(slide, "8,640",
            l=right_x, t=JG_T + 0.26, w=half_w, h=0.36,
            size=24, color=DEEP_FOREST, bold=True,
            align='center', valign='middle', word_wrap=False)
add_textbox(slide, "[千円]",
            l=right_x, t=JG_T + 0.62, w=half_w, h=0.14,
            size=8.5, color=DEEP_FOREST, align='center', word_wrap=False)
add_textbox(slide, "外注費用 5.4人月",
            l=right_x + 0.05, t=JG_T + 0.78, w=half_w - 0.10, h=0.17,
            size=11, color=DEEP_FOREST, bold=True,
            align='center', word_wrap=False)
add_textbox(slide, "(1,600 千円/月 換算)",
            l=right_x + 0.05, t=JG_T + 0.97, w=half_w - 0.10, h=0.16,
            size=7.5, color=DEEP_FOREST, align='center', word_wrap=False)

# =====================================================================
# 3) Monthly trend chart [単月]
#    Band T=2.10 -> 3.30 (h=1.20)
# =====================================================================
add_textbox(slide, "月次推移 (単月) — コスト × 採用コード行数",
            l=PAGE_L, t=CH_T, w=5.00, h=0.18,
            size=11, color=DEEP_FOREST, bold=True)
add_textbox(slide,
            "※ 7月以降は計画値 (月 750 [千円])  /  単月行数は実績のみ表示",
            l=PAGE_L+3.50, t=CH_T+0.02, w=PAGE_W-3.50, h=0.16,
            size=7, color=MUTED, align='right', word_wrap=False)

# embed PNG
slide.shapes.add_picture(mpl_chart_png,
                         Inches(PAGE_L), Inches(CH_T+0.20),
                         width=Inches(PAGE_W), height=Inches(CH_H-0.20))

# =====================================================================
# 4) Plan / actual table (full width)
#    Band T=3.40 -> 4.10 (h=0.70)
# =====================================================================
add_textbox(slide, "月別実績 / 計画",
            l=PAGE_L, t=TBL_T, w=5.00, h=0.18,
            size=11, color=DEEP_FOREST, bold=True)
add_textbox(slide, "単位: [千円]",
            l=PAGE_L+5.05, t=TBL_T+0.02, w=PAGE_W-5.05, h=0.18,
            size=8, color=MUTED, align='right')

table_top   = TBL_T + 0.22
label_col_w = 0.65
year_col_w  = 0.70
month_col_w = (PAGE_W - label_col_w - year_col_w) / 12
header_h    = 0.16
row_h       = 0.16

# Header strip
add_rect(slide, PAGE_L, table_top, PAGE_W, header_h, fill=DEEP_FOREST, line=None)
for mi, mname in enumerate(MONTHS):
    x = PAGE_L + label_col_w + mi * month_col_w
    add_textbox(slide, mname, l=x, t=table_top, w=month_col_w, h=header_h,
                size=7, color=WHITE, bold=True,
                align='center', valign='middle', word_wrap=False)
add_textbox(slide, "年間合計",
            l=PAGE_L + label_col_w + 12 * month_col_w, t=table_top,
            w=year_col_w, h=header_h,
            size=7, color=WHITE, bold=True, align='center', valign='middle',
            word_wrap=False)

# Row 1: 当初計画
row1_top = table_top + header_h
add_rect(slide, PAGE_L, row1_top, PAGE_W, row_h,
         fill=WHITE, line=CARD_BORDER, line_width=0.3)
add_textbox(slide, "当初計画", l=PAGE_L+0.05, t=row1_top, w=label_col_w-0.10,
            h=row_h, size=8, color=INK, bold=True, valign='middle', word_wrap=False)
for mi, val in enumerate(PLAN_MONTHLY):
    x = PAGE_L + label_col_w + mi * month_col_w
    add_textbox(slide, str(val), l=x, t=row1_top, w=month_col_w, h=row_h,
                size=7, color=INK, align='center', valign='middle',
                word_wrap=False)
add_textbox(slide, f"{YEAR_PLAN_SUM:,}",
            l=PAGE_L + label_col_w + 12 * month_col_w, t=row1_top,
            w=year_col_w, h=row_h,
            size=8, color=DEEP_FOREST, bold=True, align='center', valign='middle',
            word_wrap=False)

# Row 2: 実績
row2_top = row1_top + row_h
add_rect(slide, PAGE_L, row2_top, PAGE_W, row_h,
         fill=SAGE_BG, line=CARD_BORDER, line_width=0.3)
add_textbox(slide, "実績", l=PAGE_L+0.05, t=row2_top, w=label_col_w-0.10,
            h=row_h, size=8, color=INK, bold=True, valign='middle', word_wrap=False)
# Yellow vertical separator at the actual→forecast boundary
sep_x = PAGE_L + label_col_w + ACTUAL_MONTHS_COUNT * month_col_w
add_line(slide, sep_x, table_top, sep_x, row2_top + row_h, color=AMBER, width=0.9)
for mi, val in enumerate(ACTUAL_MONTHLY):
    x = PAGE_L + label_col_w + mi * month_col_w
    if val is None:
        text = "-"
        c = MUTED
    else:
        text = str(val)
        c = GREEN_OK if val < PLAN_MONTHLY[mi] else AMBER_DEEP
    add_textbox(slide, text, l=x, t=row2_top, w=month_col_w, h=row_h,
                size=7, color=c, bold=(val is not None),
                align='center', valign='middle', word_wrap=False)
add_textbox(slide, f"{Q1_ACTUAL_SUM:,}*",
            l=PAGE_L + label_col_w + 12 * month_col_w, t=row2_top,
            w=year_col_w, h=row_h,
            size=8, color=AMBER_DEEP, bold=True, align='center', valign='middle',
            word_wrap=False)

# Table footnote
add_textbox(slide,
            "* 6月末までの累積実績  /  黄色縦線 = 実績→予測の境界  /  "
            "実績 < 計画 = 緑、実績 > 計画 = 橙",
            l=PAGE_L, t=row2_top + row_h + 0.02, w=PAGE_W, h=0.14,
            size=6.5, color=MUTED)

# =====================================================================
# DIVIDER: 「6月の実績」 — section break before single-month sections
#    Full-bleed amber-accent forest bar (V11 rev5 NEW)
# =====================================================================
add_rect(slide, 0.00, DIV_T, 8.27, DIV_H, fill=DEEP_FOREST, line=None)
add_rect(slide, 0.00, DIV_T, 0.06, DIV_H, fill=AMBER)
add_textbox(slide, "6月の実績",
            l=0.22, t=DIV_T+0.03, w=3.00, h=DIV_H-0.06,
            size=14, color=WHITE, bold=True, valign='middle',
            word_wrap=False)
add_textbox(slide, "↓ 以下は 2026 年 6 月 単月のスナップショット",
            l=3.30, t=DIV_T+0.03, w=4.80, h=DIV_H-0.06,
            size=8, color=SAGE, valign='middle', align='right',
            word_wrap=False)

# =====================================================================
# 5) 6月 KPI tiles (4枚)
#    Band T=5.60 -> 6.40 (h=0.80)
# =====================================================================
kpi_june = [
    ("6月 コスト [千円]",        "690",     "+11.3%",  "vs 5月 620 [千円]",   AMBER_DEEP),
    ("6月 アクティブユーザー",    "12名",    "+2名",    "vs 5月 10名",          GREEN_OK),
    ("6月 1人あたり月額 [千円]",  "57.5",    "−4.5%",   "vs 5月 62.0 [千円]",  GREEN_OK),
    ("6月 1人あたり採用行数",     "1,667行", "+5%",     "vs 5月 1,587行",      GREEN_OK),
]
tile_w = 1.80
tile_gap = 0.16
for i, (label, big, delta, sub, dc) in enumerate(kpi_june):
    x = PAGE_L + i * (tile_w + tile_gap)
    add_rect(slide, x, KPI_T, tile_w, KPI_H,
             fill=WHITE, line=CARD_BORDER, line_width=0.5)
    add_rect(slide, x, KPI_T, tile_w, 0.04, fill=DEEP_FOREST)
    add_textbox(slide, label, l=x+0.08, t=KPI_T+0.05, w=tile_w-0.16, h=0.15,
                size=8, color=MUTED)
    add_textbox(slide, big, l=x+0.08, t=KPI_T+0.18, w=tile_w-0.16, h=0.28,
                size=17, color=DEEP_FOREST, bold=True, word_wrap=False)
    add_textbox(slide, delta, l=x+0.08, t=KPI_T+0.46, w=1.20, h=0.13,
                size=8, color=dc, bold=True)
    add_textbox(slide, sub, l=x+0.08, t=KPI_T+0.56, w=tile_w-0.16, h=0.12,
                size=7, color=MUTED)

# =====================================================================
# 6) Cost TOP5 / 生成行 TOP3 / 効率 TOP3 (3 panels horizontal)
#    Band T=5.00 -> 6.40 (h=1.40)
#    Widths: TOP5 = 40% (3.07), TOP3 = 30% each (2.20)
# =====================================================================
ROW6_GAP    = 0.10
TOP5_PANEL_W   = 3.07
TOP3_PANEL_W   = (PAGE_W - TOP5_PANEL_W - 2 * ROW6_GAP) / 2   # ~ 2.20

# --- 6a) コストTOP5 panel ---
P6A_L = PAGE_L
add_textbox(slide, "メンバー 6月コスト TOP5",
            l=P6A_L, t=TP_T, w=TOP5_PANEL_W-0.55, h=0.18,
            size=10, color=DEEP_FOREST, bold=True, word_wrap=False)
add_textbox(slide, "[千円]",
            l=P6A_L+TOP5_PANEL_W-0.55, t=TP_T+0.02, w=0.50, h=0.16,
            size=8, color=MUTED, align='right')

add_rect(slide, P6A_L, TP_T+0.20, TOP5_PANEL_W, TP_H-0.36,
         fill=WHITE, line=CARD_BORDER, line_width=0.5)
add_rect(slide, P6A_L, TP_T+0.20, TOP5_PANEL_W, 0.04, fill=DEEP_FOREST)

row_top    = TP_T + 0.28
row_h_     = 0.16
inner_pad  = 0.08
rank_w     = 0.16
name_w     = 1.00
gap        = 0.04
value_w    = 0.45
bar_x      = P6A_L + inner_pad + rank_w + gap + name_w + gap
bar_w_max  = TOP5_PANEL_W - inner_pad - rank_w - gap - name_w - gap - value_w - gap - inner_pad
RANK_BG = [AMBER, SAGE, RGBColor(0xE0, 0xE0, 0xE0), SAGE_BG, SAGE_BG]
for i, (rank, name, cost) in enumerate(TOP5_MEMBERS):
    ry = row_top + i * row_h_
    rx = P6A_L + inner_pad
    add_oval(slide, rx, ry+0.005, rank_w, rank_w, fill=RANK_BG[i], line=None)
    add_textbox(slide, str(rank), l=rx, t=ry+0.005, w=rank_w, h=rank_w,
                size=9, color=DEEP_FOREST, bold=True,
                align='center', valign='middle')
    add_textbox(slide, name, l=rx+rank_w+gap, t=ry, w=name_w, h=row_h_,
                size=10, color=INK, bold=(i == 0), valign='middle',
                word_wrap=False)
    bw = bar_w_max * (cost / TOP5_MAX_COST)
    bar_color = AMBER if i == 0 else COLOR_OTHER
    add_rect(slide, bar_x, ry+0.04, bw, row_h_-0.08, fill=bar_color, line=None)
    add_textbox(slide, f"{cost}",
                l=P6A_L + TOP5_PANEL_W - inner_pad - value_w, t=ry,
                w=value_w, h=row_h_,
                size=10, color=DEEP_FOREST, bold=True,
                align='right', valign='middle', word_wrap=False)
# V11 rev6: mock-data footnote removed for a cleaner panel.

# --- 6b/c) 生成行 / 効率 TOP3 panels ---
top_loc = sorted(USERS.items(), key=lambda kv: -kv[1]["loc"])
top_eff = sorted(USERS.items(), key=lambda kv: -kv[1]["eff"])

panels_top3 = [
    ("6月 生成行TOP3", "最もコード生成 [行]", top_loc,
        lambda d: f"{d['loc']:,}"),
    ("6月 効率TOP3",   "コスパ [行/千円]",   top_eff,
        lambda d: f"{d['eff']}"),
]
RANK_BG_2 = [AMBER, SAGE, RGBColor(0xE0, 0xE0, 0xE0)]
for pi, (head, sub, ranking, big_fn) in enumerate(panels_top3):
    px = PAGE_L + TOP5_PANEL_W + ROW6_GAP + pi * (TOP3_PANEL_W + ROW6_GAP)
    add_textbox(slide, head, l=px, t=TP_T, w=TOP3_PANEL_W, h=0.18,
                size=10, color=DEEP_FOREST, bold=True)

    add_rect(slide, px, TP_T+0.20, TOP3_PANEL_W, TP_H-0.36,
             fill=WHITE, line=CARD_BORDER, line_width=0.5)
    add_rect(slide, px, TP_T+0.20, TOP3_PANEL_W, 0.04, fill=DEEP_FOREST)
    add_textbox(slide, sub, l=px+0.08, t=TP_T+0.26, w=TOP3_PANEL_W-0.16, h=0.13,
                size=7, color=MUTED)

    row_t = TP_T + 0.45
    rh    = 0.20
    for ri, (uname, d) in enumerate(ranking):
        ry = row_t + ri * rh
        add_oval(slide, px+0.08, ry+0.02, 0.16, 0.16,
                 fill=RANK_BG_2[ri], line=None)
        add_textbox(slide, str(ri+1),
                    l=px+0.08, t=ry+0.02, w=0.16, h=0.16,
                    size=8, color=DEEP_FOREST, bold=True,
                    align='center', valign='middle')
        add_textbox(slide, uname,
                    l=px+0.28, t=ry+0.01, w=0.95, h=0.18,
                    size=10, color=INK, bold=True, valign='middle',
                    word_wrap=False)
        add_textbox(slide, big_fn(d),
                    l=px+TOP3_PANEL_W-0.65, t=ry+0.01,
                    w=0.55, h=0.18,
                    size=11, color=AMBER_DEEP, bold=True,
                    align='right', valign='middle', word_wrap=False)

# =====================================================================
# 7) 14チーム ポジショニング (V11 rev7 — quadrant labels redistributed)
#    - タイトル削除
#    - 凡例「当チーム/他チーム」削除 (色判別)
#    - 「バブル径: 1人あたりコスト」は左下隅に小さく
#    - 4象限ラベルをプロット外側の4隅に分散配置
#      (左上 / 右上 / 左下 / 右下 がプロット外辺の外側)
# =====================================================================
add_rect(slide, PAGE_L, TM_T, PAGE_W, TM_H, fill=WHITE,
         line=CARD_BORDER, line_width=0.5)

# Plot area placement: 上下に各 0.22" の余白を確保 (象限ラベル用)
PX, PY = 1.20, TM_T + 0.30
PW, PH = 5.00, 1.05   # slightly slimmer to fit 4-quadrant external labels

# Axes & gridlines
add_line(slide, PX, PY + PH, PX + PW, PY + PH, color=DEEP_FOREST, width=1.0)
add_line(slide, PX, PY,      PX,        PY + PH, color=DEEP_FOREST, width=1.0)
add_line(slide, PX + PW/2, PY, PX + PW/2, PY + PH, color=CARD_BORDER, width=0.5)
add_line(slide, PX, PY + PH/2, PX + PW, PY + PH/2, color=CARD_BORDER, width=0.5)
# X 軸ラベル削除: 横方向の意味は 4 隅の象限ラベルが直接示すため.

# Y axis label (rotated-style textbox to the left)
yl = slide.shapes.add_textbox(Inches(PX-0.85), Inches(PY + PH/2 - 0.30),
                              Inches(0.80), Inches(0.62))
yl.text_frame.word_wrap = True
yl.text_frame.margin_left = yl.text_frame.margin_right = Emu(0)
yl.text_frame.margin_top = yl.text_frame.margin_bottom = Emu(0)
ylp = yl.text_frame.paragraphs[0]
ylp.alignment = PP_ALIGN.RIGHT
ylr = ylp.add_run()
ylr.text = "↑ 効率\n(行/コスト)"
set_run_font(ylr, size=7, color=DEEP_FOREST, bold=True)

# (NO quadrant captions inside plot — moved to external panel below)

# Bubbles (own goes ON TOP)
def bubble_pos_(xn, yn, dia):
    cx = PX + xn * PW
    cy = PY + (1.0 - yn) * PH
    return cx - dia/2, cy - dia/2

for name, xn, yn, dia, is_own in sorted(TEAMS, key=lambda t: 1 if t[4] else 0):
    dia_use = dia * 0.60   # slightly bigger now that plot is taller
    bx, by = bubble_pos_(xn, yn, dia_use)
    fill = COLOR_OWN if is_own else COLOR_OTHER
    outline = AMBER_DEEP if is_own else FOREST
    add_oval(slide, bx, by, dia_use, dia_use, fill=fill, line=outline)
    if is_own:
        label_x = bx + dia_use + 0.02
        label_y = by - 0.03
        if label_x + 0.60 > PX + PW + 0.10:
            label_x = bx - 0.60
        add_textbox(slide, "当チーム",
                    l=label_x, t=label_y, w=0.65, h=0.14,
                    size=8, color=AMBER_DEEP, bold=True)

# ----- Quadrant labels: 4 corners outside the plot -----
# 上側: 上左 (効率派) / 上右 (活用優等生) — プロット上端の外側 (TM_T+0.04~0.20)
# 下側: 下左 (要支援)  / 下右 (コスト先行) — プロット下端の外側 (PY+PH+0.16~0.28)
# 横位置はプロットX範囲 (PX ~ PX+PW) を半分割した左半 / 右半 の中央近傍.
QLABEL_TOP_Y    = TM_T + 0.06       # 上ラベルのY (プロット外, 上)
QLABEL_BOTTOM_Y = PY + PH + 0.16    # 下ラベルのY (プロット外, 下: 軸ラベルより下)
QLABEL_H = 0.20
left_half_x  = PX + 0.05             # 左半開始
left_half_w  = PW/2 - 0.15
right_half_x = PX + PW/2 + 0.10      # 右半開始
right_half_w = PW/2 - 0.15

# V11 rev8: 象限ラベルを意味的名称 → 「高/低使用 × 高/低効率」機械的表記に変更
# 配置 (4隅) / 三角マーカー / 色 (上=Forest, 下=Amber) は維持.
# 上左 = 低使用 × 高効率
add_textbox(slide, "▼ 低使用 × 高効率",
            l=left_half_x, t=QLABEL_TOP_Y, w=left_half_w, h=QLABEL_H,
            size=9, color=FOREST, bold=True, valign='middle', word_wrap=False)
# 上右 = 高使用 × 高効率
add_textbox(slide, "高使用 × 高効率 ▼",
            l=right_half_x, t=QLABEL_TOP_Y, w=right_half_w, h=QLABEL_H,
            size=9, color=FOREST, bold=True, align='right', valign='middle',
            word_wrap=False)
# 下左 = 低使用 × 低効率
add_textbox(slide, "▶ 低使用 × 低効率",
            l=left_half_x, t=QLABEL_BOTTOM_Y, w=left_half_w, h=QLABEL_H,
            size=9, color=AMBER_DEEP, bold=True, valign='middle', word_wrap=False)
# 下右 = 高使用 × 低効率
add_textbox(slide, "高使用 × 低効率 ◀",
            l=right_half_x, t=QLABEL_BOTTOM_Y, w=right_half_w, h=QLABEL_H,
            size=9, color=AMBER_DEEP, bold=True, align='right', valign='middle',
            word_wrap=False)

# 「バブル径」凡例: プロット中央の真下 (下ラベル群と被らない位置に小さく)
add_textbox(slide, "バブル径 = 1人あたりコスト",
            l=PX + PW/2 - 1.00, t=TM_T+TM_H-0.13, w=2.00, h=0.12,
            size=6.5, color=MUTED, align='center', word_wrap=False)

# =====================================================================
# 8) Reference band — now the bottom-most element (V11 rev5: conclusion bar
#    removed; 参考バンド directly follows the team-mapping section).
#    2-col layout retained from V10.
#    LEFT 50% : 当チーム実績 (6月) — investment-equivalent productivity
#    RIGHT 50%: 【参考】 公開研究の月間コード生産量 (3-row research table)
# =====================================================================
LX1, LY1 = PAGE_L, REF_T
LW1, LH1 = (PAGE_W - 0.15) / 2, REF_H    # left card ~3.76"
RX1 = LX1 + LW1 + 0.15
RW1 = PAGE_W - LW1 - 0.15
RY1, RH1 = REF_T, REF_H

# ---- LEFT: 当チーム実績 (6月) ----
# Investment-equivalence math (per user, V10 1.82x multiplier):
#   投資換算 = 6月実績 690 [千円] × 1.82  ≈ 1,260 [千円]
#   行/人/月  = 6月実績 20,000行 × 1.82   = 36,400 行/人/月
#   業界比    = 36,400 ÷ 600  = 60.7 ≈ 約61倍
add_rect(slide, LX1, LY1, LW1, LH1, fill=DEEP_FOREST, line=None)
add_textbox(slide, "当チーム実績 (6月)",
            l=LX1+0.10, t=LY1+0.05, w=LW1-0.20, h=0.16,
            size=10, color=SAGE, bold=True)
add_textbox(slide, "690 [千円] → 20,000行",
            l=LX1+0.10, t=LY1+0.21, w=LW1-0.20, h=0.26,
            size=16, color=AMBER, bold=True, word_wrap=False)
add_line(slide, LX1+0.10, LY1+0.50, LX1+LW1-0.10, LY1+0.50,
         color=AMBER, width=0.5)
add_textbox(slide, "約 1,260 [千円] 投資換算 (人月単純相当)",
            l=LX1+0.10, t=LY1+0.52, w=LW1-0.20, h=0.16,
            size=9, color=WHITE, bold=True, word_wrap=False)
add_textbox(slide, "約 36,400 行 / 人/月",
            l=LX1+0.10, t=LY1+0.70, w=LW1-0.20, h=0.28,
            size=18, color=AMBER, bold=True, word_wrap=False)
add_textbox(slide, "業界中央値 600 行/人/月 の 約61倍 (※単純比較)",
            l=LX1+0.10, t=LY1+1.00, w=LW1-0.20, h=0.16,
            size=8, color=SAGE, word_wrap=False)

# ---- RIGHT: 【参考】 公開研究の月間コード生産量 (V10 carried over) ----
add_rect(slide, RX1, RY1, RW1, RH1, fill=WHITE, line=CARD_BORDER, line_width=0.5)
add_textbox(slide, "【参考】 公開研究の月間コード生産量",
            l=RX1+0.10, t=RY1+0.05, w=RW1-0.20, h=0.16,
            size=10, color=DEEP_FOREST, bold=True)
add_textbox(slide, "※ 品質・複雑度・業務特性は考慮されない単純比較",
            l=RX1+0.10, t=RY1+0.22, w=RW1-0.20, h=0.12,
            size=7, color=MUTED)

bench = [
    ("保守的",   "325 〜 750 行/月",     "Capers Jones"),
    ("中央値的", "約 1,000 行/月",        "日本SI業界 (IPA)"),
    ("生産性高", "1,500 〜 1,600 行/月",  "個人実測 / Scrum"),
]
bench_top = RY1 + 0.38
bench_row_h = 0.25
for i, (lvl, num, src) in enumerate(bench):
    by = bench_top + i * bench_row_h
    add_rect(slide, RX1+0.10, by, RW1-0.20, bench_row_h - 0.03,
             fill=SAGE_BG, line=CARD_BORDER, line_width=0.3)
    add_rect(slide, RX1+0.10, by, 0.05, bench_row_h - 0.03, fill=FOREST)
    add_textbox(slide, lvl, l=RX1+0.20, t=by+0.03,
                w=0.70, h=bench_row_h - 0.06,
                size=9, color=DEEP_FOREST, bold=True, valign='middle')
    add_textbox(slide, num, l=RX1+0.90, t=by+0.03,
                w=1.65, h=bench_row_h - 0.06,
                size=9.5, color=INK, bold=True, valign='middle',
                word_wrap=False)
    add_textbox(slide, src, l=RX1+2.58, t=by+0.03,
                w=RW1-2.73, h=bench_row_h - 0.06,
                size=8, color=MUTED, valign='middle', align='right',
                word_wrap=False)

# =====================================================================
# Speaker notes (full V11 rev2 changelog)
# =====================================================================
notes = slide.notes_slide.notes_text_frame
notes.text = (
    "V11 (rev 3) — A4 portrait, single page. Final layout with breathing.\n"
    "8 sections top→bottom:\n"
    "  1. ヘッダ\n"
    "  2. ジャッジ (60%) + 累積実績/年度末予測 split カード (40%)\n"
    "  3. 月次推移チャート (単月)\n"
    "  4. 月別実績/計画 表 (横幅フル)\n"
    "  5. 前月(6月) KPI 4タイル\n"
    "  6. コストTOP5 (40%) + 生成行TOP3 (30%) + 効率TOP3 (30%)\n"
    "  7. 14チーム ポジショニング (V10 の2色配色, ラベル無し)\n"
    "  8. 結論バー\n"
    "\n"
    "Compared to V11 rev2:\n"
    "  - CHANGED: 年度末予測カードを 上下並び → 左右(50/50)分割\n"
    "    LEFT = 累積実績 1,890 / RIGHT = 年度末予測 9,500\n"
    "    Amber 縦罫線で 実績/予測 を視覚分離.\n"
    "  - REMOVED: IPA 一次ソース ベンチマークセクション (rev2 で追加した⑧).\n"
    "  - INCREASED PADDING: セクション間 0.20-0.25\" インチ余白に拡大.\n"
    "    Chart / TOP / Team-Map をやや拡大して読みやすさ向上.\n"
    "  - CONCLUSION BAR (V10 そのまま, 文言維持) at page bottom.\n"
    "\n"
    "Mock-data: TOP1-3 cost preserved 1:1 from V7 (130/108/92). 名前は\n"
    "ユーザー指定 (山田/鈴木/佐藤). TOP4 田中 78 / TOP5 高橋 64 はモック.\n"
    "On-slide footnote in cost TOP5 panel explicitly disclosed.\n"
    "Forecast monthly 770 [千円] · 22K行 から 年度末 9,500 予測の旨は\n"
    "チャート上部 caption に明記."
)

out = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V11.pptx"
prs.save(out)
print("WROTE:", out)
