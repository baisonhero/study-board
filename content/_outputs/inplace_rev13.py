#!/usr/bin/env python3
"""
V11 rev13 — in-place: move X-axis label「→ 使用率」from plot-center to plot
right-edge (where the arrow naturally points).

Plot geometry: PX=1.20, PW=5.00 → right edge x=6.20.
New box:  L=5.65, W=0.55 (ends at 6.20), T=9.70, H=0.14.
          right-aligned so the arrow→text hugs the right end.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V11.pptx"

prs = Presentation(PPTX)
slide = prs.slides[0]

n = 0
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    try:
        t_in = sh.top / 914400
    except Exception:
        continue
    # Target band: X-axis label Y range (9.60-9.85) + has "使用率"
    if not (9.60 <= t_in < 9.85):
        continue
    if "使用率" not in sh.text_frame.text:
        continue
    # Reposition to plot right edge
    sh.left   = Inches(5.65)
    sh.width  = Inches(0.55)
    sh.top    = Inches(9.70)
    sh.height = Inches(0.14)
    # Right-align text inside
    for para in sh.text_frame.paragraphs:
        para.alignment = PP_ALIGN.RIGHT
    n += 1
    print(f"Moved X-axis label → L=5.65 T=9.70 W=0.55 H=0.14 (right-aligned)")

assert n == 1, f"Expected 1 update, got {n}"
prs.save(PPTX)
print(f"WROTE: {PPTX}")
