#!/usr/bin/env python3
"""V13_tmp contrast patch. Applies content/size/position-aware rules to
restore legibility while keeping CORAL only on the largest numbers / main title /
accent bands.

Color & weight rules
--------------------
- Body / small text → ANTHRACITE
- CORAL kept only for: main title, KPI/right-top big numbers (≥18pt), 1 accent marker
- WARM_GRAY -> ANTHRACITE wherever it sits on the cream background
- KPI / TOP5 / 参考バンドの見出しラベル → ANTHRACITE bold
- TOP5 氏名・数値 → ANTHRACITE
- LEFT-top card sub-label: stay WHITE but bold + sz +1pt
- RIGHT-top card sub-labels → ANTHRACITE (Coral fights with SAND)
- Bottom 当チーム実績 card (CORAL_DARK 地): all text → WHITE
- 月次推移グラフ / ポジショニンググラフの軸 / 凡例 → ANTHRACITE
- ヘッダー副題 (チーム名行 + 日付行) → ANTHRACITE bold
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
import os

SRC = '2026-07-01_Claude_monthly_report_V13_tmp.pptx'

ANTHRACITE = RGBColor(0x19, 0x19, 0x19)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CORAL      = RGBColor(0xCC, 0x78, 0x5C)
BRICK_RED  = RGBColor(0xA8, 0x50, 0x3D)

def hex_of(run):
    try: return str(run.font.color.rgb)
    except: return ''

def size_of(run):
    return run.font.size.pt if run.font.size else 0

def set_color(run, rgb):
    run.font.color.rgb = rgb

def set_bold(run, b=True):
    run.font.bold = b

def set_size(run, pt):
    run.font.size = Pt(pt)

def patch_run(sh, p, r, l, t):
    """Return reason string if modified."""
    text = r.text
    size = size_of(r)
    color = hex_of(r)
    notes = []

    # 1) Bottom-left 当チーム実績 card (CORAL_DARK bg) — ALL text WHITE
    if 0.39 <= l <= 0.41 and 10.40 <= t <= 11.60:
        if color != 'FFFFFF':
            set_color(r, WHITE); notes.append('team_card→WHITE')
        return ','.join(notes) or None

    # 2) Right-top LEFT card sub-label "前月までの累積 (実績)" — bump sz 9→10
    if 4.99 <= l <= 5.06 and 1.03 <= t <= 1.06:
        if size and size < 10.5:
            set_size(r, size + 1); notes.append('left_card_label_sz+1')
        return ','.join(notes) or None

    # 3) Right-top RIGHT card — keep big "1,950" CORAL; rest ANTHRACITE
    if 6.45 <= l <= 6.55 and 1.00 <= t <= 1.80:
        if size >= 20:
            return None  # big number — keep
        if color != '191919':
            set_color(r, ANTHRACITE); notes.append('right_card→ANTHRACITE')
        return ','.join(notes) or None

    # 4) Callout headlines (l~0.76, t 1.10-2.00, sz=9 CORAL bold) → ANTHRACITE bold
    if 0.74 <= l <= 0.78 and 1.10 <= t <= 2.00 and size and size <= 9.5:
        if color == 'CC785C':
            set_color(r, ANTHRACITE); notes.append('callout→ANTHRACITE')
        return ','.join(notes) or None

    # 5) Header subtitle "2026年度Q1..." (sz=9) → ANTHRACITE bold (force bold)
    if 0.40 <= l <= 0.45 and 0.73 <= t <= 0.76:
        if r.font.bold is not True:
            set_bold(r, True); notes.append('date_bold')
        if color != '191919':
            set_color(r, ANTHRACITE); notes.append('date→ANTHRACITE')
        return ','.join(notes) or None
    # 5b) "認証基盤プロジェクト" — force bold (already ANTHRACITE+bold but ensure)
    if 0.40 <= l <= 0.45 and 0.55 <= t <= 0.58:
        set_bold(r, True)
        return None

    # 6) Section band subtitle "↓ 以下は..." (l=3.30, t=5.13) SAND → WHITE
    if 3.28 <= l <= 3.34 and 5.10 <= t <= 5.20:
        set_color(r, WHITE); return 'section_sub→WHITE'

    # 7) Monthly chart heading & 月別実績/計画 heading (sz=11 CORAL) → ANTHRACITE
    if (0.29 <= l <= 0.31) and (2.43 <= t <= 2.47 or 4.18 <= t <= 4.22) and color == 'CC785C':
        set_color(r, ANTHRACITE); return 'sub_heading→ANTHRACITE'

    # 8) 全社14チームのポジショニング heading (l=3.50, t=8.10, sz=11 CORAL) → ANTHRACITE
    if 3.49 <= l <= 3.51 and 8.09 <= t <= 8.11 and color == 'CC785C':
        set_color(r, ANTHRACITE); return 'pos_heading→ANTHRACITE'

    # 9) Predicted total "2,250" in table (l=6.19, t=4.58, CC785C sz=8) → ANTHRACITE
    if 6.18 <= l <= 6.21 and 4.55 <= t <= 4.61 and color == 'CC785C':
        set_color(r, ANTHRACITE); return 'pred_total→ANTHRACITE'

    # 10) Table 単位 / dash / "-" + 消費率 dash (WARM_GRAY) → ANTHRACITE
    if color == '8B8680' and (text.strip() == '-' or text.startswith('単位')):
        set_color(r, ANTHRACITE); return 'table_misc→ANTHRACITE'

    # 11) KPI titles (sz=10 7A8450 bold) → ANTHRACITE bold
    if size == 10.0 and color == '7A8450':
        set_color(r, ANTHRACITE); set_bold(r, True); return 'kpi_label→ANTHRACITE bold'

    # 12) KPI "vs 4月..." (sz=8 8B8680) → ANTHRACITE
    if size == 8.0 and color == '8B8680' and text.startswith('vs'):
        set_color(r, ANTHRACITE); return 'kpi_vs→ANTHRACITE'

    # 13) TOP5 card headings "5月 ... TOP5" (sz=9.5 CC785C) → ANTHRACITE bold
    if size == 9.5 and color == 'CC785C':
        set_color(r, ANTHRACITE); set_bold(r, True); return 'top5_heading→ANTHRACITE'

    # 14) TOP5 unit labels "[千円]"/"[行]"/"[行/千円]" near TOP5 row (sz=8 8B8680) → ANTHRACITE
    if size == 8.0 and color == '8B8680' and 6.90 <= t <= 6.93:
        set_color(r, ANTHRACITE); return 'top5_unit→ANTHRACITE'

    # 15) TOP5 rank "1..5" (sz=9 CC785C bold, in rank columns) → ANTHRACITE bold
    if size == 9.0 and color == 'CC785C' and (
        (0.37 <= l <= 0.39) or (2.96 <= l <= 2.98) or (5.55 <= l <= 5.57)
    ) and 7.10 <= t <= 7.95:
        set_color(r, ANTHRACITE); return 'top5_rank→ANTHRACITE'

    # 16) TOP5 values (sz=11 A8503D bold) → ANTHRACITE bold
    if size == 11.0 and color == 'A8503D':
        set_color(r, ANTHRACITE); return 'top5_value→ANTHRACITE'

    # 17) Positioning chart axis labels "コスト [千円/人]" (l=0.55,t=8.40 sz=7 CC785C)
    #     and "行数 [行/人]" (l=6.45,t=9.70 sz=7 CC785C) → ANTHRACITE
    if size == 7.0 and color == 'CC785C':
        set_color(r, ANTHRACITE); return 'chart_axis→ANTHRACITE'

    # 18) Quadrant labels in positioning chart (sz=9 7A8450/A8503D) → ANTHRACITE
    if size == 9.0 and (8.40 <= t <= 8.42 or 9.85 <= t <= 9.87) and color in ('7A8450', 'A8503D'):
        set_color(r, ANTHRACITE); return 'quadrant→ANTHRACITE'

    # 19) "バブル径 = MAU..." (sz=6.5 8B8680) → ANTHRACITE
    if size == 6.5 and color == '8B8680':
        set_color(r, ANTHRACITE); return 'bubble_caption→ANTHRACITE'

    # 20) Chart tick labels (sz=6 8B8680) → ANTHRACITE
    if size == 6.0 and color == '8B8680':
        set_color(r, ANTHRACITE); return 'tick→ANTHRACITE'

    # 21) IPA card "【参考】 IPA分析" (sz=10 CC785C bold) → ANTHRACITE bold
    if 4.30 <= l <= 4.33 and 10.43 <= t <= 10.47 and color == 'CC785C':
        set_color(r, ANTHRACITE); set_bold(r, True); return 'ipa_head→ANTHRACITE'

    # 22) IPA URL (sz=5.5 CC785C) → ANTHRACITE
    if size == 5.5 and color == 'CC785C':
        set_color(r, ANTHRACITE); return 'ipa_url→ANTHRACITE'

    # 23) IPA caption ※... (sz=7.5 8B8680) → ANTHRACITE
    if size == 7.5 and color == '8B8680':
        set_color(r, ANTHRACITE); return 'ipa_cap→ANTHRACITE'

    # 24) Reference band headings 改良/新規開発 (CC785C sz=9 bold) → ANTHRACITE bold
    if 4.39 <= l <= 4.43 and (11.03 <= t <= 11.05 or 11.35 <= t <= 11.37) and color == 'CC785C':
        set_color(r, ANTHRACITE); return 'ref_band_head→ANTHRACITE'

    return None

def main():
    os.chdir(os.path.dirname(os.path.abspath(__file__)))
    prs = Presentation(SRC)
    changes = []
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            l = Emu(sh.left).inches if sh.left else 0
            t = Emu(sh.top).inches if sh.top else 0
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    res = patch_run(sh, p, r, l, t)
                    if res:
                        changes.append(f'({l:.2f},{t:.2f}) sz={size_of(r)} {res} text={r.text!r}')
    prs.save(SRC)
    for c in changes:
        print(c)
    print(f'TOTAL: {len(changes)} run-level changes')

if __name__ == '__main__':
    main()
