#!/usr/bin/env python3
"""
V11 rev14 — in-place 4 changes (preserves all other shapes):

  1. Header title:  「Claude Code 月次利用量レポート」
                  → 「Claude Code 利用量レポート」  (drop "月次")
  2. Plan/actual table footnote (Y≈4.92):  delete the「* 6月末まで…」line
  3. KPI tile #2: add 登録ユーザー (mock 18名) alongside アクティブ
  4. Team mapping title: drop subtitle 「— 当チームは右上の優良ゾーン」
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V11.pptx"

DEEP_FOREST = RGBColor(0x1F, 0x3A, 0x2E)
GREEN_OK    = RGBColor(0x2E, 0x6F, 0x40)
MUTED       = RGBColor(0x6B, 0x6B, 0x6B)
JP_FONT, EN_FONT = "メイリオ", "Segoe UI"

def set_run_props(run, *, size=None, color=None, bold=None, ea=JP_FONT, latin=EN_FONT):
    if size is not None:
        run.font.size = Pt(size)
    if color is not None:
        run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold
    if latin is not None:
        run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    if ea is not None:
        for old in rPr.findall(qn('a:ea')):
            rPr.remove(old)
        ea_elem = etree.SubElement(rPr, qn('a:ea'))
        ea_elem.set('typeface', ea)

def replace_text_keep_props(shape, new_text):
    """Replace all run text with new_text, preserving the FIRST run's font."""
    tf = shape.text_frame
    # Keep first run's properties
    p0 = tf.paragraphs[0]
    if not p0.runs:
        return
    first = p0.runs[0]
    first.text = new_text
    # Clear other runs
    for run in list(p0.runs)[1:]:
        run._r.getparent().remove(run._r)
    # Clear additional paragraphs
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)

prs = Presentation(PPTX)
slide = prs.slides[0]

changes = []

for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    text = sh.text_frame.text
    try:
        t_in = sh.top / 914400
        l_in = sh.left / 914400
    except Exception:
        continue

    # --- 1. Header title ---
    if "Claude Code 月次利用量レポート" in text and t_in < 0.50:
        replace_text_keep_props(sh, "Claude Code 利用量レポート")
        changes.append("✓ 1. Header title: '月次' removed")

    # --- 2. Table footnote (Y ~4.92) ---
    if 4.85 <= t_in < 5.05 and "6月末までの累積実績" in text:
        # Delete the entire shape
        sp = sh._element
        sp.getparent().remove(sp)
        changes.append("✓ 2. Table footnote shape deleted")

    # --- 3a. KPI #2 label ---
    if 5.55 <= t_in < 5.75 and 2.20 <= l_in < 2.50 and text.strip() == "6月 アクティブユーザー":
        replace_text_keep_props(sh, "ユーザー数 (アクティブ/登録)")
        changes.append("✓ 3a. KPI#2 label updated")

    # --- 3b. KPI #2 big value ---
    if 5.70 <= t_in < 5.95 and 2.20 <= l_in < 2.50 and text.strip() == "12名":
        replace_text_keep_props(sh, "12 / 18名")
        changes.append("✓ 3b. KPI#2 big value: '12名' → '12 / 18名'")

    # --- 3c. KPI #2 delta ---
    if 5.95 <= t_in < 6.15 and 2.20 <= l_in < 2.50 and text.strip() == "+2名":
        replace_text_keep_props(sh, "+2 / +1名")
        changes.append("✓ 3c. KPI#2 delta: '+2名' → '+2 / +1名'")

    # --- 3d. KPI #2 sub ---
    if 6.10 <= t_in < 6.25 and 2.20 <= l_in < 2.50 and "vs 5月" in text and "10名" in text:
        replace_text_keep_props(sh, "vs 5月 10 / 17名")
        changes.append("✓ 3d. KPI#2 sub: 'vs 5月 10名' → 'vs 5月 10 / 17名'")

    # --- 4. Team mapping title ---
    if "全社14チームのポジショニング" in text:
        replace_text_keep_props(sh, "全社14チームのポジショニング")
        changes.append("✓ 4. Team-map title shortened")

print("\n".join(changes))
assert len(changes) == 7, f"Expected 7 changes (1 + 1 + 4 + 1), got {len(changes)}"

prs.save(PPTX)
print(f"\nWROTE: {PPTX}")
