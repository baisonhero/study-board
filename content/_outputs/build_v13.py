#!/usr/bin/env python3
"""
V13 build — reads the V12-state pptx (saved at the V11 path) and writes
V13 as a NEW file. Preserves V11/V12 on disk.

5 changes vs V12:

  1. Q1 framework (3 months only)
     - Header sub-title "2026年度Q1 (4-6月) / 報告日: 2026-07-01"
     - Version pill V11 → V13
     - Chart re-rendered with 3 months (4-6月) only
     - Plan/actual table → 3 months + Q1合計
     - Forecast card RIGHT becomes "Q1当初計画 vs Q1実績消費率"

  2. Plan/actual table: add 消費率 row (3 rows total)
     消費率 = 実績 ÷ 計画 × 100  (84.0% for Q1)

  3. KPI ユーザー数 → MAU only
     "ユーザー数 (アクティブ/登録) / 12 / 18名 …" → "MAU [6月] / 12名 / +2名 / vs 5月 10名"

  4. KPI 4 cards → 3 cards (delete 採用行数 card)
     Resize remaining 3 cards to 1/3 width each.

  5. Team mapping new axes
     X = 1人当たり生成行数 (行/人)
     Y = 1人当たりコスト (千円/人)
     Bubble = MAU
     Quadrant labels updated; own team relocated to 下右 (低コスト×高生産)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import MaxNLocator

plt.rcParams["font.family"] = ["DejaVu Sans", "Droid Sans Fallback",
                                "Noto Sans CJK JP", "Hiragino Sans",
                                "Yu Gothic", "Meiryo"]
plt.rcParams["axes.unicode_minus"] = False

SRC = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V11.pptx"
DST = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V13.pptx"

# Palette (from house style)
DEEP_FOREST = RGBColor(0x1F, 0x3A, 0x2E)
FOREST      = RGBColor(0x2E, 0x6F, 0x40)
SAGE        = RGBColor(0xB5, 0xC9, 0xA8)
SAGE_BG     = RGBColor(0xE8, 0xEF, 0xE3)
AMBER       = RGBColor(0xF4, 0xB9, 0x42)
AMBER_DEEP  = RGBColor(0xC8, 0x8E, 0x1F)
INK         = RGBColor(0x22, 0x28, 0x24)
MUTED       = RGBColor(0x6B, 0x6B, 0x6B)
CARD_BORDER = RGBColor(0xD7, 0xDC, 0xCF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_OK    = RGBColor(0x2E, 0x6F, 0x40)
RED         = RGBColor(0xC0, 0x39, 0x2B)
COLOR_OWN   = AMBER
COLOR_OTHER = RGBColor(0xA8, 0xB8, 0xA0)
JP_FONT, EN_FONT = "メイリオ", "Segoe UI"

forest_hex   = "#1F3A2E"
muted_hex    = "#6B6B6B"
burnt_orange = "#D97706"

# ---------------- helpers ----------------
def set_run_font(run, size=12, color=INK, bold=False, latin=EN_FONT, ea=JP_FONT):
    run.font.name = latin
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea_elem = etree.SubElement(rPr, qn('a:ea'))
    ea_elem.set('typeface', ea)

def replace_text(shape, new_text, *, size=None, color=None, bold=None):
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        return False
    r0 = p0.runs[0]
    r0.text = new_text
    if size is not None: r0.font.size = Pt(size)
    if color is not None: r0.font.color.rgb = color
    if bold is not None: r0.font.bold = bold
    for r in list(p0.runs)[1:]:
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    return True

def add_textbox(slide, text, l, t, w, h, *, size=12, color=INK, bold=False,
                align='left', valign='top', word_wrap=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if valign == 'middle': tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if align == 'center': p.alignment = PP_ALIGN.CENTER
    elif align == 'right': p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    set_run_font(r, size=size, color=color, bold=bold)
    return tb

def add_rect(slide, l, t, w, h, *, fill=WHITE, line=None, line_width=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                                 Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        if line_width: shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_oval(slide, l, t, w, h, *, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t),
                                 Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

# ---------------- 1. Render new chart (3 months) ----------------
MONTHS = ['4月','5月','6月']
COST_ACT  = [580, 620, 690]
LOC_ACT   = [15, 17, 20]

fig, ax_left = plt.subplots(figsize=(7.67, 1.30), dpi=220)
fig.patch.set_facecolor("white")
x_pos = list(range(len(MONTHS)))
ax_left.bar(x_pos, COST_ACT, color=forest_hex, edgecolor=forest_hex,
            width=0.55, linewidth=0.5, label="単月コスト (実績)", zorder=2)
ax_left.set_xticks(x_pos)
ax_left.set_xticklabels(MONTHS, fontsize=7, color=muted_hex)
ax_left.set_ylabel("単月コスト [千円]", color=forest_hex, fontsize=7, labelpad=2)
ax_left.tick_params(axis="y", labelcolor=forest_hex, labelsize=6.5, length=2, pad=1)
ax_left.tick_params(axis="x", length=2, pad=1)
ax_left.grid(True, axis="y", linestyle="--", linewidth=0.4,
             color=muted_hex, alpha=0.35, zorder=1)
ax_left.set_axisbelow(True)
for spine in ("top", "right"):
    ax_left.spines[spine].set_visible(False)
ax_left.spines["left"].set_color(forest_hex)
ax_left.spines["left"].set_linewidth(0.6)
ax_left.spines["bottom"].set_color(muted_hex)
ax_left.spines["bottom"].set_linewidth(0.6)
ax_left.yaxis.set_major_locator(MaxNLocator(nbins=5, integer=True))
ax_left.set_ylim(0, max(COST_ACT) * 1.30)

ax_right = ax_left.twinx()
ax_right.plot(x_pos, LOC_ACT, color=burnt_orange, marker="o", markersize=4,
              markerfacecolor=burnt_orange, markeredgecolor=burnt_orange,
              linewidth=2.0, label="単月行数 (実績)", zorder=4)
ax_right.set_ylabel("単月行数 [千行]", color=burnt_orange, fontsize=7, labelpad=2)
ax_right.tick_params(axis="y", labelcolor=burnt_orange, labelsize=6.5,
                     length=2, pad=1)
for spine in ("top", "left"):
    ax_right.spines[spine].set_visible(False)
ax_right.spines["right"].set_color(burnt_orange)
ax_right.spines["right"].set_linewidth(0.8)
ax_right.set_ylim(0, max(LOC_ACT) * 1.40)

# Legend external
h1, l1 = ax_left.get_legend_handles_labels()
h2, l2 = ax_right.get_legend_handles_labels()
ax_left.legend(h1 + h2, l1 + l2,
               loc="lower center", bbox_to_anchor=(0.5, 1.02), ncol=2,
               fontsize=6.5, frameon=False,
               handlelength=1.4, handletextpad=0.35,
               columnspacing=1.0, borderpad=0.0,
               labelcolor=[forest_hex, burnt_orange])
plt.subplots_adjust(left=0.06, right=0.95, top=0.85, bottom=0.16)
CHART_PNG = "/tmp/v13_trend.png"
fig.savefig(CHART_PNG, dpi=220, facecolor="white", edgecolor="none")
plt.close(fig)
print("✓ Chart rendered (4-6月 only)")

# ---------------- Open + apply ----------------
prs = Presentation(SRC)
slide = prs.slides[0]

# === HEADER: version pill V11 → V13, add Q1 subtitle ===
header_subtitle_exists = False
for sh in slide.shapes:
    if not sh.has_text_frame: continue
    try:
        t_in = sh.top/914400
    except: continue
    if t_in < 0.50 and "V11" in sh.text_frame.text:
        replace_text(sh, "V13")
        print("✓ Version pill V11 → V13")
    if 0.55 <= t_in < 0.75 and "2026" in sh.text_frame.text:
        header_subtitle_exists = True

# Add Q1 subtitle below title (if not present)
if not header_subtitle_exists:
    add_textbox(slide, "2026年度Q1 (4-6月)  /  報告日: 2026-07-01",
                l=0.42, t=0.62, w=7.50, h=0.20,
                size=9, color=MUTED)
    print("✓ Added Q1 subtitle")

# === JUDGE card: 投資規模 line update ===
for sh in slide.shapes:
    if sh.has_text_frame and "外注費用 2.6人月相当" in sh.text_frame.text:
        replace_text(sh, "Q1消費率 84% (実績 1,890 / 計画 2,250 [千円])",
                     size=9, color=DEEP_FOREST, bold=True)
        print("✓ Judge 投資規模 line → Q1消費率")

# === FORECAST CARD RIGHT: reshape to Q1当初計画 vs 実績 ===
# Current shapes (positions from V11/rev17):
#   "9月までの予測コスト" / "4,140" / "[千円]" / "外注費用 2.6人月" / "(1,600 千円/月 換算)"
# Targets:
#   "Q1当初計画" / "2,250" / "[千円]" / "計画消費率 84%" / "(実績 1,890 / 計画 2,250)"
RIGHT_HALF_X = 6.485
forecast_swap = {
    "9月までの予測コスト": "Q1当初計画",
    "4,140": "2,250",
    "外注費用 2.6人月": "計画消費率 84%",
    "(1,600 千円/月 換算)": "(実績 1,890 / 計画 2,250)",
}
# Also update LEFT card title from "前月までの累積 (実績)" to "Q1実績累積"
left_swap = {
    "前月までの累積 (実績)": "Q1実績累積",
    "1,890": "1,890",  # no change but ensure
}
ALL_SWAPS = {**forecast_swap, **left_swap}
for sh in slide.shapes:
    if not sh.has_text_frame: continue
    txt = sh.text_frame.text.strip()
    if txt in ALL_SWAPS and ALL_SWAPS[txt] != txt:
        replace_text(sh, ALL_SWAPS[txt])
        print(f"  ✓ '{txt}' → '{ALL_SWAPS[txt]}'")

# === CHART picture swap ===
for sh in slide.shapes:
    if sh.shape_type == 13:
        l, t, w, h = sh.left, sh.top, sh.width, sh.height
        sh._element.getparent().remove(sh._element)
        slide.shapes.add_picture(CHART_PNG, l, t, width=w, height=h)
        print("✓ Chart picture replaced (3 months)")
        break

# === KPI restructure: delete 採用行数 card (4 shapes at L=6.18), resize 3 ===
# Identify cards by their L:
#   Card 1 (コスト):     L=0.30
#   Card 2 (ユーザー):   L=2.26
#   Card 3 (1人あたり月額): L=4.22
#   Card 4 (採用行数):   L=6.18  ← DELETE
# After: each card W = (7.67 - 2*0.20) / 3 = 2.4233 ≈ 2.42
PAGE_W = 7.67
n_cards = 3
gap = 0.20
new_w = (PAGE_W - (n_cards-1)*gap) / n_cards   # 2.42
new_starts = [0.30 + i * (new_w + gap) for i in range(n_cards)]   # [0.30, 2.92, 5.54]

# Map old L → new L for kept cards (and contents)
KPI_L_MAP = {
    0.30: 0.30,   # card 1
    2.26: 2.92,   # card 2 → moves right
    4.22: 5.54,   # card 3 → moves right
    # 6.18 → DELETE
}
# Inner text L offsets: tile bg at card L, inner texts at card L + 0.08
INNER_OFFSET = 0.08

kpi_to_delete = []
for sh in slide.shapes:
    try:
        t_in = sh.top/914400; l_in = sh.left/914400
    except: continue
    if not (5.55 <= t_in < 6.30):
        continue
    # Is this in a card column?
    # Card BG L values: {0.30, 2.26, 4.22, 6.18}; inner texts offset +0.08
    card_l = None
    for kl in (0.30, 2.26, 4.22, 6.18):
        if abs(l_in - kl) < 0.02:
            card_l = kl
            inner_offset = 0
            break
        if abs(l_in - (kl + 0.08)) < 0.02:
            card_l = kl
            inner_offset = 0.08
            break
    if card_l is None:
        continue

    if card_l == 6.18:
        # Delete entire card
        kpi_to_delete.append(sh)
        continue

    # Reposition + resize
    new_card_l = KPI_L_MAP[card_l]
    sh.left = Inches(new_card_l + inner_offset)
    # Resize bg rects (width = 1.80 → new_w)
    cur_w_in = sh.width / 914400
    if abs(cur_w_in - 1.80) < 0.05:
        sh.width = Inches(new_w)
    elif abs(cur_w_in - 1.64) < 0.05:
        sh.width = Inches(new_w - 0.16)
    elif abs(cur_w_in - 1.20) < 0.05:
        sh.width = Inches(new_w - 0.40)

for sh in kpi_to_delete:
    sh._element.getparent().remove(sh._element)
print(f"✓ KPI deleted 採用行数 card ({len(kpi_to_delete)} shapes), repositioned remaining 3 cards")

# === MAU card text update (was "ユーザー数 (アクティブ/登録)" etc) ===
mau_swap = {
    "ユーザー数 (アクティブ/登録)": "MAU [6月]",
    "12 / 18名": "12名",
    "+2 / +1名": "+2名",
    "vs 5月 10 / 17名": "vs 5月 10名",
}
for sh in slide.shapes:
    if not sh.has_text_frame: continue
    txt = sh.text_frame.text.strip()
    if txt in mau_swap:
        replace_text(sh, mau_swap[txt])
        print(f"  ✓ MAU card: '{txt}' → '{mau_swap[txt]}'")

# === PLAN/ACTUAL table: delete 7-9月 cells, change "6ヶ月合計"→"Q1合計", values, add 消費率 row ===
# Current table (rev18b): 6 months 4-9月 at L positions 0.95, 1.86, 2.77, 3.68, 4.59, 5.50
# Each W=0.91. 合計 col at L=6.41 W=1.56.
# Delete 7-9月 month cells (L in {3.68, 4.59, 5.50}, T in {4.42, 4.58, 4.74})
DELETE_TABLE_L = {3.68, 4.59, 5.50}
DELETE_TABLE_T = {4.42, 4.58, 4.74}
table_to_delete = []
for sh in slide.shapes:
    try:
        t_in = round(sh.top/914400, 2); l_in = round(sh.left/914400, 2)
    except: continue
    if t_in in DELETE_TABLE_T and l_in in DELETE_TABLE_L:
        table_to_delete.append(sh)
    # Also delete the vertical separator line (between 6月 and 7月) at L=3.68
    if sh.shape_type == 9 and abs(l_in - 3.68) < 0.05 and 4.30 <= t_in < 5.00:
        table_to_delete.append(sh)
# Dedupe
table_to_delete = list({id(s): s for s in table_to_delete}.values())
for sh in table_to_delete:
    sh._element.getparent().remove(sh._element)
print(f"✓ Deleted {len(table_to_delete)} table cells (7-9月 + separator)")

# Update table totals (rev18 had "6ヶ月合計"/"4,500"; V13 needs "Q1合計"/"2,250")
totals_swap = {
    "6ヶ月合計": "Q1合計",
    "4,500": "2,250",
    "1,890*": "1,890",   # asterisk no longer needed (no future months)
}
# But also resize totals column (currently L=6.41 W=1.56) — move to L=4.20 W=3.77 to fill?
# Actually let's keep 4-6月 columns at original L (0.95, 1.86, 2.77, W=0.91) and just move
# the 合計 col closer (since 7-9月 columns are gone).
# New layout: 4月..6月 same, then Q1合計 at L=3.85 W=2.50
TOTALS_L = 3.85
TOTALS_W = 4.12   # fills remaining width to 7.97
for sh in slide.shapes:
    if not sh.has_text_frame: continue
    txt = sh.text_frame.text.strip()
    if txt in totals_swap:
        replace_text(sh, totals_swap[txt])
        print(f"  ✓ Totals: '{txt}' → '{totals_swap[txt]}'")
    # Reposition totals shapes (were at L=6.41 W=1.56)
    try:
        t_in = round(sh.top/914400, 2); l_in = round(sh.left/914400, 2)
        if l_in == 6.41 and t_in in {4.42, 4.58, 4.74}:
            sh.left = Inches(TOTALS_L)
            sh.width = Inches(TOTALS_W)
    except: pass

# Add 消費率 row at T=4.90, h=0.16
ROW3_T = 4.90
add_rect(slide, 0.30, ROW3_T, 7.67, 0.16, fill=WHITE, line=CARD_BORDER, line_width=0.3)
# Row label
add_textbox(slide, "消費率", l=0.35, t=ROW3_T, w=0.55, h=0.16,
            size=8, color=INK, bold=True, valign='middle')
# Month values
month_lvals = {"4月": 0.95, "5月": 1.86, "6月": 2.77}
consumption = {"4月": "77.3%", "5月": "82.7%", "6月": "92.0%"}
for m in ["4月","5月","6月"]:
    val = consumption[m]
    pct = float(val.rstrip("%"))
    color = GREEN_OK if pct < 100 else AMBER_DEEP
    add_textbox(slide, val, l=month_lvals[m], t=ROW3_T, w=0.91, h=0.16,
                size=7, color=color, bold=True, align='center', valign='middle')
# Q1合計 column 消費率: 84.0%
add_textbox(slide, "84.0%", l=TOTALS_L, t=ROW3_T, w=TOTALS_W, h=0.16,
            size=8, color=GREEN_OK, bold=True, align='center', valign='middle')
print("✓ 消費率 row added (84.0%)")

# === TEAM MAPPING: new axes + new bubbles ===
PX, PY, PW, PH = 1.20, 8.65, 5.00, 1.05

# Update Y axis label "↑ 効率\n(行/コスト)" → "↑ 1人当たり\nコスト (千円/人)"
# Update X axis label "→ 使用率" → "→ 行数/人"
for sh in slide.shapes:
    if not sh.has_text_frame: continue
    txt = sh.text_frame.text.strip()
    try:
        t_in = sh.top/914400; l_in = sh.left/914400
    except: continue
    # Y axis
    if 8.30 <= t_in < 10.20 and l_in < 0.50 and "効率" in txt:
        # Keep the rotated multi-line style
        tf = sh.text_frame
        p0 = tf.paragraphs[0]
        if p0.runs:
            p0.runs[0].text = "↑ 1人当たり\nコスト (千円/人)"
            for r in list(p0.runs)[1:]:
                r._r.getparent().remove(r._r)
        print("  ✓ Y axis label updated")
    # X axis
    if 9.60 <= t_in < 9.85 and "使用率" in txt:
        replace_text(sh, "→ 行数/人")
        print("  ✓ X axis label updated")
    # Quadrant labels (Y in 8.40-8.45 and 9.85-9.90)
    if 8.38 <= t_in < 8.50:
        if "低使用 × 高効率" in txt:
            replace_text(sh, "▼ 高コスト × 低生産")
        elif "高使用 × 高効率" in txt:
            replace_text(sh, "高コスト × 高生産 ▼")
    if 9.83 <= t_in < 9.92:
        if "低使用 × 低効率" in txt:
            replace_text(sh, "▲ 低コスト × 低生産")
        elif "高使用 × 低効率" in txt:
            replace_text(sh, "低コスト × 高生産 ▲")
    # バブル径 caption
    if 9.95 <= t_in < 10.05 and "バブル径" in txt:
        replace_text(sh, "バブル径 = MAU (アクティブユーザー数)")

# Delete existing bubbles and own-team label
bubbles_to_delete = []
for sh in slide.shapes:
    if sh.shape_type not in (MSO_SHAPE.OVAL, 9):  # ovals + lines? lines stay
        if sh.shape_type != 9:  # not line
            try:
                t_in = sh.top/914400; l_in = sh.left/914400
                # Bubbles are ovals in 8.60-9.70, with small dimensions
                # But add_oval creates MSO_SHAPE.OVAL = 9... wait that's LINE
                # Actually MSO_SHAPE.OVAL is enum; bubbles in V11 added by add_oval
                # Let me identify by type AUTO_SHAPE (1) and small dims
                w_in = sh.width/914400; h_in = sh.height/914400
                if (sh.shape_type == 1 and 0.05 <= w_in <= 0.35
                    and abs(w_in - h_in) < 0.01
                    and 8.55 <= t_in < 9.70):
                    bubbles_to_delete.append(sh)
            except: pass
# Also remove "当チーム" label (will re-add at new position)
for sh in slide.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip() == "当チーム":
        try:
            t_in = sh.top/914400
            if 8.50 <= t_in < 9.80:
                bubbles_to_delete.append(sh)
        except: pass

for sh in bubbles_to_delete:
    sh._element.getparent().remove(sh._element)
print(f"✓ Deleted {len(bubbles_to_delete)} old bubbles + own label")

# Add new bubbles with new axes
# Mock data: 14 teams (name, X=行/人, Y=千円/人, MAU, is_own)
TEAMS = [
    ("当チーム", 1667, 57.5, 12, True),
    ("B",       1900, 70.0, 15, False),
    ("C",       1400, 80.0, 18, False),
    ("D",       1200, 55.0, 10, False),
    ("E",       2000, 85.0, 20, False),
    ("F",        800, 95.0, 14, False),
    ("G",        600,105.0,  8, False),
    ("H",        900,110.0, 12, False),
    ("I",       1100, 65.0, 22, False),
    ("J",        500, 45.0,  7, False),
    ("K",        400, 50.0,  9, False),
    ("L",        700, 75.0, 11, False),
    ("M",       1500, 65.0, 16, False),
    ("N",       2200, 95.0, 25, False),
]
X_MIN, X_MAX = 300, 2400     # 行/人
Y_MIN, Y_MAX = 40, 115        # 千円/人
M_MIN, M_MAX = 7, 25          # MAU bubble dia 0.10–0.22 in

def map_xy(x_val, y_val):
    nx = (x_val - X_MIN) / (X_MAX - X_MIN)
    ny = (y_val - Y_MIN) / (Y_MAX - Y_MIN)
    vx = PX + nx * PW
    # Y inverted (high cost = visual top)
    vy = PY + (1 - ny) * PH
    return vx, vy

def map_dia(mau):
    return 0.10 + (mau - M_MIN) / (M_MAX - M_MIN) * 0.12

# Draw bubbles (own last so it's on top)
for name, xv, yv, mau, is_own in sorted(TEAMS, key=lambda t: 1 if t[4] else 0):
    cx, cy = map_xy(xv, yv)
    d = map_dia(mau)
    bx, by = cx - d/2, cy - d/2
    fill = COLOR_OWN if is_own else COLOR_OTHER
    line = AMBER_DEEP if is_own else FOREST
    add_oval(slide, bx, by, d, d, fill=fill, line=line)
    if is_own:
        # Label "当チーム" near own bubble
        label_x = bx + d + 0.02
        label_y = by - 0.04
        if label_x + 0.65 > PX + PW + 0.10:
            label_x = bx - 0.65
        add_textbox(slide, "当チーム",
                    l=label_x, t=label_y, w=0.65, h=0.14,
                    size=8, color=AMBER_DEEP, bold=True)
print(f"✓ Drew {len(TEAMS)} new team bubbles (own team in 下右)")

# === SAVE as V13 ===
prs.save(DST)
print(f"\n✓ WROTE: {DST}")
