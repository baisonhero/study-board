#!/usr/bin/env python3
"""
V11 rev10 — in-place 3 changes (preserves all other shapes):
  1. Monthly trend chart: AMBER #F4B942 (薄い) → BURNT_ORANGE #D97706 (濃い)
     - LOC line color
     - Right Y axis label + ticks + spine
     - Legend amber items color
     Picture is re-rendered then swapped in place at the same coordinates.

  2. Team mapping: add a title ABOVE the outer rect
     "全社14チームのポジショニング — 当チームは右上の優良ゾーン"
     14pt bold, DEEP_FOREST, placed at T=8.05 (in the gap above the team-map rect).

  3. Quadrant markers (bottom row): ▶ / ◀ → ▲ / ▲
     - "▶ 低使用 × 低効率" → "▲ 低使用 × 低効率"
     - "高使用 × 低効率 ◀" → "高使用 × 低効率 ▲"
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import MaxNLocator

plt.rcParams["font.family"] = [
    "DejaVu Sans", "Droid Sans Fallback",
    "Noto Sans CJK JP", "Hiragino Sans", "Yu Gothic", "Meiryo",
]
plt.rcParams["axes.unicode_minus"] = False

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V11.pptx"

# ---------- Palette ----------
DEEP_FOREST = RGBColor(0x1F, 0x3A, 0x2E)
MUTED       = RGBColor(0x6B, 0x6B, 0x6B)

JP_FONT = "メイリオ"
EN_FONT = "Segoe UI"

# Chart colors
forest_hex   = "#1F3A2E"
muted_hex    = "#6B6B6B"
burnt_orange = "#D97706"   # ← NEW: darker orange in place of AMBER #F4B942

# ---------- Chart data (carried over from build_v11.py) ----------
MONTHS = ['4月','5月','6月','7月','8月','9月','10月','11月','12月','1月','2月','3月']
COST_MONTHLY = [580,620,690, 770,770,770,770,770,770,770,770,770]
LOC_MONTHLY  = [15,17,20,    22,22,22,22,22,22,22,22,22]
PLAN_MONTHLY = [750]*12
ACTUAL_MONTHS_COUNT = 3

# ---------- 1) Re-render chart with darker orange ----------
chart_w_in, chart_h_in = 7.67, 1.30
fig, ax_left = plt.subplots(figsize=(chart_w_in, chart_h_in), dpi=220)
fig.patch.set_facecolor("white")

x_pos = list(range(len(MONTHS)))
ax_left.bar(x_pos[:ACTUAL_MONTHS_COUNT],
            COST_MONTHLY[:ACTUAL_MONTHS_COUNT],
            color=forest_hex, edgecolor=forest_hex, width=0.55, linewidth=0.5,
            label="単月コスト (実績)", zorder=2)
ax_left.bar(x_pos[ACTUAL_MONTHS_COUNT:],
            PLAN_MONTHLY[ACTUAL_MONTHS_COUNT:],
            facecolor="white", edgecolor=forest_hex, hatch="///",
            width=0.55, linewidth=0.6,
            label="単月コスト (計画)", zorder=2)
ax_left.set_xticks(x_pos)
ax_left.set_xticklabels(MONTHS, fontsize=6.5, color=muted_hex)
ax_left.set_ylabel("単月コスト [千円]", color=forest_hex, fontsize=7, labelpad=2)
ax_left.tick_params(axis="y", labelcolor=forest_hex, labelsize=6.5, length=2, pad=1)
ax_left.tick_params(axis="x", length=2, pad=1)
ax_left.grid(True, axis="y", linestyle="--", linewidth=0.4,
             color=muted_hex, alpha=0.35, zorder=1)
ax_left.set_axisbelow(True)
for spine in ("top", "right"):
    ax_left.spines[spine].set_visible(False)
ax_left.spines["left"].set_color(forest_hex)
ax_left.spines["left"].set_linewidth(0.6)
ax_left.spines["bottom"].set_color(muted_hex)
ax_left.spines["bottom"].set_linewidth(0.6)
ax_left.yaxis.set_major_locator(MaxNLocator(nbins=5, integer=True))
ax_left.set_ylim(0, max(max(COST_MONTHLY[:ACTUAL_MONTHS_COUNT]),
                         max(PLAN_MONTHLY[ACTUAL_MONTHS_COUNT:])) * 1.25)

# Right axis with BURNT ORANGE (was AMBER)
ax_right = ax_left.twinx()
ax_right.plot(x_pos[:ACTUAL_MONTHS_COUNT], LOC_MONTHLY[:ACTUAL_MONTHS_COUNT],
              color=burnt_orange, marker="o", markersize=4,
              markerfacecolor=burnt_orange, markeredgecolor=burnt_orange,
              linewidth=2.0, label="単月行数 (実績)", zorder=4)
ax_right.set_ylabel("単月行数 [千行]", color=burnt_orange,
                    fontsize=7, labelpad=2)
ax_right.tick_params(axis="y", labelcolor=burnt_orange, labelsize=6.5,
                     length=2, pad=1)
for spine in ("top", "left"):
    ax_right.spines[spine].set_visible(False)
ax_right.spines["right"].set_color(burnt_orange)
ax_right.spines["right"].set_linewidth(0.8)
ax_right.set_ylim(0, max(LOC_MONTHLY[:ACTUAL_MONTHS_COUNT]) * 1.40)

# Vertical separator (実績 / 計画 境界)
ax_left.axvline(x=ACTUAL_MONTHS_COUNT - 0.5, color=muted_hex,
                linestyle=":", linewidth=0.6, alpha=0.7, zorder=1)
boundary_y = max(max(COST_MONTHLY[:ACTUAL_MONTHS_COUNT]),
                  max(PLAN_MONTHLY[ACTUAL_MONTHS_COUNT:])) * 1.13
ax_left.text(ACTUAL_MONTHS_COUNT - 0.6, boundary_y,
             "実績→", color=forest_hex, fontsize=6, ha="right", va="center")
ax_left.text(ACTUAL_MONTHS_COUNT - 0.4, boundary_y,
             "←計画", color=muted_hex, fontsize=6, ha="left", va="center")

# Legend external — top of plot, single row
h1, l1 = ax_left.get_legend_handles_labels()
h2, l2 = ax_right.get_legend_handles_labels()
ax_left.legend(h1 + h2, l1 + l2,
               loc="lower center", bbox_to_anchor=(0.5, 1.02), ncol=4,
               fontsize=6.5, frameon=False,
               handlelength=1.4, handletextpad=0.35,
               columnspacing=1.0, borderpad=0.0,
               labelcolor=[forest_hex, forest_hex, burnt_orange])

plt.subplots_adjust(left=0.06, right=0.95, top=0.85, bottom=0.14)
NEW_CHART_PNG = "/tmp/v11r10_trend.png"
fig.savefig(NEW_CHART_PNG, dpi=220, facecolor="white", edgecolor="none")
plt.close(fig)
print(f"Chart re-rendered (darker orange): {NEW_CHART_PNG}")

# ---------- Open pptx and apply 3 changes ----------
prs = Presentation(PPTX)
slide = prs.slides[0]

# 1a) Find chart picture and replace it (shape_type == 13)
picture_to_replace = None
for sh in slide.shapes:
    if sh.shape_type == 13:
        picture_to_replace = sh
        break

if picture_to_replace is None:
    raise RuntimeError("Chart picture not found in slide.")

pic_l = picture_to_replace.left
pic_t = picture_to_replace.top
pic_w = picture_to_replace.width
pic_h = picture_to_replace.height
print(f"Replacing picture at L={pic_l/914400:.2f} T={pic_t/914400:.2f} "
      f"W={pic_w/914400:.2f} H={pic_h/914400:.2f}")

# Delete old picture
sp_el = picture_to_replace._element
sp_el.getparent().remove(sp_el)
# Insert new
slide.shapes.add_picture(NEW_CHART_PNG, pic_l, pic_t, width=pic_w, height=pic_h)

# 2) Add team-mapping title above the outer rect
def set_run_font(run, latin=EN_FONT, ea=JP_FONT, size=12, color=DEEP_FOREST, bold=False):
    run.font.name = latin
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea_elem = etree.SubElement(rPr, qn('a:ea'))
    ea_elem.set('typeface', ea)

title_text = "全社14チームのポジショニング — 当チームは右上の優良ゾーン"
tb = slide.shapes.add_textbox(Inches(0.30), Inches(8.05),
                              Inches(7.67), Inches(0.22))
tf = tb.text_frame
tf.word_wrap = False
tf.margin_left = tf.margin_right = Emu(0)
tf.margin_top = tf.margin_bottom = Emu(0)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = title_text
set_run_font(r, size=14, color=DEEP_FOREST, bold=True)
print(f"Added team-map title at T=8.05: {title_text!r}")

# 3) Quadrant markers: ▶ / ◀ → ▲ for bottom labels
n_changed = 0
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    for para in sh.text_frame.paragraphs:
        for run in para.runs:
            t = run.text
            if "▶ 低使用 × 低効率" in t:
                run.text = t.replace("▶", "▲")
                n_changed += 1
                print(f"  ✓ bottom-left marker: '▶' → '▲'")
            elif "高使用 × 低効率 ◀" in t:
                run.text = t.replace("◀", "▲")
                n_changed += 1
                print(f"  ✓ bottom-right marker: '◀' → '▲'")
print(f"Changed {n_changed} markers")
assert n_changed == 2, f"Expected 2 marker changes, got {n_changed}"

# Save
prs.save(PPTX)
print(f"\nWROTE: {PPTX}")
