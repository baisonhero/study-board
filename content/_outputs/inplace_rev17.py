#!/usr/bin/env python3
"""
V11 rev17 — in-place: align the two big numbers (1,890 / 8,640) in the
forecast card so they share the same Y, height, and font size.

Unified layout inside the JG band (T=1.00–2.20):
  Header (combined)   T=1.04 H=0.18  9pt
  Big number          T=1.24 H=0.40  28pt   ← LEFT and RIGHT identical
  [千円] unit          T=1.66 H=0.14  8pt
  RIGHT only:
    外注費用 5.4人月    T=1.82 H=0.16  11pt
    (1,600 千円/月 換算) T=1.99 H=0.16  8pt
"""
from pptx import Presentation
from pptx.util import Inches, Pt

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V11.pptx"

LEFT_X  = 5.00
RIGHT_X = 6.485
HALF_W  = 1.49
JG_T    = 1.00

# Target Y/H/Pt by current text content
TARGETS = {
    # LEFT side
    "前月までの累積コスト":  dict(t=1.04, h=0.18, w=HALF_W, l=LEFT_X,  pt=9,  text="前月までの累積コスト (実績)"),
    "1,890":                dict(t=1.24, h=0.40, w=HALF_W, l=LEFT_X,  pt=28),
    # RIGHT side
    "年度末着地予測 (予測)": dict(t=1.04, h=0.18, w=HALF_W, l=RIGHT_X, pt=9),
    "8,640":                dict(t=1.24, h=0.40, w=HALF_W, l=RIGHT_X, pt=28),
    "外注費用 5.4人月":      dict(t=1.82, h=0.16, w=HALF_W, l=RIGHT_X, pt=11),
    "(1,600 千円/月 換算)":  dict(t=1.99, h=0.16, w=HALF_W, l=RIGHT_X, pt=8),
}

prs = Presentation(PPTX)
slide = prs.slides[0]

# Find both [千円] units (need to differentiate LEFT vs RIGHT)
n_ok = 0
to_delete = []
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    try:
        t_in = sh.top/914400; l_in = sh.left/914400
    except Exception:
        continue
    if not (1.00 <= t_in < 2.20 and 5.00 <= l_in < 8.00):
        continue
    txt = sh.text_frame.text.strip()
    # Delete the standalone "(実績)" — merging into title
    if txt == "(実績)":
        to_delete.append(sh)
        continue
    # Both [千円] units → new T=1.66 H=0.14 size 8pt
    if txt == "[千円]":
        sh.top = Inches(1.66)
        sh.height = Inches(0.14)
        sh.width = Inches(HALF_W)
        # left already at correct half, but ensure
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(8)
        n_ok += 1
        continue
    # Other targeted updates
    if txt in TARGETS:
        spec = TARGETS[txt]
        sh.top    = Inches(spec['t'])
        sh.height = Inches(spec['h'])
        sh.width  = Inches(spec['w'])
        sh.left   = Inches(spec['l'])
        # Update font size
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(spec['pt'])
        # Update text if needed
        if 'text' in spec:
            # replace first run text, clear others
            p0 = sh.text_frame.paragraphs[0]
            if p0.runs:
                p0.runs[0].text = spec['text']
                for r in list(p0.runs)[1:]:
                    r._r.getparent().remove(r._r)
        n_ok += 1
        print(f"  ✓ {txt!r} → T={spec['t']} H={spec['h']} pt={spec['pt']}")

for sh in to_delete:
    sp = sh._element
    sp.getparent().remove(sp)
    print(f"  ✓ Deleted '(実績)' (merged into title)")

prs.save(PPTX)
print(f"\nUpdated {n_ok} shapes, deleted {len(to_delete)}.\nWROTE: {PPTX}")
