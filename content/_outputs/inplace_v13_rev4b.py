#!/usr/bin/env python3
"""
V13 rev4b — right top card rewrite per user mockup.

Target layout:
  LEFT panel (DEEP_FOREST bg, white text):
    [22] '前月までの累積 (実績)'
    [23] '1,890'
    [24] '[千円]'

  RIGHT panel (AMBER bg, deep text):
    [25] '9月までの予測コスト'
    [26] '4,140'
    [27] '[千円]'
    [28] '外注費用 2.6人月'
    [29] '(1,600 千円/月 換算)'

Numbers are Q1 (4-6月) based:
  累積 = 580+620+690 = 1,890
  予測 = 1,890 + 750×3 = 4,140
  換算 = 4,140 / 1,600 = 2.59 → 2.6 人月
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V13.pptx"
JP_FONT, EN_FONT = "メイリオ", "Segoe UI"

def set_ea(run):
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', JP_FONT)

def replace_text_only(shape, new_text):
    """Replace text content only; keep all original font properties."""
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        return
    p0.runs[0].text = new_text
    for r in list(p0.runs)[1:]:
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    r0 = p0.runs[0]
    r0.font.name = EN_FONT
    set_ea(r0)

prs = Presentation(PPTX)
slide = prs.slides[0]

# Direct text-based replacement targeted within right top card region (T 1.00-2.20, L > 4.95)
REPLACEMENTS = {
    "Q1進捗実績 (4-5月)":          "前月までの累積 (実績)",
    "1,200":                       "1,890",
    "Q1当初計画":                  "9月までの予測コスト",
    "2,250":                       "4,140",
    "計画進捗率 53%":              "外注費用 2.6人月",
    "(実績 1,200 / 計画 2,250)":   "(1,600 千円/月 換算)",
}

hits = 0
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    try:
        t = sh.top/914400; l = sh.left/914400
    except:
        continue
    # Confine to right top card area
    if not (1.00 <= t <= 2.20 and l > 4.95):
        continue
    txt = sh.text_frame.text.strip()
    if txt in REPLACEMENTS:
        new_text = REPLACEMENTS[txt]
        replace_text_only(sh, new_text)
        print(f"  ✓ '{txt}' → '{new_text}'  (T={t:.2f} L={l:.2f})")
        hits += 1

assert hits >= 6, f"Expected 6 edits in right top card, got {hits}"

prs.save(PPTX)
print(f"\nWROTE: {PPTX}")
