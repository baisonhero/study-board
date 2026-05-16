#!/usr/bin/env python3
"""
V13 rev6 — KPI 3カードを V12 縦置き中央揃えパターンに復元 + 6月データ反映.

レイアウト目標（各カード内）:
  ラベル    T=5.64 H=0.16  size=10  bold  Forest Green  CENTER
  大きな数字 T=5.80 H=0.42  size=28  bold  DEEP_FOREST   CENTER
  変化率     T=6.22 H=0.18  size=13  bold  GREEN/RED     CENTER
  比較詳細   T=6.40 H=0.10  size=8        MUTED_GRAY    CENTER

カード高さ: 0.80 → 0.90 (bottom 6.50, 直下の TOP3 段に接する)

3カードのデータ更新:
  Card1: '5月 コスト [千円]' / 620 / '+6.9%' / 'vs 4月 580 [千円]'
      →  '6月 コスト [千円]' / 690 / '+11.3%' / 'vs 5月 620 [千円]'      [+ → GREEN]
  Card2: 'MAU [5月]'         / '10名' / '+2名' / 'vs 4月 8名'
      →  'ユーザー数 (MAU)'  / '12名' / '+2名' / 'vs 5月 10名'           [+ → GREEN]
  Card3: '5月 1人あたり月額 [千円]' / '62.0' / '−14.5%' / 'vs 4月 72.5 [千円]'
      →  '6月 1人あたり月額 [千円]' / '57.5' / '-4.5%'  / 'vs 5月 62.0 [千円]' [- → RED]
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V13.pptx"

DEEP_FOREST = RGBColor(0x1F, 0x3A, 0x2E)
FOREST      = RGBColor(0x2E, 0x6F, 0x40)
RED_CRIMSON = RGBColor(0xC0, 0x2A, 0x2A)
MUTED       = RGBColor(0x6B, 0x6B, 0x6B)
JP_FONT, EN_FONT = "メイリオ", "Segoe UI"

def set_ea(run):
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', JP_FONT)

def restyle(shape, new_text, *, size, color, bold, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p0 = tf.paragraphs[0]
    if not p0.runs:
        r = p0.add_run()
    else:
        p0.runs[0].text = new_text
        for r in list(p0.runs)[1:]:
            r._r.getparent().remove(r._r)
    if p0.runs:
        p0.runs[0].text = new_text
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    r0 = p0.runs[0]
    r0.font.name = EN_FONT
    r0.font.size = Pt(size)
    r0.font.bold = bold
    r0.font.color.rgb = color
    set_ea(r0)
    p0.alignment = align

prs = Presentation(PPTX)
slide = prs.slides[0]

# Card geometry (3 cards side by side)
CARDS = [
    {  # Card 1: 6月 コスト
        "L": 0.30,
        "label":  ("6月 コスト [千円]",   10, FOREST,      True),
        "value":  ("690",                 28, DEEP_FOREST, True),
        "delta":  ("+11.3%",              13, FOREST,      True),
        "sub":    ("vs 5月 620 [千円]",    8, MUTED,       False),
        "old_label": "5月 コスト [千円]",
        "old_value": "620",
        "old_delta": "+6.9%",
        "old_sub":   "vs 4月 580 [千円]",
    },
    {  # Card 2: MAU
        "L": 2.92,
        "label":  ("ユーザー数 (MAU)",    10, FOREST,      True),
        "value":  ("12名",                28, DEEP_FOREST, True),
        "delta":  ("+2名",                13, FOREST,      True),
        "sub":    ("vs 5月 10名",          8, MUTED,       False),
        "old_label": "MAU [5月]",
        "old_value": "10名",
        "old_delta": "+2名",
        "old_sub":   "vs 4月 8名",
    },
    {  # Card 3: 6月 1人あたり月額
        "L": 5.54,
        "label":  ("6月 1人あたり月額 [千円]", 10, FOREST,      True),
        "value":  ("57.5",                    28, DEEP_FOREST, True),
        "delta":  ("-4.5%",                   13, RED_CRIMSON, True),
        "sub":    ("vs 5月 62.0 [千円]",       8, MUTED,       False),
        "old_label": "5月 1人あたり月額 [千円]",
        "old_value": "62.0",
        "old_delta": "−14.5%",
        "old_sub":   "vs 4月 72.5 [千円]",
    },
]

# Layout Y/H within card (centered vertical stack)
LABEL_T, LABEL_H = 5.64, 0.16
VALUE_T, VALUE_H = 5.80, 0.42
DELTA_T, DELTA_H = 6.22, 0.20
SUB_T,   SUB_H   = 6.40, 0.10
CARD_W = 2.42
CARD_H_NEW = 0.90  # 5.60 + 0.90 = 6.50 (touches TOP3 row above)

# 1) Resize card background + accent stripe to new height (just background, not stripe)
for sh in slide.shapes:
    try:
        t = sh.top/914400; l = sh.left/914400; w = sh.width/914400; h = sh.height/914400
    except: continue
    if not (5.55 <= t <= 5.65 and h > 0.50):  # the big background AUTO_SHAPE
        continue
    for card in CARDS:
        if abs(l - card["L"]) < 0.05 and abs(w - CARD_W) < 0.05:
            sh.height = Inches(CARD_H_NEW)
            print(f"  ✓ Card bg @ L={l:.2f} → H={CARD_H_NEW}")

# 2) Restyle each text shape in each card
for card in CARDS:
    card_l = card["L"]
    inner_l = card_l + 0.08
    inner_w = CARD_W - 0.16
    hits = 0
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        try:
            t = sh.top/914400; l = sh.left/914400
        except: continue
        if not (5.60 <= t <= 6.30): continue
        # Match by card's left coordinate range
        if not (card_l - 0.05 <= l <= card_l + CARD_W):
            continue
        txt = sh.text_frame.text.strip()
        which = None
        if txt == card["old_label"]: which = "label"
        elif txt == card["old_value"]: which = "value"
        elif txt == card["old_delta"]: which = "delta"
        elif txt == card["old_sub"]: which = "sub"
        if which is None: continue

        new_text, size, color, bold = card[which]
        # Restyle text
        restyle(sh, new_text, size=size, color=color, bold=bold, align=PP_ALIGN.CENTER)
        # Reposition + resize
        sh.left = Inches(inner_l)
        sh.width = Inches(inner_w)
        if which == "label":
            sh.top = Inches(LABEL_T); sh.height = Inches(LABEL_H)
        elif which == "value":
            sh.top = Inches(VALUE_T); sh.height = Inches(VALUE_H)
        elif which == "delta":
            sh.top = Inches(DELTA_T); sh.height = Inches(DELTA_H)
        elif which == "sub":
            sh.top = Inches(SUB_T); sh.height = Inches(SUB_H)
        hits += 1
        print(f"  ✓ Card L={card_l:.2f}  {which:6s} → '{new_text}' ({size}pt)")
    assert hits == 4, f"Card L={card_l} expected 4 hits, got {hits}"

prs.save(PPTX)
print(f"\nWROTE: {PPTX}")
