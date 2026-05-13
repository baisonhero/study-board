#!/usr/bin/env python3
"""
Build AWS Cross-Account Private API Gateway proposal deck.

House-style: Calm Forest + Amber, メイリオ + Segoe UI, top title 28pt, bottom
conclusion message 18pt. 16:9 widescreen (13.333" x 7.5") layout.

Slides (10):
  1. 表紙
  2. 目次
  3. 背景 / 目的
  4. アーキテクチャ全体図
  5. Account A がやること
  6. Account B がやること
  7. 認証方式の比較
  8. Lambda Python サンプルコード
  9. Cognito トリガーとしての注意点
 10. コスト / 監視 / 追加考慮
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- House-style: Calm Forest + Amber ----------
DEEP_FOREST  = RGBColor(0x1F, 0x3A, 0x2E)
FOREST       = RGBColor(0x2E, 0x6F, 0x40)
SAGE         = RGBColor(0xB5, 0xC9, 0xA8)
SAGE_BG      = RGBColor(0xE8, 0xEF, 0xE3)
AMBER        = RGBColor(0xF4, 0xB9, 0x42)
AMBER_DEEP   = RGBColor(0xC8, 0x8E, 0x1F)
AMBER_TINT   = RGBColor(0xFD, 0xF1, 0xD9)
INK          = RGBColor(0x22, 0x28, 0x24)
MUTED        = RGBColor(0x6B, 0x6B, 0x6B)
CARD_BORDER  = RGBColor(0xD7, 0xDC, 0xCF)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG      = RGBColor(0x16, 0x1F, 0x1A)
CODE_INK     = RGBColor(0xE6, 0xED, 0xE2)
CODE_KEY     = RGBColor(0xF4, 0xB9, 0x42)
CODE_STR     = RGBColor(0xA8, 0xCB, 0x8A)
CODE_COMMENT = RGBColor(0x8A, 0x9A, 0x8E)

JP_FONT  = "メイリオ"
EN_FONT  = "Segoe UI"
MONO     = "Consolas"

# ---------- helpers ----------
def set_run_font(run, latin=EN_FONT, ea=JP_FONT, size=12, color=INK, bold=False, italic=False):
    run.font.name = latin
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea_elem = etree.SubElement(rPr, qn('a:ea'))
    ea_elem.set('typeface', ea)

def add_textbox(slide, text, l, t, w, h, size=12, color=INK, bold=False,
                align='left', valign='top', latin=EN_FONT, ea=JP_FONT,
                line_spacing=None, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = {'top': MSO_ANCHOR.TOP,
                          'middle': MSO_ANCHOR.MIDDLE,
                          'bottom': MSO_ANCHOR.BOTTOM}[valign]
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT,
                   'center': PP_ALIGN.CENTER,
                   'right': PP_ALIGN.RIGHT}[align]
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    set_run_font(r, latin=latin, ea=ea, size=size, color=color,
                 bold=bold, italic=italic)
    return tb

def add_multiline(slide, lines, l, t, w, h, size=12, color=INK,
                  align='left', valign='top', latin=EN_FONT, ea=JP_FONT,
                  line_spacing=1.15):
    """lines: list of (text, dict_of_overrides_or_None)."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = {'top': MSO_ANCHOR.TOP,
                          'middle': MSO_ANCHOR.MIDDLE,
                          'bottom': MSO_ANCHOR.BOTTOM}[valign]
    for i, item in enumerate(lines):
        text, ov = item if isinstance(item, tuple) else (item, None)
        ov = ov or {}
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = {'left': PP_ALIGN.LEFT,
                       'center': PP_ALIGN.CENTER,
                       'right': PP_ALIGN.RIGHT}[ov.get('align', align)]
        p.line_spacing = ov.get('line_spacing', line_spacing)
        if 'space_after' in ov:
            p.space_after = Pt(ov['space_after'])
        r = p.add_run()
        r.text = text
        set_run_font(r,
                     latin=ov.get('latin', latin),
                     ea=ov.get('ea', ea),
                     size=ov.get('size', size),
                     color=ov.get('color', color),
                     bold=ov.get('bold', False),
                     italic=ov.get('italic', False))
    return tb

def add_rect(slide, l, t, w, h, fill=WHITE, line=None, line_width=None,
             corner=None):
    if corner is not None:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(l), Inches(t), Inches(w), Inches(h))
        try:
            shp.adjustments[0] = corner
        except Exception:
            pass
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        if line_width:
            shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_arrow(slide, x1, y1, x2, y2, color=DEEP_FOREST, width=2.0,
              connector_type=MSO_CONNECTOR.STRAIGHT, dash=False):
    conn = slide.shapes.add_connector(connector_type,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    line = conn.line
    line.color.rgb = color
    line.width = Pt(width)
    # set arrow head
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    if dash:
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')
    return conn

def add_line(slide, x1, y1, x2, y2, color=CARD_BORDER, width=0.75, dash=False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dash:
        ln = line.line._get_or_add_ln()
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')
    return line


# ---------- presentation setup: 16:9 widescreen ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = 13.333
SH = 7.5
BLANK = prs.slide_layouts[6]


def add_page_chrome(slide, page_num, total_pages, title, conclusion):
    """Common chrome: amber accent + title (top, 28pt) + bottom message (18pt)."""
    # top amber accent bar (left)
    add_rect(slide, 0.50, 0.50, 0.08, 0.65, fill=AMBER)
    # title
    add_textbox(slide, title,
                l=0.70, t=0.45, w=11.5, h=0.75,
                size=28, color=DEEP_FOREST, bold=True, valign='middle')
    # page no.
    add_textbox(slide, f"{page_num} / {total_pages}",
                l=12.20, t=0.55, w=0.80, h=0.40,
                size=10, color=MUTED, align='right', valign='middle')
    # divider
    add_line(slide, 0.50, 1.28, 12.83, 1.28, color=CARD_BORDER, width=1.0)
    # bottom conclusion strip
    add_rect(slide, 0.50, 6.45, 12.33, 0.70, fill=SAGE_BG)
    add_rect(slide, 0.50, 6.45, 0.08, 0.70, fill=AMBER)
    add_multiline(slide, [
        ("このページの結論  ", {'size': 11, 'bold': True, 'color': AMBER_DEEP}),
        (conclusion,         {'size': 14, 'color': DEEP_FOREST, 'bold': True}),
    ], l=0.70, t=6.50, w=12.10, h=0.60, valign='middle', line_spacing=1.1)
    # Actually use single-line layout for conclusion - rewrite
    # Clear by removing the multiline we just added — easier to use two textboxes
    pass


def make_conclusion_strip(slide, conclusion):
    """Replace previous chrome conclusion: 2-textbox approach for clean baseline."""
    # NB: leave the rectangle from chrome in place
    pass


# Simpler chrome: header + footer label "このページの結論" inline before the message
def add_chrome(slide, page_num, total_pages, title, conclusion, dark=False):
    if dark:
        # Title slide variant: dark forest background handled outside
        return
    # top amber accent
    add_rect(slide, 0.50, 0.50, 0.08, 0.65, fill=AMBER)
    add_textbox(slide, title,
                l=0.72, t=0.45, w=11.0, h=0.75,
                size=28, color=DEEP_FOREST, bold=True, valign='middle')
    add_textbox(slide, f"{page_num} / {total_pages}",
                l=12.10, t=0.55, w=0.80, h=0.40,
                size=10, color=MUTED, align='right', valign='middle')
    add_line(slide, 0.50, 1.30, 12.83, 1.30, color=CARD_BORDER, width=1.0)
    # bottom strip
    add_rect(slide, 0.50, 6.55, 12.33, 0.65, fill=SAGE_BG)
    add_rect(slide, 0.50, 6.55, 0.08, 0.65, fill=AMBER)
    add_textbox(slide, "このページの結論",
                l=0.72, t=6.58, w=2.10, h=0.30,
                size=10, color=AMBER_DEEP, bold=True, valign='middle')
    add_textbox(slide, conclusion,
                l=0.72, t=6.88, w=12.00, h=0.30,
                size=14, color=DEEP_FOREST, bold=True, valign='middle')


TOTAL = 11


# =====================================================================
# Slide 1: 表紙
# =====================================================================
s1 = prs.slides.add_slide(BLANK)
# dark forest background
add_rect(s1, 0, 0, SW, SH, fill=DEEP_FOREST)
# amber side accent
add_rect(s1, 0.80, 1.80, 0.12, 4.20, fill=AMBER)
# kicker
add_textbox(s1, "ARCHITECTURE PROPOSAL",
            l=1.10, t=1.85, w=8.0, h=0.40,
            size=14, color=AMBER, bold=True,
            latin=EN_FONT, ea=JP_FONT)
# title
add_multiline(s1, [
    ("AWS Cross-Account",                            {'size': 40, 'bold': True, 'color': WHITE, 'space_after': 4}),
    ("Private API Gateway 呼び出し",                 {'size': 40, 'bold': True, 'color': WHITE, 'space_after': 4}),
    ("アーキテクチャ提案",                           {'size': 40, 'bold': True, 'color': WHITE}),
], l=1.10, t=2.30, w=11.0, h=2.50, line_spacing=1.15)
# subtitle
add_textbox(s1, "Account B Lambda（Cognito Trigger）から Account A Private API を呼ぶ最小構成",
            l=1.10, t=4.95, w=11.0, h=0.40,
            size=16, color=SAGE, valign='middle')
# bottom meta
add_line(s1, 1.10, 6.20, 12.20, 6.20, color=AMBER, width=1.0)
add_textbox(s1, "報告者: Ogata", l=1.10, t=6.35, w=4.0, h=0.30,
            size=12, color=WHITE, valign='middle')
add_textbox(s1, "2026-05-11", l=5.20, t=6.35, w=3.0, h=0.30,
            size=12, color=WHITE, valign='middle')
add_textbox(s1, "対象: 単一リージョン / 同一 VPC（B 側）",
            l=8.20, t=6.35, w=4.0, h=0.30,
            size=12, color=SAGE, align='right', valign='middle')


# =====================================================================
# Slide 2: 目次
# =====================================================================
s2 = prs.slides.add_slide(BLANK)
add_chrome(s2, 2, TOTAL, "目次", "本資料は背景→構成図→各アカウントの作業→認証→実装→運用の順で読み進められます。")

toc = [
    ("01", "背景 / 目的",                                  "要件・制約・ゴール"),
    ("02", "アーキテクチャ全体図",                          "VPCE 経由のプライベート完結構成"),
    ("03", "Account A がやること",                          "Resource Policy 更新・VPCE-ID 受け入れ"),
    ("04", "Account B がやること",                          "VPCE 作成・Lambda VPC アタッチ・IAM"),
    ("05", "認証方式の比較",                                "AWS_IAM / API Key / なし / Authorizer"),
    ("06", "Lambda Python サンプルコード",                  "SigV4 署名 + VPCE 固有 URL"),
    ("07", "Cognito トリガーとしての注意点",                "タイムアウト・Cold start・エラー"),
    ("08", "コスト / 監視 / 追加考慮",                      "VPCE 費用・CloudWatch・X-Ray"),
]
# 2列レイアウト
col_left_x  = 1.20
col_right_x = 7.00
row_h = 0.55
for i, (num, jp, desc) in enumerate(toc):
    col = i % 2
    row = i // 2
    x = col_left_x if col == 0 else col_right_x
    y = 1.70 + row * (row_h + 0.45)
    add_rect(s2, x, y, 5.30, row_h, fill=WHITE, line=CARD_BORDER, line_width=0.5)
    add_rect(s2, x, y, 0.06, row_h, fill=AMBER)
    add_textbox(s2, num, l=x+0.20, t=y, w=0.60, h=row_h,
                size=18, color=AMBER_DEEP, bold=True, valign='middle')
    add_textbox(s2, jp, l=x+0.85, t=y, w=4.30, h=row_h*0.55,
                size=14, color=DEEP_FOREST, bold=True, valign='middle')
    add_textbox(s2, desc, l=x+0.85, t=y+row_h*0.50, w=4.30, h=row_h*0.50,
                size=10, color=MUTED, valign='middle')


# =====================================================================
# Slide 3: 背景 / 目的
# =====================================================================
s3 = prs.slides.add_slide(BLANK)
add_chrome(s3, 3, TOTAL, "背景 / 目的",
           "Cognito 認証フローの中で、別アカウントの Private API を AWS バックボーン経由で呼ぶ最小構成を確立する。")

# 3カラム: 要件 / 制約 / ゴール
def card(slide, x, y, w, h, label, body_lines, accent=AMBER):
    add_rect(slide, x, y, w, h, fill=WHITE, line=CARD_BORDER, line_width=0.5)
    add_rect(slide, x, y, 0.08, h, fill=accent)
    add_textbox(slide, label, l=x+0.25, t=y+0.15, w=w-0.40, h=0.35,
                size=14, color=AMBER_DEEP, bold=True)
    add_multiline(slide, body_lines, l=x+0.25, t=y+0.60, w=w-0.40, h=h-0.70,
                  size=12, color=INK, line_spacing=1.25)

req_lines = [
    ("● Account A:", {'size': 12, 'bold': True, 'color': DEEP_FOREST, 'space_after': 1}),
    ("    Private API Gateway が既に存在", {'size': 12, 'space_after': 6}),
    ("● Account B:", {'size': 12, 'bold': True, 'color': DEEP_FOREST, 'space_after': 1}),
    ("    Lambda（Python）から A を呼ぶ", {'size': 12, 'space_after': 6}),
    ("● Lambda B は Cognito User Pool の", {'size': 12, 'space_after': 1}),
    ("    認証トリガー（例: Pre Auth）", {'size': 12}),
]
con_lines = [
    ("● API はインターネット非公開", {'size': 12, 'bold': True, 'color': DEEP_FOREST, 'space_after': 6}),
    ("● 通信は AWS バックボーン完結", {'size': 12, 'space_after': 6}),
    ("● 同一リージョン・既存 VPC を利用", {'size': 12, 'space_after': 6}),
    ("● Private DNS の cross-account", {'size': 12, 'space_after': 1}),
    ("    共有は前提とせず VPCE 固有 URL", {'size': 12, 'space_after': 1}),
    ("    を採用する", {'size': 12}),
]
goal_lines = [
    ("● B → A をプライベート経路で", {'size': 12, 'bold': True, 'color': DEEP_FOREST, 'space_after': 1}),
    ("    確実に到達できる構成を確立", {'size': 12, 'space_after': 6}),
    ("● 認証/認可を明確化し", {'size': 12, 'space_after': 1}),
    ("    監査可能な状態にする", {'size': 12, 'space_after': 6}),
    ("● Cognito Trigger の運用制約", {'size': 12, 'space_after': 1}),
    ("    （5秒タイムアウト等）に耐える", {'size': 12}),
]

cw = 4.00
gap = 0.15
card_x0 = 0.50 + 0.20
card_y  = 1.55
card_h  = 4.65
card(s3, card_x0,                 card_y, cw, card_h, "要件", req_lines)
card(s3, card_x0 + cw + gap,      card_y, cw, card_h, "制約", con_lines)
card(s3, card_x0 + (cw+gap)*2,    card_y, cw, card_h, "ゴール", goal_lines)


# =====================================================================
# Slide 4: アーキテクチャ全体図
# =====================================================================
s4 = prs.slides.add_slide(BLANK)
add_chrome(s4, 4, TOTAL, "アーキテクチャ全体図",
           "Lambda B → Interface VPC Endpoint → execute-api → Account A の Private API、すべて AWS 内で完結。")

# Layout: two big account boxes left/right
# Account B (left)
B_x, B_y, B_w, B_h = 0.70, 1.55, 6.20, 4.70
A_x, A_y, A_w, A_h = 7.30, 1.55, 5.50, 4.70

# B box
add_rect(s4, B_x, B_y, B_w, B_h, fill=SAGE_BG, line=FOREST, line_width=1.5)
add_rect(s4, B_x, B_y, B_w, 0.40, fill=FOREST)
add_textbox(s4, "Account B  （呼び出し側）", l=B_x+0.15, t=B_y+0.05, w=B_w-0.30, h=0.30,
            size=12, color=WHITE, bold=True, valign='middle')

# Cognito (top)
cog_x, cog_y, cog_w, cog_h = B_x+0.40, B_y+0.65, 2.40, 0.85
add_rect(s4, cog_x, cog_y, cog_w, cog_h, fill=WHITE, line=AMBER, line_width=1.5, corner=0.15)
add_textbox(s4, "Cognito User Pool", l=cog_x, t=cog_y+0.08, w=cog_w, h=0.30,
            size=12, color=DEEP_FOREST, bold=True, align='center', valign='middle')
add_textbox(s4, "Pre Authentication Trigger", l=cog_x, t=cog_y+0.40, w=cog_w, h=0.30,
            size=10, color=MUTED, align='center', valign='middle')

# VPC box inside B
vpc_x, vpc_y, vpc_w, vpc_h = B_x+0.30, B_y+1.75, B_w-0.60, B_h-2.05
add_rect(s4, vpc_x, vpc_y, vpc_w, vpc_h, fill=WHITE, line=FOREST, line_width=1.0)
add_textbox(s4, "VPC  /  Private Subnet", l=vpc_x+0.15, t=vpc_y+0.05, w=vpc_w-0.30, h=0.28,
            size=10, color=FOREST, bold=True, valign='middle')

# Lambda
lam_x, lam_y, lam_w, lam_h = vpc_x+0.35, vpc_y+0.55, 2.30, 0.90
add_rect(s4, lam_x, lam_y, lam_w, lam_h, fill=AMBER_TINT, line=AMBER_DEEP, line_width=1.2, corner=0.15)
add_textbox(s4, "Lambda  (Python)", l=lam_x, t=lam_y+0.08, w=lam_w, h=0.30,
            size=12, color=AMBER_DEEP, bold=True, align='center', valign='middle')
add_textbox(s4, "SigV4 署名 + 環境変数", l=lam_x, t=lam_y+0.40, w=lam_w, h=0.30,
            size=10, color=MUTED, align='center', valign='middle')

# VPC Endpoint
vpe_x, vpe_y, vpe_w, vpe_h = vpc_x+vpc_w-2.65, vpc_y+0.55, 2.30, 0.90
add_rect(s4, vpe_x, vpe_y, vpe_w, vpe_h, fill=WHITE, line=FOREST, line_width=1.2, corner=0.15)
add_textbox(s4, "Interface VPCE", l=vpe_x, t=vpe_y+0.08, w=vpe_w, h=0.30,
            size=12, color=FOREST, bold=True, align='center', valign='middle')
add_textbox(s4, "com.amazonaws.{region}", l=vpe_x, t=vpe_y+0.36, w=vpe_w, h=0.24,
            size=9, color=MUTED, align='center', valign='middle')
add_textbox(s4, ".execute-api", l=vpe_x, t=vpe_y+0.58, w=vpe_w, h=0.24,
            size=9, color=MUTED, align='center', valign='middle')

# IAM Role / SG note
add_textbox(s4, "Lambda 実行ロール  +  Security Group（443 outbound → VPCE）",
            l=vpc_x+0.30, t=vpc_y+vpc_h-0.45, w=vpc_w-0.60, h=0.30,
            size=10, color=MUTED, align='center', valign='middle', italic=True)

# Arrows in B
# Cognito → Lambda
add_arrow(s4, cog_x+cog_w/2, cog_y+cog_h, cog_x+cog_w/2, lam_y, color=DEEP_FOREST, width=1.8)
add_textbox(s4, "invoke", l=cog_x+cog_w/2+0.10, t=cog_y+cog_h+0.05, w=0.80, h=0.22,
            size=9, color=DEEP_FOREST, align='left', italic=True)
# Lambda → VPCE
add_arrow(s4, lam_x+lam_w, lam_y+lam_h/2, vpe_x, vpe_y+vpe_h/2, color=AMBER_DEEP, width=2.0)

# A box
add_rect(s4, A_x, A_y, A_w, A_h, fill=SAGE_BG, line=FOREST, line_width=1.5)
add_rect(s4, A_x, A_y, A_w, 0.40, fill=FOREST)
add_textbox(s4, "Account A  （API 提供側）", l=A_x+0.15, t=A_y+0.05, w=A_w-0.30, h=0.30,
            size=12, color=WHITE, bold=True, valign='middle')

# API Gateway (Private)
apigw_x, apigw_y, apigw_w, apigw_h = A_x+0.50, A_y+1.70, A_w-1.00, 1.20
add_rect(s4, apigw_x, apigw_y, apigw_w, apigw_h, fill=AMBER_TINT, line=AMBER_DEEP, line_width=1.5, corner=0.10)
add_textbox(s4, "API Gateway  (Private)", l=apigw_x, t=apigw_y+0.10, w=apigw_w, h=0.35,
            size=13, color=AMBER_DEEP, bold=True, align='center', valign='middle')
add_textbox(s4, "Endpoint Type: PRIVATE", l=apigw_x, t=apigw_y+0.50, w=apigw_w, h=0.25,
            size=10, color=MUTED, align='center', valign='middle')
add_textbox(s4, "Resource Policy で許可 Principal / VPCE-ID を制御",
            l=apigw_x, t=apigw_y+0.78, w=apigw_w, h=0.30,
            size=10, color=DEEP_FOREST, align='center', valign='middle', italic=True)

# Backend (Lambda or whatever)
bk_x, bk_y, bk_w, bk_h = A_x+0.50, A_y+3.20, A_w-1.00, 0.85
add_rect(s4, bk_x, bk_y, bk_w, bk_h, fill=WHITE, line=FOREST, line_width=1.0, corner=0.10)
add_textbox(s4, "Backend  (Lambda / その他)", l=bk_x, t=bk_y+0.10, w=bk_w, h=0.30,
            size=12, color=FOREST, bold=True, align='center', valign='middle')
add_textbox(s4, "本資料の対象外", l=bk_x, t=bk_y+0.42, w=bk_w, h=0.30,
            size=9, color=MUTED, align='center', valign='middle', italic=True)

# APIGW → Backend
add_arrow(s4, apigw_x+apigw_w/2, apigw_y+apigw_h, bk_x+bk_w/2, bk_y, color=DEEP_FOREST, width=1.5, dash=True)

# Cross-account arrow: VPCE → APIGW
add_arrow(s4, vpe_x+vpe_w, vpe_y+vpe_h/2, apigw_x, apigw_y+apigw_h/2,
          color=AMBER_DEEP, width=3.0)
# Label on the cross-account arrow
add_rect(s4, 6.55, vpe_y+vpe_h/2-0.50, 1.55, 0.45, fill=WHITE, line=AMBER, line_width=0.75, corner=0.30)
add_textbox(s4, "VPCE 固有 URL", l=6.55, t=vpe_y+vpe_h/2-0.48, w=1.55, h=0.20,
            size=9, color=AMBER_DEEP, bold=True, align='center', valign='middle')
add_textbox(s4, "SigV4 / Host 上書き", l=6.55, t=vpe_y+vpe_h/2-0.28, w=1.55, h=0.20,
            size=8, color=MUTED, align='center', valign='middle')

# Backbone label (between account boxes, on cross-account arrow band)
add_textbox(s4, "── AWS Backbone（インターネット非経由）──",
            l=4.40, t=6.28, w=4.50, h=0.22,
            size=9, color=AMBER_DEEP, align='center', valign='middle', italic=True, bold=True)


# =====================================================================
# Slide 5: Account A がやること
# =====================================================================
s5 = prs.slides.add_slide(BLANK)
add_chrome(s5, 5, TOTAL, "Account A がやること",
           "Resource Policy で「誰が（Principal）どこから（VPCE-ID）」呼べるかを制御し、認証方式に応じた追加対応を行う。")

a_tasks = [
    ("01", "Resource Policy の更新",
     "Private API の Resource Policy に Account B の VPCE-ID を許可する。\n"
     "（Condition: aws:SourceVpce = vpce-xxxxxxxx）"),
    ("02", "Principal の追加（AWS_IAM 採用時）",
     "Action: execute-api:Invoke の Principal に\n"
     "arn:aws:iam::<account-b>:role/<lambda-role> を許可。"),
    ("03", "VPCE-ID の受け入れ調整",
     "Account B 側で作成された VPCE-ID を共有してもらい、\n"
     "上記 Resource Policy へ反映する（事前合意）。"),
    ("04", "認証方式に応じた追加対応",
     "API Key を使うなら Usage Plan + Key 発行 / カスタム Authorizer なら Lambda 配備。\n"
     "認証なしの場合は Resource Policy だけが信頼境界となる点に注意。"),
]
ay = 1.55
ah = 1.10
for i, (num, head, body) in enumerate(a_tasks):
    y = ay + i * (ah + 0.08)
    add_rect(s5, 0.70, y, 12.10, ah, fill=WHITE, line=CARD_BORDER, line_width=0.5)
    add_rect(s5, 0.70, y, 0.08, ah, fill=AMBER)
    # number badge
    add_rect(s5, 0.95, y+0.18, 0.65, 0.65, fill=DEEP_FOREST, corner=0.50)
    add_textbox(s5, num, l=0.95, t=y+0.18, w=0.65, h=0.65,
                size=14, color=AMBER, bold=True, align='center', valign='middle')
    add_textbox(s5, head, l=1.75, t=y+0.08, w=10.80, h=0.38,
                size=15, color=DEEP_FOREST, bold=True, valign='middle')
    add_textbox(s5, body, l=1.75, t=y+0.50, w=10.80, h=ah-0.55,
                size=11, color=INK, valign='top', line_spacing=1.25)


# =====================================================================
# Slide 6: Account B がやること
# =====================================================================
s6 = prs.slides.add_slide(BLANK)
add_chrome(s6, 6, TOTAL, "Account B がやること",
           "VPCE 作成（Private DNS OFF）と Lambda の VPC アタッチ。SigV4 用 IAM ロールと環境変数を整える。")

b_tasks_left = [
    ("VPC / Subnet", "Lambda を配置する VPC を確認。\nプライベートサブネット ×2 以上推奨（マルチ AZ）。"),
    ("Interface VPCE 作成", "Service: com.amazonaws.{region}.execute-api\n★ Private DNS = OFF（cross-account の鉄則）"),
    ("Security Group", "VPCE 用 SG: 443 inbound from Lambda SG。\nLambda SG: 443 outbound to VPCE SG。"),
]
b_tasks_right = [
    ("Lambda VPC アタッチ", "上記サブネット + Lambda SG に紐付ける。\nNAT 不要（VPCE で完結）。"),
    ("Lambda 実行ロール（IAM）", "execute-api:Invoke を該当 API の ARN に許可。\nAWS_IAM 採用時のみ必要（他方式では不要）。"),
    ("環境変数", "API_ID / VPCE_ID / STAGE / API_PATH を設定。\nリージョンは AWS_REGION（予約済み）を使用。"),
]
def b_card(slide, x, y, w, h, head, body):
    add_rect(slide, x, y, w, h, fill=WHITE, line=CARD_BORDER, line_width=0.5)
    add_rect(slide, x, y, 0.08, h, fill=FOREST)
    add_textbox(slide, head, l=x+0.20, t=y+0.10, w=w-0.30, h=0.38,
                size=14, color=DEEP_FOREST, bold=True, valign='middle')
    add_textbox(slide, body, l=x+0.20, t=y+0.52, w=w-0.30, h=h-0.62,
                size=11, color=INK, line_spacing=1.30, valign='top')

cw2 = 5.85
ch2 = 1.45
gx = 0.70
gx2 = 0.70 + cw2 + 0.30
gy0 = 1.55
for i, (head, body) in enumerate(b_tasks_left):
    b_card(s6, gx, gy0 + i*(ch2+0.10), cw2, ch2, head, body)
for i, (head, body) in enumerate(b_tasks_right):
    b_card(s6, gx2, gy0 + i*(ch2+0.10), cw2, ch2, head, body)


# =====================================================================
# Slide 7: 認証方式の比較
# =====================================================================
s7 = prs.slides.add_slide(BLANK)
add_chrome(s7, 7, TOTAL, "認証方式の比較",
           "AWS_IAM を推奨。証跡（CloudTrail）と最小権限制御の両立、追加コンポーネント不要が決め手。")

# Table 5 cols x 5 rows (header + 4)
table_x, table_y = 0.70, 1.55
table_w = 12.10
col_widths = [2.30, 4.40, 1.80, 1.80, 1.80]  # sums to 12.10
header_h = 0.55
row_h = 0.95
headers = ["方式", "概要", "セキュリティ", "実装難易度", "監査容易性"]

# header
cx = table_x
for i, hd in enumerate(headers):
    add_rect(s7, cx, table_y, col_widths[i], header_h, fill=DEEP_FOREST)
    add_textbox(s7, hd, l=cx, t=table_y, w=col_widths[i], h=header_h,
                size=12, color=WHITE, bold=True, align='center', valign='middle')
    cx += col_widths[i]

rows = [
    {
        "name": "(a) AWS_IAM ★推奨",
        "desc": "Account A の Resource Policy + Principal に B の Lambda 実行ロール ARN を許可。Lambda 側は SigV4 署名で呼ぶ。",
        "sec": "◎ 強い",
        "impl": "○ 中程度",
        "audit": "◎ CloudTrail",
        "highlight": True,
    },
    {
        "name": "(b) API Key",
        "desc": "API Gateway の Usage Plan + API Key。Lambda は x-api-key ヘッダで呼ぶ。Key の Secrets Manager 管理が前提。",
        "sec": "△ Key 漏洩リスク",
        "impl": "◎ 容易",
        "audit": "○ 利用計画ログ",
        "highlight": False,
    },
    {
        "name": "(c) 認証なし",
        "desc": "Resource Policy（VPCE-ID + Principal）だけで信頼。VPC 越境がない前提で運用できるが、A の Resource Policy のみが境界となる。",
        "sec": "△ 単一防御層",
        "impl": "◎ 最小",
        "audit": "△ 識別性低",
        "highlight": False,
    },
    {
        "name": "(d) カスタム Lambda Authorizer",
        "desc": "Account A 側に Authorizer Lambda を配置し、独自トークンや JWT を検証。柔軟だが運用負荷が増える。",
        "sec": "○ 設計依存",
        "impl": "△ 高い",
        "audit": "○ 自前ログ",
        "highlight": False,
    },
]

row_y = table_y + header_h
for r in rows:
    cx = table_x
    fill = AMBER_TINT if r["highlight"] else WHITE
    # row backgrounds per column
    cells = [r["name"], r["desc"], r["sec"], r["impl"], r["audit"]]
    for i, val in enumerate(cells):
        add_rect(s7, cx, row_y, col_widths[i], row_h, fill=fill, line=CARD_BORDER, line_width=0.4)
        if i == 0:
            color = AMBER_DEEP if r["highlight"] else DEEP_FOREST
            add_textbox(s7, val, l=cx+0.12, t=row_y, w=col_widths[i]-0.24, h=row_h,
                        size=12, color=color, bold=True, valign='middle')
        elif i == 1:
            add_textbox(s7, val, l=cx+0.15, t=row_y+0.05, w=col_widths[i]-0.30, h=row_h-0.10,
                        size=10, color=INK, valign='middle', line_spacing=1.25)
        else:
            add_textbox(s7, val, l=cx, t=row_y, w=col_widths[i], h=row_h,
                        size=11, color=DEEP_FOREST, bold=True, align='center', valign='middle')
        cx += col_widths[i]
    row_y += row_h

# Judgement criteria caption
add_textbox(s7, "判断基準: ① セキュリティ（漏洩耐性・最小権限）  ② 実装容易性（Lambda コードの薄さ）  ③ 監査容易性（誰が呼んだかを後追いできるか）",
            l=table_x, t=row_y+0.15, w=table_w, h=0.30,
            size=10, color=MUTED, valign='middle', italic=True)


# =====================================================================
# Slide 8: Lambda Python サンプルコード
# =====================================================================
s8 = prs.slides.add_slide(BLANK)
add_chrome(s8, 8, TOTAL, "Lambda Python サンプルコード",
           "VPCE 固有 URL + SigV4 署名 + Host ヘッダ上書きの 3 点が要点。event をそのまま返せば認証継続。")

# Code panel
cx, cy, cw_, ch_ = 0.70, 1.55, 12.10, 4.35
add_rect(s8, cx, cy, cw_, ch_, fill=CODE_BG, line=DEEP_FOREST, line_width=0.5, corner=0.02)
# code title bar
add_rect(s8, cx, cy, cw_, 0.32, fill=DEEP_FOREST)
add_textbox(s8, "  lambda_function.py    —    Cognito Pre Authentication Trigger",
            l=cx, t=cy, w=cw_, h=0.32,
            size=10, color=AMBER, bold=True, valign='middle', latin=MONO)

# Code body (slightly trimmed for slide fit; full version in build script comment)
code_lines = [
    ("import os, json, boto3, urllib3",                                  CODE_INK),
    ("from botocore.auth import SigV4Auth",                              CODE_INK),
    ("from botocore.awsrequest import AWSRequest",                       CODE_INK),
    ("",                                                                  CODE_INK),
    ("http = urllib3.PoolManager()",                                     CODE_INK),
    ("API_ID  = os.environ[\"API_ID\"]",                                 CODE_INK),
    ("VPCE_ID = os.environ[\"VPCE_ID\"]",                                CODE_INK),
    ("REGION  = os.environ[\"AWS_REGION\"]",                             CODE_INK),
    ("STAGE   = os.environ[\"STAGE\"]",                                  CODE_INK),
    ("API_PATH = os.environ.get(\"API_PATH\", \"/users\")",              CODE_INK),
    ("",                                                                  CODE_INK),
    ("def _signed_request(method, url, body=\"\"):",                     CODE_KEY),
    ("    creds = boto3.Session().get_credentials().get_frozen_credentials()", CODE_INK),
    ("    req = AWSRequest(method=method, url=url, data=body)",          CODE_INK),
    ("    SigV4Auth(creds, \"execute-api\", REGION).add_auth(req)",      CODE_INK),
    ("    return dict(req.headers.items())",                             CODE_INK),
    ("",                                                                  CODE_INK),
    ("def lambda_handler(event, context):",                              CODE_KEY),
    ("    user_email = event[\"request\"][\"userAttributes\"].get(\"email\", \"unknown\")", CODE_INK),
    ("    url = (f\"https://{API_ID}-{VPCE_ID}.execute-api.{REGION}\"", CODE_INK),
    ("           f\".amazonaws.com/{STAGE}{API_PATH}?email={user_email}\")", CODE_INK),
    ("    headers = _signed_request(\"GET\", url)",                      CODE_INK),
    ("    # API Gateway は Host で API 識別 → 標準 URL の Host で上書き",   CODE_COMMENT),
    ("    headers[\"Host\"] = f\"{API_ID}.execute-api.{REGION}.amazonaws.com\"", CODE_INK),
    ("    resp = http.request(\"GET\", url, headers=headers, timeout=urllib3.Timeout(3.0))", CODE_INK),
    ("    if resp.status != 200:",                                        CODE_INK),
    ("        raise Exception(f\"upstream {resp.status}: {resp.data[:200]}\")", CODE_INK),
    ("    return event   # Cognito は event を返すと認証継続",             CODE_COMMENT),
]

# Render line-by-line
code_x = cx + 0.20
code_y = cy + 0.42
code_w = cw_ - 0.40
line_h = 0.135
code_size = 9
for i, (txt, col) in enumerate(code_lines):
    add_textbox(s8, txt if txt else " ",
                l=code_x, t=code_y + i*line_h, w=code_w, h=line_h+0.05,
                size=code_size, color=col, latin=MONO, ea=MONO, valign='top')

# right caption (positioned between code panel and bottom strip)
add_textbox(s8, "★ 3 点セット:  ① VPCE 固有 URL   ② SigV4 署名   ③ Host ヘッダの上書き",
            l=0.70, t=6.00, w=12.10, h=0.32,
            size=11, color=AMBER_DEEP, bold=True, align='center', valign='middle')


# =====================================================================
# Slide 9: Cognito トリガーとしての注意点
# =====================================================================
s9 = prs.slides.add_slide(BLANK)
add_chrome(s9, 9, TOTAL, "Cognito トリガーとしての注意点",
           "5 秒のハードタイムアウトを軸に、Cold start とエラー伝播の方針を先に決めてから実装する。")

notes = [
    ("01", "タイムアウト 5 秒の壁",
     "Cognito Trigger は 5 秒以内に Lambda が応答しないと認証失敗。\n"
     "API 呼び出しは 3 秒タイムアウト（urllib3）で打ち切り、再試行 1 回までに留める。"),
    ("02", "Cold start 対策",
     "VPC アタッチ Lambda は初回起動が遅い。Provisioned Concurrency を最小 1 で常駐させる、\n"
     "またはランタイムを軽量化（不要ライブラリを除く）。"),
    ("03", "エラーハンドリング方針の事前合意",
     "Trigger で例外を投げる = 認証拒否。要件次第で「ログを残して握り潰し→event 返却」も選択肢。\n"
     "失敗時の挙動はビジネス判断（拒否 / 通過）として事前に文書化する。"),
    ("04", "冪等性 / 副作用の整理",
     "Pre Auth は同一セッションで再実行されうる。書き込み系を呼ぶ場合は idempotency key で重複を防ぐ。"),
]
ny = 1.55
nh = 1.10
for i, (num, head, body) in enumerate(notes):
    y = ny + i*(nh + 0.08)
    add_rect(s9, 0.70, y, 12.10, nh, fill=WHITE, line=CARD_BORDER, line_width=0.5)
    add_rect(s9, 0.70, y, 0.08, nh, fill=AMBER)
    # number badge
    add_rect(s9, 0.95, y+0.18, 0.75, 0.75, fill=DEEP_FOREST, corner=0.50)
    add_textbox(s9, num, l=0.95, t=y+0.18, w=0.75, h=0.75,
                size=14, color=AMBER, bold=True, align='center', valign='middle')
    add_textbox(s9, head, l=1.85, t=y+0.08, w=10.80, h=0.38,
                size=15, color=DEEP_FOREST, bold=True, valign='middle')
    add_textbox(s9, body, l=1.85, t=y+0.50, w=10.80, h=nh-0.55,
                size=11, color=INK, valign='top', line_spacing=1.30)


# =====================================================================
# Slide 10: コスト / 監視 / 追加考慮
# =====================================================================
s10 = prs.slides.add_slide(BLANK)
add_chrome(s10, 10, TOTAL, "コスト / 監視 / 追加考慮",
           "VPCE の常時課金を許容できることが構成成立の前提。監視はメトリクス＋構造化ログ＋トレースの三段構え。")

items_left = [
    ("コスト",
     "● Interface VPCE: AZ × 時間 + データ処理料金\n"
     "● 同一 AZ 内呼び出しでデータ転送は最小化\n"
     "● Lambda: VPC アタッチによる ENI 維持コスト"),
    ("監視",
     "● CloudWatch Logs (Lambda) → エラーレート / レイテンシ\n"
     "● API Gateway 4xx/5xx メトリクス\n"
     "● VPC Flow Logs（Reject トラフィック検知）"),
]
items_right = [
    ("トレース",
     "● X-Ray 有効化（Lambda + API GW）\n"
     "● サブセグメント: SigV4 署名 / HTTP 呼び出し\n"
     "● 相関 ID は Cognito event['request']['userAttributes'] から付与"),
    ("フォールバック / 拡張",
     "● 失敗時挙動の合意（拒否 / 通過）を IaC コメントで明記\n"
     "● Account A 側が同一 VPC 内なら VPCE は省略可\n"
     "● 将来クロスリージョンに広げる場合は Private DNS 戦略を再検討"),
]

def info_card(slide, x, y, w, h, head, body):
    add_rect(slide, x, y, w, h, fill=WHITE, line=CARD_BORDER, line_width=0.5)
    add_rect(slide, x, y, 0.08, h, fill=FOREST)
    add_textbox(slide, head, l=x+0.25, t=y+0.10, w=w-0.40, h=0.42,
                size=15, color=DEEP_FOREST, bold=True, valign='middle')
    add_textbox(slide, body, l=x+0.25, t=y+0.58, w=w-0.40, h=h-0.70,
                size=11, color=INK, valign='top', line_spacing=1.35)

cw3 = 5.85
ch3 = 2.25
gx = 0.70
gx2 = 0.70 + cw3 + 0.30
gy0 = 1.55
for i, (head, body) in enumerate(items_left):
    info_card(s10, gx, gy0 + i*(ch3+0.18), cw3, ch3, head, body)
for i, (head, body) in enumerate(items_right):
    info_card(s10, gx2, gy0 + i*(ch3+0.18), cw3, ch3, head, body)


# =====================================================================
# Slide 11: Appendix — VPC Peering の要否判断
# =====================================================================
s11 = prs.slides.add_slide(BLANK)
# Custom chrome — title on one line at slightly smaller size to fit subtitle
add_rect(s11, 0.50, 0.50, 0.08, 0.65, fill=AMBER)
# Title (main) - bold deep forest
add_textbox(s11, "Appendix: VPC Peering の要否判断",
            l=0.72, t=0.45, w=6.40, h=0.75,
            size=24, color=DEEP_FOREST, bold=True, valign='middle')
# Title (subtitle/dash) - muted, smaller
add_textbox(s11, "—  今回は不要、根拠を 3 点で示す",
            l=7.10, t=0.50, w=4.90, h=0.65,
            size=15, color=AMBER_DEEP, italic=True, valign='middle')
# page no
add_textbox(s11, f"{11} / {TOTAL}",
            l=12.10, t=0.55, w=0.80, h=0.40,
            size=10, color=MUTED, align='right', valign='middle')
add_line(s11, 0.50, 1.30, 12.83, 1.30, color=CARD_BORDER, width=1.0)
# Bottom conclusion strip
add_rect(s11, 0.50, 6.55, 12.33, 0.65, fill=SAGE_BG)
add_rect(s11, 0.50, 6.55, 0.08, 0.65, fill=AMBER)
add_textbox(s11, "このページの結論",
            l=0.72, t=6.58, w=2.10, h=0.30,
            size=10, color=AMBER_DEEP, bold=True, valign='middle')
add_textbox(s11, "Private API GW を叩くだけなら VPCE で完結。VPC Peering は将来 A の VPC 内リソース直接接続が必要になった時に再検討。",
            l=0.72, t=6.88, w=12.00, h=0.30,
            size=13, color=DEEP_FOREST, bold=True, valign='middle')

# --------- LEFT half: 「なぜ不要なのか（通信経路の本質）」 ---------
L_x, L_y, L_w, L_h = 0.70, 1.55, 5.95, 4.85
add_rect(s11, L_x, L_y, L_w, L_h, fill=WHITE, line=CARD_BORDER, line_width=0.5)
add_rect(s11, L_x, L_y, 0.08, L_h, fill=AMBER)
add_textbox(s11, "なぜ不要なのか（通信経路の本質）",
            l=L_x+0.25, t=L_y+0.12, w=L_w-0.40, h=0.40,
            size=14, color=AMBER_DEEP, bold=True, valign='middle')

# Diagram nodes - vertical stack of 5 boxes connecting top→bottom
node_w, node_h = L_w-1.00, 0.50
node_x = L_x + 0.50
nodes = [
    ("Lambda  (Account B)",                AMBER_TINT, AMBER_DEEP),
    ("VPCE  (B / execute-api)",            WHITE,      FOREST),
    ("AWS バックボーン",                   SAGE_BG,    DEEP_FOREST),
    ("API Gateway サービス（AWS マネージド）", WHITE,   FOREST),
    ("Account A  の  Private API",         AMBER_TINT, AMBER_DEEP),
]
y_top = L_y + 0.65
gap = 0.18
for i, (label, bg, fg) in enumerate(nodes):
    y = y_top + i * (node_h + gap)
    add_rect(s11, node_x, y, node_w, node_h, fill=bg, line=fg, line_width=1.0, corner=0.15)
    add_textbox(s11, label, l=node_x, t=y, w=node_w, h=node_h,
                size=11, color=fg, bold=True, align='center', valign='middle')
    # arrow between i and i+1
    if i < len(nodes)-1:
        arrow_y1 = y + node_h
        arrow_y2 = arrow_y1 + gap
        add_arrow(s11, node_x + node_w/2, arrow_y1,
                  node_x + node_w/2, arrow_y2, color=DEEP_FOREST, width=1.5)

# Key message callout below diagram
diag_bottom = y_top + len(nodes)*node_h + (len(nodes)-1)*gap  # ≈ 5.42
key_y = diag_bottom + 0.05
key_h = 0.90
add_rect(s11, L_x+0.25, key_y, L_w-0.50, key_h, fill=AMBER_TINT, line=AMBER_DEEP, line_width=0.8, corner=0.10)
add_multiline(s11, [
    ("VPCE は AWS サービスへの専用入口。",
                                                    {'size': 11, 'bold': True, 'color': AMBER_DEEP, 'space_after': 1}),
    ("相手 VPC への接続線ではない。",
                                                    {'size': 11, 'bold': True, 'color': AMBER_DEEP, 'space_after': 3}),
    ("Private API は AWS マネージドの execute-api 上にホスト。",
                                                    {'size': 9, 'color': MUTED, 'italic': True}),
], l=L_x+0.40, t=key_y+0.08, w=L_w-0.80, h=key_h-0.16, line_spacing=1.10)

# --------- RIGHT TOP: 「VPC Peering が必要なケース vs 不要なケース」 ---------
R_x = 6.85
RT_y = 1.55
RT_w = 5.95
RT_h = 2.25
add_rect(s11, R_x, RT_y, RT_w, RT_h, fill=WHITE, line=CARD_BORDER, line_width=0.5)
add_rect(s11, R_x, RT_y, 0.08, RT_h, fill=AMBER)
add_textbox(s11, "VPC Peering が必要 vs 不要",
            l=R_x+0.25, t=RT_y+0.10, w=RT_w-0.40, h=0.34,
            size=14, color=AMBER_DEEP, bold=True, valign='middle')

# table: 2 columns
tbl_x = R_x + 0.25
tbl_y = RT_y + 0.52
col1_w = 3.85
col2_w = 1.75
row_h11 = 0.38
hdr_h11 = 0.30

# header
add_rect(s11, tbl_x,         tbl_y, col1_w, hdr_h11, fill=DEEP_FOREST)
add_rect(s11, tbl_x+col1_w,  tbl_y, col2_w, hdr_h11, fill=DEEP_FOREST)
add_textbox(s11, "  ケース", l=tbl_x, t=tbl_y, w=col1_w, h=hdr_h11,
            size=10, color=WHITE, bold=True, valign='middle')
add_textbox(s11, "必要？", l=tbl_x+col1_w, t=tbl_y, w=col2_w, h=hdr_h11,
            size=10, color=WHITE, bold=True, align='center', valign='middle')

rows11 = [
    ("Private API GW を叩く（今回）", "不要 ★", True),
    ("A の VPC 内 RDS / EC2 直接接続", "必要", False),
    ("大量サービス・複数アカウント", "TGW / PrivateLink 推奨", False),
    ("Private DNS 共有のため", "不要（RAM で R53 共有）", False),
]

ry = tbl_y + hdr_h11
for case, need, hl in rows11:
    fill = AMBER_TINT if hl else WHITE
    add_rect(s11, tbl_x,        ry, col1_w, row_h11, fill=fill, line=CARD_BORDER, line_width=0.4)
    add_rect(s11, tbl_x+col1_w, ry, col2_w, row_h11, fill=fill, line=CARD_BORDER, line_width=0.4)
    color = AMBER_DEEP if hl else INK
    bold = hl
    add_textbox(s11, case, l=tbl_x+0.10, t=ry, w=col1_w-0.15, h=row_h11,
                size=10, color=color, bold=bold, valign='middle')
    add_textbox(s11, need, l=tbl_x+col1_w, t=ry, w=col2_w, h=row_h11,
                size=10, color=color, bold=bold, align='center', valign='middle')
    ry += row_h11

# --------- RIGHT BOTTOM: Pros vs Cons ---------
RB_y = 4.00
RB_h = 2.45
# two columns side by side, each half of right half
half_w = (RT_w - 0.15) / 2.0
add_rect(s11, R_x,            RB_y, half_w, RB_h, fill=SAGE_BG, line=FOREST, line_width=0.8)
add_rect(s11, R_x+half_w+0.15, RB_y, half_w, RB_h, fill=AMBER_TINT, line=AMBER_DEEP, line_width=0.8)

# Pros header
add_textbox(s11, "Pros（Peering なし）",
            l=R_x+0.20, t=RB_y+0.12, w=half_w-0.30, h=0.34,
            size=12, color=DEEP_FOREST, bold=True, valign='middle')
pros = [
    "● 構成シンプル",
    "● CIDR 重複リスク回避",
    "● 最小権限（VPCE-ID で縛れる）",
]
for i, line in enumerate(pros):
    add_textbox(s11, line,
                l=R_x+0.20, t=RB_y+0.50+i*0.40, w=half_w-0.30, h=0.36,
                size=11, color=INK, valign='middle')

# Cons header
add_textbox(s11, "Cons（Peering ありの負荷）",
            l=R_x+half_w+0.35, t=RB_y+0.12, w=half_w-0.30, h=0.34,
            size=12, color=AMBER_DEEP, bold=True, valign='middle')
cons = [
    "● 両アカウント網設定",
    "● CIDR 重複 NG",
    "● ブラスト半径拡大",
    "● クロス AZ データ転送コスト",
]
for i, line in enumerate(cons):
    add_textbox(s11, line,
                l=R_x+half_w+0.35, t=RB_y+0.50+i*0.40, w=half_w-0.40, h=0.36,
                size=11, color=INK, valign='middle')


# ---------- save ----------
out_path = "/sessions/exciting-eloquent-clarke/mnt/_outputs/2026-05-11_AWS_CrossAccount_PrivateAPI_proposal.pptx"
prs.save(out_path)
print(f"saved: {out_path}")
