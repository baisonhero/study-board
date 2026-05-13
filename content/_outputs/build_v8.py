#!/usr/bin/env python3
"""
Build V8 of Claude Code 月次利用量レポート (A4 portrait, single-page dashboard).

Changes from V7 (per user feedback):
  1) Remove all Opus/Sonnet/Haiku model-ratio sections (chart + Opus distribution panel)
  2) Replace single TOP3 with three TOP3 panels: コスト / 生成行 / 効率
  3) Remove "📝 チーム代表まとめ" panel (and its sub-fields)
  4) Re-design team mapping:
       - one focused message ("活用優等生 vs 未活用 — 導入支援余地")
       - axes: x = Claude利用量(相対), y = 効率(LOC/コスト)
       - team-name labels directly on bubbles
       - quadrants explicitly named with one-line meanings
       - 4 colors only (own team / 優等生 / 標準 / 要支援)

House-style override per user:
  - Calm Forest base + Amber accent #F4B942
  - Meiryo (JP) + Segoe UI (EN/numbers)
  - Top title 28pt, bottom conclusion 18pt non-bold
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# matplotlib for the dual-axis monthly trend chart (rendered to PNG, embedded as picture)
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import MaxNLocator
from matplotlib import font_manager as fm

# Per-glyph font fallback (matplotlib >=3.6): set font.family to a LIST directly.
# DejaVu Sans handles Latin/digits/punctuation; Droid Sans Fallback covers CJK.
# Note: setting font.sans-serif (alone) does NOT enable per-glyph fallback —
# we have to pass the list to font.family itself.
plt.rcParams["font.family"] = [
    "DejaVu Sans",          # Latin & digits & punctuation
    "Droid Sans Fallback",  # CJK fallback (Linux sandbox)
    "Noto Sans CJK JP",     # preferred CJK if available
    "Hiragino Sans", "Yu Gothic", "Meiryo",
]
plt.rcParams["axes.unicode_minus"] = False

# ---------- House-style: Calm Forest + Amber ----------
DEEP_FOREST  = RGBColor(0x1F, 0x3A, 0x2E)  # primary
FOREST       = RGBColor(0x2E, 0x6F, 0x40)  # secondary
SAGE         = RGBColor(0xB5, 0xC9, 0xA8)  # tertiary (light fill)
SAGE_BG      = RGBColor(0xE8, 0xEF, 0xE3)  # very light card bg
AMBER        = RGBColor(0xF4, 0xB9, 0x42)  # accent
AMBER_DEEP   = RGBColor(0xC8, 0x8E, 0x1F)  # accent text
INK          = RGBColor(0x22, 0x28, 0x24)  # body
MUTED        = RGBColor(0x6B, 0x6B, 0x6B)  # caption
CARD_BORDER  = RGBColor(0xD7, 0xDC, 0xCF)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
RED          = RGBColor(0xC0, 0x39, 0x2B)  # for ↑ negative ratios
GREEN_OK     = RGBColor(0x2E, 0x6F, 0x40)  # for ↑ positive ratios

# bubble colors (2 colors only, per user revision)
COLOR_OWN    = AMBER                     # 当チーム — focal point
COLOR_OTHER  = RGBColor(0xA8, 0xB8, 0xA0)  # Sage Green — all 13 other teams (uniform)

JP_FONT  = "メイリオ"
EN_FONT  = "Segoe UI"

# ---------- helpers ----------
def set_run_font(run, latin=EN_FONT, ea=JP_FONT, size=12, color=INK, bold=False):
    run.font.name = latin
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea_elem = etree.SubElement(rPr, qn('a:ea'))
    ea_elem.set('typeface', ea)

def add_textbox(slide, text, l, t, w, h, size=12, color=INK, bold=False,
                align='left', valign='top', latin=EN_FONT, ea=JP_FONT,
                line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if valign == 'top':
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif valign == 'middle':
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == 'bottom':
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    if align == 'left':
        p.alignment = PP_ALIGN.LEFT
    elif align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    set_run_font(r, latin=latin, ea=ea, size=size, color=color, bold=bold)
    return tb

def add_rect(slide, l, t, w, h, fill=WHITE, line=None, line_width=None,
             corner=None):
    if corner:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(l), Inches(t), Inches(w), Inches(h))
        # set corner radius via XML adjust
        try:
            shp.adjustments[0] = corner
        except Exception:
            pass
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        if line_width:
            shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_oval(slide, l, t, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_line(slide, x1, y1, x2, y2, color=CARD_BORDER, width=0.75):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line

# ---------- slide setup (A4 portrait) ----------
prs = Presentation()
prs.slide_width  = Inches(8.27)
prs.slide_height = Inches(11.69)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# white background (default), draw nothing

# =====================================================================
# 1) HEADER (top)  T: 0.25 - 0.95
# =====================================================================
# left amber accent bar
add_rect(slide, 0.30, 0.30, 0.06, 0.50, fill=AMBER)

add_textbox(slide, "Claude Code 月次利用量レポート",
            l=0.42, t=0.20, w=6.30, h=0.45,
            size=22, color=DEEP_FOREST, bold=True)

# right-aligned version pill
add_rect(slide, 6.95, 0.24, 1.02, 0.34, fill=DEEP_FOREST, corner=0.30)
add_textbox(slide, "V8",
            l=6.95, t=0.27, w=1.02, h=0.28,
            size=14, color=AMBER, bold=True, align='center', valign='middle')

add_textbox(slide, "2026年3月度  /  報告者: Ogata  /  報告日: 2026-05-08  /  対象期間: 2026/03/01 - 03/31",
            l=0.42, t=0.65, w=7.50, h=0.20,
            size=9, color=MUTED)

# divider
add_line(slide, 0.30, 0.92, 7.97, 0.92, color=CARD_BORDER, width=1.0)

# =====================================================================
# 2) JUDGE box (left) + 年度末予測 (right)   T: 1.00 - 2.85
# =====================================================================
add_textbox(slide, "今月のジャッジ",
            l=0.30, t=1.00, w=4.70, h=0.22,
            size=14, color=DEEP_FOREST, bold=True)
add_textbox(slide, "数字から見えること — チームでの議論材料",
            l=0.30, t=1.20, w=4.70, h=0.18,
            size=9, color=MUTED)

judge_items = [
    ("✓", GREEN_OK,
     "効率性は改善傾向 (継続観測中)",
     "過去6ヶ月でコスト2.3倍に対し採用行数3.7倍"),
    ("▶", AMBER_DEEP,
     "来期の投資規模: ミドル6.6人月相当",
     "ミドル増員 vs Claude活用拡大、議論事項"),
    ("⚠", RED,
     "即対応: 認証基盤PJのモデル選定見直し",
     "局所的な過剰利用の兆候 (詳細は別ダッシュボード)"),
]
for i, (icon, icon_color, head, sub) in enumerate(judge_items):
    y = 1.42 + i * 0.47
    add_rect(slide, 0.30, y, 4.70, 0.43, fill=SAGE_BG, line=CARD_BORDER, line_width=0.5)
    add_rect(slide, 0.30, y, 0.05, 0.43, fill=icon_color)
    add_textbox(slide, icon, l=0.43, t=y+0.06, w=0.32, h=0.32,
                size=18, color=icon_color, bold=True, align='center', valign='middle')
    add_textbox(slide, head, l=0.82, t=y+0.03, w=4.10, h=0.21,
                size=12, color=DEEP_FOREST, bold=True)
    add_textbox(slide, sub, l=0.82, t=y+0.22, w=4.10, h=0.19,
                size=9, color=INK)

# right: 年度末着地予測
add_rect(slide, 5.10, 1.42, 2.87, 1.45, fill=DEEP_FOREST, line=None)
add_textbox(slide, "年度末着地予測",
            l=5.20, t=1.48, w=2.67, h=0.20,
            size=11, color=SAGE, bold=True)
add_textbox(slide, "¥9.9M",
            l=5.10, t=1.66, w=2.87, h=0.55,
            size=40, color=AMBER, bold=True, align='center')
add_line(slide, 5.40, 2.27, 7.67, 2.27, color=AMBER, width=0.5)
add_textbox(slide, "≈ ミドル人材 6.6人月分",
            l=5.10, t=2.30, w=2.87, h=0.22,
            size=12, color=WHITE, bold=True, align='center')
add_textbox(slide, "¥1.5M/月換算 / 現行ペース×12ヶ月",
            l=5.10, t=2.51, w=2.87, h=0.18,
            size=9, color=SAGE, align='center')
add_textbox(slide, "投資判断: 人を雇うか / Claude活用拡大か",
            l=5.10, t=2.69, w=2.87, h=0.18,
            size=9, color=AMBER, align='center', bold=True)

# =====================================================================
# 3) KPI TILES   T: 3.00 - 3.78
# =====================================================================
kpi = [
    ("総コスト",        "¥822k",   "+19.0%",  "vs 前月 ¥691k",    AMBER_DEEP),
    ("アクティブユーザー", "12名",    "+2名",    "vs 前月 10名",     GREEN_OK),
    ("1人あたり月額",     "¥68.5k",  "−1.2%",   "vs 前月 ¥69.1k",   GREEN_OK),
    ("1人あたり採用行数",  "2,180行", "+18%",    "全社平均 1,520行",  GREEN_OK),
]
tile_w = 1.80
tile_gap = 0.16
start_x = 0.30
for i, (label, big, delta, sub, dc) in enumerate(kpi):
    x = start_x + i * (tile_w + tile_gap)
    add_rect(slide, x, 3.00, tile_w, 0.78, fill=WHITE, line=CARD_BORDER, line_width=0.5)
    add_rect(slide, x, 3.00, tile_w, 0.04, fill=DEEP_FOREST)
    add_textbox(slide, label, l=x+0.08, t=3.05, w=tile_w-0.16, h=0.17,
                size=9, color=MUTED)
    add_textbox(slide, big, l=x+0.08, t=3.20, w=tile_w-0.16, h=0.30,
                size=19, color=DEEP_FOREST, bold=True)
    add_textbox(slide, delta, l=x+0.08, t=3.50, w=0.99, h=0.16,
                size=9, color=dc, bold=True)
    add_textbox(slide, sub, l=x+0.08, t=3.62, w=tile_w-0.16, h=0.14,
                size=8, color=MUTED)

# =====================================================================
# 4) MONTHLY TREND CHART  T: 3.92 - 5.50
# =====================================================================
add_textbox(slide, "月次推移 — コスト × 採用コード行数 (過去6ヶ月)",
            l=0.30, t=3.92, w=5.20, h=0.22,
            size=12, color=DEEP_FOREST, bold=True)
add_textbox(slide, "コスト 2.3倍 vs 採用行数 3.7倍 = 効率改善継続",
            l=5.50, t=3.92, w=2.47, h=0.22,
            size=10, color=AMBER_DEEP, bold=True, align='right')

# --- Render dual-axis chart with matplotlib (bar=cost on left, line=LOC on right) ---
cats = ['25/10', '25/11', '25/12', '26/01', '26/02', '26/03']
cost_k  = [357, 420, 510, 612, 691, 822]      # 千円
loc_k   = [7.0, 10.5, 12.0, 17.5, 22.0, 26.0]  # 千行

forest_hex = "#1F3A2E"
amber_hex  = "#F4B942"
muted_hex  = "#6B6B6B"
border_hex = "#D7DCCF"

# Match the slide's chart placeholder area: 7.67 x 1.33 inches
chart_w_in, chart_h_in = 7.67, 1.33
fig, ax_left = plt.subplots(figsize=(chart_w_in, chart_h_in), dpi=200)
fig.patch.set_facecolor("white")

# Left Y axis: cost as bars
bars = ax_left.bar(cats, cost_k, color=forest_hex, width=0.55,
                   label="コスト [千円]", zorder=2)
ax_left.set_ylabel("コスト [千円]", color=forest_hex, fontsize=7.5, labelpad=2)
ax_left.tick_params(axis="y", labelcolor=forest_hex, labelsize=7,
                    length=2, pad=1)
ax_left.tick_params(axis="x", labelsize=7, colors=muted_hex,
                    length=2, pad=1)
ax_left.grid(True, axis="y", linestyle="--", linewidth=0.4,
             color=muted_hex, alpha=0.35, zorder=1)
ax_left.set_axisbelow(True)
for spine in ("top", "right"):
    ax_left.spines[spine].set_visible(False)
ax_left.spines["left"].set_color(forest_hex)
ax_left.spines["left"].set_linewidth(0.6)
ax_left.spines["bottom"].set_color(muted_hex)
ax_left.spines["bottom"].set_linewidth(0.6)
ax_left.yaxis.set_major_locator(MaxNLocator(nbins=5, integer=True))
ax_left.set_ylim(0, max(cost_k) * 1.18)

# Right Y axis: LOC as line + markers (Amber)
ax_right = ax_left.twinx()
line = ax_right.plot(cats, loc_k, color=amber_hex, marker="o",
                     markersize=5, markeredgecolor=amber_hex,
                     markerfacecolor=amber_hex, linewidth=1.8,
                     label="生成行数 [千行]", zorder=3)
ax_right.set_ylabel("生成行数 [千行]", color=amber_hex,
                    fontsize=7.5, labelpad=2)
ax_right.tick_params(axis="y", labelcolor=amber_hex, labelsize=7,
                     length=2, pad=1)
for spine in ("top", "left"):
    ax_right.spines[spine].set_visible(False)
ax_right.spines["right"].set_color(amber_hex)
ax_right.spines["right"].set_linewidth(0.6)
ax_right.set_ylim(0, max(loc_k) * 1.25)

# Combined legend in upper-left
handles = [bars, line[0]]
labels  = ["コスト [千円]", "生成行数 [千行]"]
ax_left.legend(handles, labels, loc="upper left",
               fontsize=7, frameon=False, handlelength=1.6,
               handletextpad=0.4, borderpad=0.0,
               labelcolor=[forest_hex, amber_hex])

# Tight margins so the rendered PNG fills the picture box
plt.subplots_adjust(left=0.07, right=0.93, top=0.93, bottom=0.16)

mpl_chart_png = "/tmp/v8_monthly_trend.png"
fig.savefig(mpl_chart_png, dpi=200, facecolor="white", edgecolor="none")
plt.close(fig)

# Embed the rendered PNG as a picture in the same slot the native chart was in
slide.shapes.add_picture(mpl_chart_png,
                         Inches(0.30), Inches(4.17),
                         width=Inches(7.67), height=Inches(1.33))

# =====================================================================
# 5) TEAM MAPPING (REDESIGNED, larger, full-width)  T: 5.65 - 8.50
# =====================================================================
# ---- panel ----
add_rect(slide, 0.30, 5.65, 7.67, 2.85, fill=WHITE, line=CARD_BORDER, line_width=0.5)
# title
add_textbox(slide,
            "全社14チームのポジショニング — 当チームは右上の理想ゾーン、下位3チームに導入支援余地",
            l=0.40, t=5.70, w=7.50, h=0.22,
            size=11, color=DEEP_FOREST, bold=True)
add_textbox(slide,
            "X軸: Claude利用量 (相対) →   /   Y軸: 効率 [採用行数 ÷ コスト] ↑   /   バブル径: 1人あたりコスト",
            l=0.40, t=5.92, w=7.50, h=0.16,
            size=8, color=MUTED)

# ---- plot box ----
PX, PY = 1.10, 6.15      # plot origin (top-left of plot area)
PW, PH = 5.20, 2.05      # plot dimensions (inches)
# axis lines + arrows (forest)
add_line(slide, PX, PY+PH, PX+PW, PY+PH, color=DEEP_FOREST, width=1.0)  # x
add_line(slide, PX, PY,    PX,     PY+PH, color=DEEP_FOREST, width=1.0)  # y
# 4-quadrant grid lines (light)
add_line(slide, PX+PW/2, PY, PX+PW/2, PY+PH, color=CARD_BORDER, width=0.5)
add_line(slide, PX, PY+PH/2, PX+PW, PY+PH/2, color=CARD_BORDER, width=0.5)

# axis labels
add_textbox(slide, "→ Claude利用量 (大)",
            l=PX+PW-1.50, t=PY+PH+0.05, w=1.50, h=0.18,
            size=8, color=DEEP_FOREST, bold=True, align='right')
add_textbox(slide, "(小)",
            l=PX-0.05, t=PY+PH+0.05, w=0.30, h=0.18,
            size=8, color=MUTED, align='left')
# Y axis label (vertical not supported simply; place rotated text via shape)
yl = slide.shapes.add_textbox(Inches(PX-0.95), Inches(PY+PH/2-0.40),
                              Inches(0.90), Inches(0.80))
yl.text_frame.word_wrap = True
yl.text_frame.margin_left = yl.text_frame.margin_right = Emu(0)
yl.text_frame.margin_top = yl.text_frame.margin_bottom = Emu(0)
ylp = yl.text_frame.paragraphs[0]
ylp.alignment = PP_ALIGN.RIGHT
ylr = ylp.add_run()
ylr.text = "↑ 効率\n(行/コスト)\n高"
set_run_font(ylr, size=8, color=DEEP_FOREST, bold=True)

# Quadrant captions (subtle, in the corners of plot area)
quad_specs = [
    # (text, sub, l_offset_frac, t_offset_frac, color)
    # All quadrant captions in muted gray so they don't compete with the focal Amber.
    ("活用優等生",        "高利用 × 高効率",   0.55, 0.04, MUTED),
    ("効率派・規模小",    "低利用 × 高効率",   0.04, 0.04, MUTED),
    ("要支援 (未活用)",   "低利用 × 低効率",   0.04, 0.78, MUTED),
    ("コスト先行・低効率", "高利用 × 低効率",  0.55, 0.78, MUTED),
]
for txt, sub, lf, tf, col in quad_specs:
    qx = PX + PW * lf
    qy = PY + PH * tf
    add_textbox(slide, txt, l=qx, t=qy, w=PW*0.42, h=0.18,
                size=9, color=col, bold=True)
    add_textbox(slide, sub, l=qx, t=qy+0.16, w=PW*0.42, h=0.16,
                size=7, color=MUTED)

# ---- bubble data (positions preserved from V7; team labels are placeholders
# because V7 has no team names; clearly noted in footnote) ----
# Format: (label, x_norm 0..1, y_norm 0..1 (1=top), size_in, is_own)
# All non-own teams use the same uniform sage color so the focal point (当チーム) pops.
TEAMS = [
    ("当チーム",  0.83, 0.95, 0.20, True),
    ("チームB",   0.86, 0.92, 0.20, False),
    ("チームC",   0.90, 0.78, 0.26, False),
    ("チームD",   0.80, 0.81, 0.21, False),
    ("チームE",   0.76, 0.84, 0.18, False),
    ("チームF",   0.70, 0.73, 0.19, False),
    ("チームG",   0.68, 0.66, 0.16, False),
    ("チームH",   0.58, 0.64, 0.18, False),
    ("チームI",   0.53, 0.60, 0.23, False),
    ("チームJ",   0.48, 0.56, 0.24, False),
    ("チームK",   0.42, 0.55, 0.26, False),
    ("チームL",   0.37, 0.48, 0.28, False),
    ("チームM",   0.32, 0.53, 0.31, False),
    ("チームN",   0.22, 0.39, 0.29, False),
]

# Helper to convert normalized (xn, yn_top1) into inches top-left for bubble
def bubble_pos(xn, yn, dia):
    cx = PX + xn * PW
    cy = PY + (1.0 - yn) * PH   # invert: yn=1 means top of plot
    return cx - dia/2, cy - dia/2

# draw bubbles — own team in Amber (focal), all others in uniform Sage
# Draw "other" bubbles first so the own-team bubble sits on top visually.
for name, xn, yn, dia, is_own in sorted(TEAMS, key=lambda t: 1 if t[4] else 0):
    bx, by = bubble_pos(xn, yn, dia)
    fill = COLOR_OWN if is_own else COLOR_OTHER
    outline = AMBER_DEEP if is_own else FOREST
    add_oval(slide, bx, by, dia, dia, fill=fill, line=outline)
    # label: above-right of bubble
    label_x = bx + dia + 0.02
    label_y = by - 0.04
    # if label would clip the right edge, place to left
    if label_x + 0.55 > PX + PW + 0.10:
        label_x = bx - 0.55
    add_textbox(slide, name,
                l=label_x, t=label_y, w=0.60, h=0.16,
                size=8,
                color=AMBER_DEEP if is_own else MUTED,
                bold=is_own)

# ---- legend (right side of plot) — simplified to 2 items ----
LX = PX + PW + 0.30
LY = PY + 0.10
legend_items = [
    ("● 当チーム",  COLOR_OWN,   AMBER_DEEP, True),
    ("● 他チーム",  COLOR_OTHER, MUTED,      False),
]
for i, (lbl, _swatch_col, text_col, bold) in enumerate(legend_items):
    add_textbox(slide, lbl, l=LX, t=LY + i*0.26, w=1.30, h=0.22,
                size=10, color=text_col, bold=bold)

# Key insight callout
add_rect(slide, LX, LY + 0.75, 1.30, 1.20, fill=SAGE_BG, line=AMBER, line_width=0.75)
add_textbox(slide, "示唆",
            l=LX+0.06, t=LY+0.78, w=1.20, h=0.18,
            size=9, color=AMBER_DEEP, bold=True)
add_textbox(slide,
            "当チームは右上に位置し\n高利用×高効率を実現。\n\n横展開支援で\n左下チームの底上げ\n余地あり。",
            l=LX+0.06, t=LY+0.99, w=1.20, h=0.94,
            size=7, color=INK, line_spacing=1.20)

# footnote
add_textbox(slide,
            "※ チーム名は匿名化表記 (V7も同様にチーム名は非公開)。バブル位置はV7と一致する相対値。具体数値はダッシュボード参照。",
            l=0.40, t=8.32, w=7.50, h=0.16,
            size=7, color=MUTED)

# =====================================================================
# 6) THREE TOP3 PANELS  T: 8.55 - 9.85
# =====================================================================
# Header row
add_textbox(slide,
            "メンバー別 TOP3 — 当チーム内 (V7記載の3名を3視点で再ランキング)",
            l=0.30, t=8.55, w=7.67, h=0.20,
            size=12, color=DEEP_FOREST, bold=True)

# Three columns
panel_w = 2.49
panel_gap = 0.10
panel_left = 0.30

# Data: 3 users from V7 (efficiency top3)
USERS = {
    "A": dict(cost_k=130, loc=4200),
    "B": dict(cost_k=108, loc=4500),
    "C": dict(cost_k=92,  loc=2800),
}
for u, d in USERS.items():
    d["eff"] = round(d["loc"] / d["cost_k"], 1)  # 行 / 千円

# rankings (descending)
top_cost = sorted(USERS.items(), key=lambda kv: -kv[1]["cost_k"])
top_loc  = sorted(USERS.items(), key=lambda kv: -kv[1]["loc"])
top_eff  = sorted(USERS.items(), key=lambda kv: -kv[1]["eff"])

panels = [
    ("コストTOP3",  "最も使った人 (千円)",   top_cost,
        lambda d: f"¥{d['cost_k']}k",
        lambda d: f"{d['loc']:,}行 / {d['eff']}行/¥1k"),
    ("生成行TOP3",  "最もコード生成 (行)",   top_loc,
        lambda d: f"{d['loc']:,}行",
        lambda d: f"¥{d['cost_k']}k / {d['eff']}行/¥1k"),
    ("効率TOP3",    "コスパ (行/¥1k)",      top_eff,
        lambda d: f"{d['eff']}",
        lambda d: f"¥{d['cost_k']}k → {d['loc']:,}行"),
]

panel_top = 8.78
panel_h = 1.07
RANK_BG = [AMBER, SAGE, RGBColor(0xE0, 0xE0, 0xE0)]  # 1st, 2nd, 3rd
for pi, (head, sub, ranking, big_fn, tail_fn) in enumerate(panels):
    px = panel_left + pi * (panel_w + panel_gap)
    add_rect(slide, px, panel_top, panel_w, panel_h,
             fill=WHITE, line=CARD_BORDER, line_width=0.5)
    # header strip
    add_rect(slide, px, panel_top, panel_w, 0.04, fill=DEEP_FOREST)
    add_textbox(slide, head, l=px+0.08, t=panel_top+0.06, w=panel_w-0.16, h=0.18,
                size=11, color=DEEP_FOREST, bold=True)
    add_textbox(slide, sub, l=px+0.08, t=panel_top+0.24, w=panel_w-0.16, h=0.14,
                size=8, color=MUTED)
    # rows
    row_top = panel_top + 0.40
    row_h = 0.22
    for ri, (uname, d) in enumerate(ranking):
        ry = row_top + ri * row_h
        # rank badge
        add_oval(slide, px+0.08, ry+0.02, 0.18, 0.18,
                 fill=RANK_BG[ri], line=None)
        add_textbox(slide, str(ri+1),
                    l=px+0.08, t=ry+0.02, w=0.18, h=0.18,
                    size=10, color=DEEP_FOREST, bold=True,
                    align='center', valign='middle')
        # name
        add_textbox(slide, f"ユーザー{uname}",
                    l=px+0.30, t=ry+0.01, w=0.85, h=0.18,
                    size=10, color=INK, bold=True, valign='middle')
        # big metric (right)
        add_textbox(slide, big_fn(d),
                    l=px+1.00, t=ry+0.01, w=panel_w-1.05, h=0.18,
                    size=12, color=AMBER_DEEP, bold=True,
                    align='right', valign='middle')
        # tail under
        # (omitted to keep panel clean — replaced with single summary line)
    # bottom mini-summary
    add_textbox(slide, tail_fn(ranking[0][1]) if ranking else "",
                l=px+0.08, t=panel_top+panel_h-0.18, w=panel_w-0.16, h=0.16,
                size=7, color=MUTED, align='right')

# =====================================================================
# 7) 当チーム実績 + 公開研究  T: 9.95 - 11.30
# =====================================================================
# left: 当チーム実績
LX1, LY1 = 0.30, 9.95
LW1, LH1 = 3.65, 1.35
add_rect(slide, LX1, LY1, LW1, LH1, fill=DEEP_FOREST, line=None)
add_textbox(slide, "当チーム実績 (3月)",
            l=LX1+0.10, t=LY1+0.06, w=LW1-0.20, h=0.18,
            size=11, color=SAGE, bold=True)
add_textbox(slide, "¥822k → 26,000行",
            l=LX1+0.10, t=LY1+0.24, w=LW1-0.20, h=0.30,
            size=18, color=AMBER, bold=True)
add_line(slide, LX1+0.10, LY1+0.55, LX1+LW1-0.10, LY1+0.55, color=AMBER, width=0.5)
add_textbox(slide, "¥1.5M投資換算 (人月単価相当)",
            l=LX1+0.10, t=LY1+0.58, w=LW1-0.20, h=0.18,
            size=10, color=WHITE, bold=True)
add_textbox(slide, "約 47,400 行 / 人月",
            l=LX1+0.10, t=LY1+0.78, w=LW1-0.20, h=0.30,
            size=20, color=AMBER, bold=True)
add_textbox(slide, "業界中央値 600行/人月の約79倍 (※単純比較)",
            l=LX1+0.10, t=LY1+1.10, w=LW1-0.20, h=0.18,
            size=8, color=SAGE)

# right: 公開研究
RX1, RY1 = 4.10, 9.95
RW1, RH1 = 3.87, 1.35
add_rect(slide, RX1, RY1, RW1, RH1, fill=WHITE, line=CARD_BORDER, line_width=0.5)
add_textbox(slide, "【参考】 公開研究の月間コード生産量",
            l=RX1+0.10, t=RY1+0.06, w=RW1-0.20, h=0.18,
            size=11, color=DEEP_FOREST, bold=True)
add_textbox(slide, "※ 品質・複雑度・業務特性は考慮されない単純比較",
            l=RX1+0.10, t=RY1+0.24, w=RW1-0.20, h=0.14,
            size=8, color=MUTED)

bench = [
    ("保守的",   "325 〜 750 行/月",   "Capers Jones"),
    ("中央値的", "約 1,000 行/月",    "日本SI業界 (IPA)"),
    ("生産性高", "1,500 〜 1,600 行/月","個人実測 / Scrum"),
]
for i, (lvl, num, src) in enumerate(bench):
    by = RY1 + 0.43 + i * 0.30
    add_rect(slide, RX1+0.10, by, RW1-0.20, 0.27, fill=SAGE_BG,
             line=CARD_BORDER, line_width=0.3)
    add_rect(slide, RX1+0.10, by, 0.05, 0.27, fill=FOREST)
    add_textbox(slide, lvl, l=RX1+0.20, t=by+0.04, w=0.70, h=0.18,
                size=9, color=DEEP_FOREST, bold=True, valign='middle')
    add_textbox(slide, num, l=RX1+0.95, t=by+0.04, w=1.55, h=0.18,
                size=10, color=INK, bold=True, valign='middle')
    add_textbox(slide, src, l=RX1+2.50, t=by+0.04, w=RW1-2.65, h=0.18,
                size=8, color=MUTED, valign='middle', align='right')

# =====================================================================
# 8) BOTTOM CONCLUSION BAR  T: 11.30 - 11.55
# =====================================================================
add_rect(slide, 0.00, 11.30, 8.27, 0.39, fill=DEEP_FOREST)
# left amber accent
add_rect(slide, 0.00, 11.30, 0.06, 0.39, fill=AMBER)
add_textbox(slide,
            "効率は改善継続、年度末予測 ¥9.9M=ミドル6.6人月相当。次手は『増員 vs Claude活用拡大』の経営判断",
            l=0.20, t=11.34, w=8.00, h=0.32,
            size=11, color=WHITE, align='left', valign='middle')

# footer (very thin, below)
# (page is full; place footnote inside the conclusion bar by omission)

# =====================================================================
# Speaker notes (sources)
# =====================================================================
notes = slide.notes_slide.notes_text_frame
notes.text = (
    "V8 (revision 2) — based on V7 (claude_usage_report_202603.pptx, 2026-04-29).\n"
    "Changes vs V7:\n"
    "  - REMOVED: 当チーム モデル別利用比率 chart (Opus 52→38%)\n"
    "  - REMOVED: 全社14チーム Opus比率分布 panel\n"
    "  - REMOVED: チーム代表まとめ (リーダーコメント) panel\n"
    "  - SPLIT:   single 効率Top3 → 3 panels (コスト/生成行/効率)\n"
    "  - REDESIGNED: 14チームマッピング — 1メッセージ集中、明確な軸ラベル、\n"
    "                チーム名直接ラベル、4象限の意味文を併記\n"
    "V8 rev2 update (configured per user feedback):\n"
    "  - Team mapping colors collapsed from 4 to 2:\n"
    "      当チーム = Amber #F4B942 (focal),\n"
    "      他13チーム = Sage Green #A8B8A0 (uniform).\n"
    "  - Quadrant captions kept; recolored to muted gray.\n"
    "  - Legend simplified to 2 items (当チーム / 他チーム).\n"
    "Data preserved 1:1 from V7: 月次推移6点, KPI 4枚, ¥9.9M予測, ¥822k/26,000行,\n"
    "Top3 (B 41.7 / A 32.3 / C 30.4 行/¥1k), 公開研究3水準。\n"
    "Team mapping bubble positions match V7 exactly; team names anonymized\n"
    "as チームA-N because V7 also did not show real names."
)

out = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-05-08_Claude_monthly_report_V8.pptx"
prs.save(out)
print("WROTE:", out)
