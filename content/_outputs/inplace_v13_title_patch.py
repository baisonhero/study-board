#!/usr/bin/env python3
"""Targeted title-bar patches per spec.
- 認証基盤プロジェクト (l=0.42, t=0.56) → ANTHRACITE
- 2026年度Q1... (l=0.42, t=0.74) → ANTHRACITE
- V13 badge text (l=6.85, t=0.22) → WHITE
"""
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
import os

SRC = '2026-07-01_Claude_monthly_report_V13_tmp.pptx'
ANTHRACITE = RGBColor(0x19, 0x19, 0x19)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

def main():
    os.chdir(os.path.dirname(os.path.abspath(__file__)))
    prs = Presentation(SRC)
    patched = []
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            l = Emu(sh.left).inches if sh.left else 0
            t = Emu(sh.top).inches if sh.top else 0
            target = None
            if 0.40 <= l <= 0.45 and 0.55 <= t <= 0.60:
                target = ANTHRACITE; label = 'TEAM'
            elif 0.40 <= l <= 0.45 and 0.73 <= t <= 0.76:
                target = ANTHRACITE; label = 'DATE'
            elif 6.83 <= l <= 6.87 and 0.20 <= t <= 0.25:
                target = WHITE; label = 'V13'
            if target is None:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = target
                    patched.append(f'{label}({l:.2f},{t:.2f}) text={r.text!r}')
    prs.save(SRC)
    for p in patched:
        print(p)

if __name__ == '__main__':
    main()
