#!/usr/bin/env python3
"""
V13 rev3 — in-place 4-point micro-adjustments:

  1. Header underline/title overlap fixed
     - Title up to T=0.15, amber bar matches
     - Team name T=0.56, subtitle T=0.74
     - Divider line moved to T=0.94 (below subtitle, above judge T=1.00)
     - V13 pill aligned to new title Y

  2. Team-map Y axis title moved to top, 2-line, [ ] brackets
     "コスト\n(千円/人)" → "コスト\n[千円/人]" at T≈8.42 above the "120" tick

  3. Team-map X axis title to right of "2000" tick, [ ] brackets
     "行数 (行/人)" → "行数 [行/人]"

  4. Axis ticks reduced to 3 each:
     X: keep 0/1000/2000, delete 500/1500
     Y: keep 0/60/120,   delete 30/90
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V13.pptx"

prs = Presentation(PPTX)
slide = prs.slides[0]

# ============================================================
# 1) Header layout — fix overlap of underline/divider with text
# ============================================================
# New target positions:
#   Title (shape[1])            T=0.15 H=0.40  (was 0.20)
#   Amber bar (shape[0])        T=0.20 H=0.50  (taller, matching title)
#   V13 pill rect (shape[2])    T=0.19 H=0.32  (aligned)
#   V13 pill text (shape[3])    T=0.22 H=0.26
#   Team name (shape[199])      T=0.56 H=0.18
#   Subtitle (shape[177])       T=0.74 H=0.18
#   Divider line (shape[4])     T=0.94 (below subtitle bottom 0.92)
#   Judge area starts at        T=1.00 (unchanged)

for sh in slide.shapes:
    try:
        t_in = sh.top/914400; l_in = sh.left/914400
    except: continue

    # Title text
    if sh.has_text_frame and "Claude Code 利用量レポート" in sh.text_frame.text:
        sh.top = Inches(0.15)
        print("  ✓ Title moved to T=0.15")
    # Amber accent bar (small narrow rect on the left)
    elif sh.shape_type == 1 and 0.28 <= t_in < 0.34 and l_in < 0.40 and (sh.width/914400) < 0.10:
        sh.top = Inches(0.20)
        sh.height = Inches(0.50)
        print("  ✓ Amber accent bar → T=0.20 H=0.50")
    # V13 pill rect
    elif sh.shape_type == 1 and 0.20 <= t_in < 0.30 and 6.80 < l_in < 7.00:
        sh.top = Inches(0.19)
        print("  ✓ V13 pill rect → T=0.19")
    # V13 pill text
    elif sh.has_text_frame and sh.text_frame.text.strip() == "V13":
        sh.top = Inches(0.22)
        print("  ✓ V13 pill text → T=0.22")
    # Team name
    elif sh.has_text_frame and sh.text_frame.text.strip() == "認証基盤プロジェクト":
        sh.top = Inches(0.56)
        print("  ✓ Team name → T=0.56")
    # Subtitle
    elif sh.has_text_frame and "2026年度Q1" in sh.text_frame.text and t_in < 1.0:
        sh.top = Inches(0.74)
        print("  ✓ Subtitle → T=0.74")
    # Divider line
    elif sh.shape_type == 9 and 0.80 <= t_in < 0.95 and l_in < 0.40 and (sh.width/914400) > 5.0:
        sh.top = Inches(0.94)
        print("  ✓ Divider line → T=0.94 (between subtitle and judge)")

# ============================================================
# 2) Y axis title: move to top of plot, center-aligned 2-line
# ============================================================
# Plot left edge PX=1.20. "120" tick at L=0.90 T=8.59.
# Place axis title centered around plot left edge (≈0.95) above the 120 tick.
for sh in slide.shapes:
    if not sh.has_text_frame: continue
    txt = sh.text_frame.text
    if "コスト" in txt and ("千円" in txt or "[" in txt):
        try:
            t_in = sh.top/914400
        except: continue
        if t_in > 8.0:
            # Move to top of plot
            sh.left = Inches(0.55)
            sh.top = Inches(8.40)
            sh.width = Inches(0.80)
            sh.height = Inches(0.30)
            # Update text and styling
            tf = sh.text_frame
            tf.word_wrap = True
            for para in list(tf.paragraphs)[1:]:
                para._p.getparent().remove(para._p)
            p0 = tf.paragraphs[0]
            # Update first run text → "コスト\n[千円/人]"
            if p0.runs:
                p0.runs[0].text = "コスト\n[千円/人]"
                # delete extra runs in p0
                for r in list(p0.runs)[1:]:
                    r._r.getparent().remove(r._r)
            # Center-align
            p0.alignment = PP_ALIGN.CENTER
            print(f"  ✓ Y axis title → top, centered 2-line 'コスト\\n[千円/人]'")
            break

# ============================================================
# 3) X axis title: to right of "2000" tick, square brackets
# ============================================================
for sh in slide.shapes:
    if not sh.has_text_frame: continue
    txt = sh.text_frame.text.strip()
    if txt == "行数 (行/人)" or "行数" in txt and "行/人" in txt:
        try:
            t_in = sh.top/914400
        except: continue
        if t_in > 9.0:
            # Move to right of "2000" tick (which ends at ~L=5.57)
            sh.left = Inches(5.70)
            sh.top = Inches(9.70)
            sh.width = Inches(1.10)
            sh.height = Inches(0.14)
            tf = sh.text_frame
            p0 = tf.paragraphs[0]
            if p0.runs:
                p0.runs[0].text = "行数 [行/人]"
                for r in list(p0.runs)[1:]:
                    r._r.getparent().remove(r._r)
            p0.alignment = PP_ALIGN.LEFT
            print(f"  ✓ X axis title → 'L=5.70 right of 2000-tick / 行数 [行/人]'")
            break

# ============================================================
# 4) Reduce axis ticks to 3 each
# ============================================================
# Delete: X-ticks "500", "1500"; Y-ticks "30", "90"
DELETE_TICKS = {"500", "1500", "30", "90"}
to_del = []
for sh in slide.shapes:
    if not sh.has_text_frame: continue
    txt = sh.text_frame.text.strip()
    if txt not in DELETE_TICKS:
        continue
    try:
        t_in = sh.top/914400
        # X-axis ticks: T≈9.71, Y-axis ticks: T≈8.5-9.7 with L≈0.90
    except: continue
    if 8.50 <= t_in < 9.80:
        # double-check by L
        l_in = sh.left/914400
        if abs(l_in - 0.90) < 0.05 or 0.80 <= l_in < 5.50:
            to_del.append((sh, txt))

for sh, txt in to_del:
    sh._element.getparent().remove(sh._element)
    print(f"  ✓ Deleted tick '{txt}'")

prs.save(PPTX)
print(f"\nWROTE: {PPTX}")
