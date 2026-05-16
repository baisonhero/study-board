#!/usr/bin/env python3
"""
V13 -> V13_tmp color-only swap to Claude brand palette.
Does NOT touch data, text, layout, or fonts. Only color values.

Mapping policy
==============
1F3A2E DEEP_FOREST  → geometry-dependent:
  - h<=0.05            → CORAL          (thin underline)
  - l<=0.05 & w>=7.5   → CORAL          (full-width section band)
  - h<=0.22 & w>=5     → SAND           (thin wide band, e.g. table header)
  - else               → CORAL_DARK     (cards / badges)
2E6F40 FOREST       → fill→CORAL, text→WARM_OLIVE  (delta indicator)
222824 INK          → ANTHRACITE
C88E1F AMBER_DEEP   → fill→CORAL_DARK, text→BRICK_RED
6B6B6B MUTED        → WARM_GRAY
D7DCCF CARD_BORDER  → SAND
A8B8A0 SAGE_OTHER   → WARM_GRAY        (other-team bubbles)
B5C9A8 SAGE_BRIGHT  → SAND             (#2 rank badges)
F4B942 AMBER        → CORAL            (own team / accents)
E8EFE3 SAGE_BG      → CREAM
C0392B RED          → BRICK_RED
"""
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from lxml import etree
import shutil, sys, os

SRC = '2026-07-01_Claude_monthly_report_V13_tmp.pptx'

# Claude brand palette
CORAL       = 'CC785C'
CORAL_DARK  = 'A85F4A'
CREAM       = 'F5F4EE'
SAND        = 'E8DCC8'
WARM_OLIVE  = '7A8450'
BRICK_RED   = 'A8503D'
ANTHRACITE  = '191919'
WARM_GRAY   = '8B8680'

# Simple direct map (uppercase hex)
SIMPLE = {
    '222824': ANTHRACITE,
    '6B6B6B': WARM_GRAY,
    'D7DCCF': SAND,
    'A8B8A0': WARM_GRAY,
    'B5C9A8': SAND,
    'F4B942': CORAL,
    'E8EFE3': CREAM,
    'C0392B': BRICK_RED,
}

# Colors needing context (decided in walker)
CONTEXT_FILL = {
    '2E6F40': CORAL,        # fill use of FOREST → CORAL (bars/accents)
    'C88E1F': CORAL_DARK,   # fill use of AMBER_DEEP → CORAL_DARK
}
CONTEXT_TEXT = {
    '2E6F40': WARM_OLIVE,   # text use of FOREST → WARM_OLIVE (positive delta)
    'C88E1F': BRICK_RED,    # text use of AMBER_DEEP → BRICK_RED (>100%)
}

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def deep_forest_for_shape(l_in, t_in, w_in, h_in):
    if h_in <= 0.05:
        return CORAL
    if l_in <= 0.05 and w_in >= 7.5:
        return CORAL
    if h_in <= 0.22 and w_in >= 5:
        return SAND
    return CORAL_DARK

def replace_srgb(el_color_parent, old_to_new_hex_func):
    """Walk down el_color_parent looking for <a:srgbClr val='...'/>.
    old_to_new_hex_func: callable old_upper_hex -> new_upper_hex or None to skip.
    Returns number of replacements made."""
    n = 0
    for srgb in el_color_parent.iter(f'{{{A_NS}}}srgbClr'):
        val = srgb.get('val', '')
        if not val:
            continue
        new = old_to_new_hex_func(val.upper())
        if new and new.upper() != val.upper():
            srgb.set('val', new.upper())
            n += 1
    return n

def walk_shape(shape, summary):
    # Determine geometry
    l = Emu(shape.left).inches if shape.left is not None else 0
    t = Emu(shape.top).inches if shape.top is not None else 0
    w = Emu(shape.width).inches if shape.width is not None else 0
    h = Emu(shape.height).inches if shape.height is not None else 0

    el = shape._element  # CT_Shape / CT_Pic / etc.

    # 1) SHAPE-LEVEL fill / line: find <p:spPr> and replace shape-tier colors
    # spPr is the shape properties container — fills and line live there
    sp_pr_list = el.findall(f'.//{{http://schemas.openxmlformats.org/drawingml/2006/main}}solidFill')
    # We need to distinguish fill of the shape vs fill of text run vs fill of line.
    # Use direct parent context.

    # Walk every srgbClr inside the shape
    for srgb in el.iter(f'{{{A_NS}}}srgbClr'):
        val = srgb.get('val', '').upper()
        if not val:
            continue
        # Detect context by climbing parents
        parent = srgb.getparent()
        ancestors = []
        node = parent
        while node is not None:
            tag = etree.QName(node).localname
            ancestors.append(tag)
            node = node.getparent()
        # is this a text run color?  (parent chain contains 'rPr' or 'defRPr')
        is_text = any(a in ('rPr', 'defRPr', 'endParaRPr') for a in ancestors)
        # is this a line color?
        is_line = any(a in ('ln',) for a in ancestors)

        new = None
        if val in SIMPLE:
            new = SIMPLE[val]
        elif val == '1F3A2E':
            if is_text:
                new = CORAL          # any DEEP_FOREST text → CORAL
            elif is_line:
                new = CORAL          # rare; underline lines
            else:
                new = deep_forest_for_shape(l, t, w, h)
        elif val in CONTEXT_FILL:
            if is_text:
                new = CONTEXT_TEXT[val]
            else:
                new = CONTEXT_FILL[val]
        if new and new.upper() != val:
            srgb.set('val', new.upper())
            summary[(val, new.upper())] = summary.get((val, new.upper()), 0) + 1

def main():
    os.chdir(os.path.dirname(os.path.abspath(__file__)))
    if not os.path.exists(SRC):
        print(f'missing: {SRC}', file=sys.stderr)
        sys.exit(1)
    prs = Presentation(SRC)
    summary = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            walk_shape(shape, summary)
    prs.save(SRC)
    print('replacement counts (old -> new : count):')
    for (o, n), c in sorted(summary.items()):
        print(f'  {o} -> {n} : {c}')

if __name__ == '__main__':
    main()
