#!/usr/bin/env python3
"""
V11 rev15 — in-place 3 changes:

  1. Delete subtitle 「2026年度Q1 累積  /  報告日: ... /  対象期間: ...」 (header)
  2. Delete chart subtitle 「※ 7月以降は計画値 (月 750 [千円])  /  単月行数は実績のみ表示」
  3. Replace right reference panel with content from STALE_LOCKED 版:
       - Title with IPA URL inline
       - Updated caveat
       - 2-row layout: 改良開発上位25% / 約800 行/月
                     新規開発上位25% / 約1,000 行/月
       (Left panel "当チーム実績 (6月)" is identical between current V11 and
        STALE — no change needed.)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

PPTX = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-07-01_Claude_monthly_report_V11.pptx"

# Palette
DEEP_FOREST = RGBColor(0x1F, 0x3A, 0x2E)
FOREST      = RGBColor(0x2E, 0x6F, 0x40)
SAGE_BG     = RGBColor(0xE8, 0xEF, 0xE3)
INK         = RGBColor(0x22, 0x28, 0x24)
MUTED       = RGBColor(0x6B, 0x6B, 0x6B)
CARD_BORDER = RGBColor(0xD7, 0xDC, 0xCF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
JP_FONT, EN_FONT = "メイリオ", "Segoe UI"

def set_run_font(run, size, color, bold=False, ea=JP_FONT, latin=EN_FONT):
    run.font.name = latin
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea_elem = etree.SubElement(rPr, qn('a:ea'))
    ea_elem.set('typeface', ea)

def add_textbox(slide, text, l, t, w, h, size=12, color=INK, bold=False,
                align='left', valign='top', word_wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if valign == 'middle': tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == 'bottom': tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    if align == 'center': p.alignment = PP_ALIGN.CENTER
    elif align == 'right': p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    set_run_font(r, size=size, color=color, bold=bold)
    return tb

def add_rect(slide, l, t, w, h, fill=WHITE, line=None, line_width=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        if line_width: shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

prs = Presentation(PPTX)
slide = prs.slides[0]

n_deleted = 0
shapes_to_delete = []

for sh in slide.shapes:
    if not sh.has_text_frame:
        # Item 3 — right reference band area: also delete autoshapes (track bgs etc)
        try:
            t_in = sh.top/914400; l_in = sh.left/914400
            if t_in >= 10.30 and l_in >= 4.05:
                shapes_to_delete.append(sh)
        except Exception:
            pass
        continue

    text = sh.text_frame.text
    try:
        t_in = sh.top/914400; l_in = sh.left/914400
    except Exception:
        continue

    # Item 1: subtitle "2026年度Q1 累積..."
    if "2026年度Q1" in text and t_in < 1.0:
        shapes_to_delete.append(sh)
        continue
    # Item 2: chart subtitle "※ 7月以降は計画値..."
    if "7月以降" in text and 2.30 <= t_in < 2.80:
        shapes_to_delete.append(sh)
        continue
    # Item 3: right reference panel text shapes (L >= 4.05, Y >= 10.30)
    if t_in >= 10.30 and l_in >= 4.05:
        shapes_to_delete.append(sh)
        continue

# Delete collected shapes
for sh in shapes_to_delete:
    sp = sh._element
    sp.getparent().remove(sp)
    n_deleted += 1
print(f"Deleted {n_deleted} shapes (item1 subtitle + item2 chart cap + item3 right panel).")

# ---------- Item 3: build new right panel from STALE 版 ----------
# Geometry from STALE inspection: L=4.21 (slight gap), W=3.76, H=1.20+
RX, RY = 4.21, 10.40
RW, RH = 3.76, 1.20

# Outer card (white bg, light border)
add_rect(slide, RX, RY, RW, RH, fill=WHITE, line=CARD_BORDER, line_width=0.5)

# Title: 「【参考】 IPA分析」 + URL (mixed run sizes)
title_tb = slide.shapes.add_textbox(Inches(RX+0.10), Inches(RY+0.05),
                                     Inches(RW-0.20), Inches(0.30))
title_tf = title_tb.text_frame
title_tf.word_wrap = True
title_tf.margin_left = title_tf.margin_right = Emu(0)
title_tf.margin_top = title_tf.margin_bottom = Emu(0)
p1 = title_tf.paragraphs[0]
p1.alignment = PP_ALIGN.LEFT
r1 = p1.add_run()
r1.text = "【参考】 IPA分析  "
set_run_font(r1, size=10, color=DEEP_FOREST, bold=True)
r2 = p1.add_run()
r2.text = "https://www.ipa.go.jp/digital/software-survey/metrics/hjuojm000000c6it-att/000102171.pdf"
set_run_font(r2, size=5.5, color=DEEP_FOREST, bold=False)

# Caveat
add_textbox(slide,
            "※ 品質・複雑度・業務特性・言語差分は考慮されない単純比較",
            l=RX+0.10, t=RY+0.40, w=RW-0.20, h=0.13,
            size=7.5, color=MUTED)

# Two rows: 改良開発上位25% (約800 行/月) and 新規開発上位25% (約1,000 行/月)
# Bar bg + small marker, label left, value right.
ROWS = [
    ("改良開発上位25%", 800,    "約 800 行/月"),
    ("新規開発上位25%", 1000,   "約 1,000 行/月"),
]
ROW_TOPS = [RY + 0.60, RY + 0.92]
ROW_H    = 0.24
TRACK_L  = RX + 0.10
TRACK_W  = RW - 0.20
# Marker proportional: max value is 1,000 in this set → bar fills relative
MAX_VAL = max(v for _, v, _ in ROWS)
for (label, val, value_text), top in zip(ROWS, ROW_TOPS):
    # Track background (light sage)
    add_rect(slide, TRACK_L, top, TRACK_W, ROW_H, fill=SAGE_BG, line=None)
    # Forest accent (left edge, 0.05 wide as in STALE)
    add_rect(slide, TRACK_L, top, 0.05, ROW_H, fill=FOREST)
    # Optional: filled portion proportional to value
    bar_fill_w = TRACK_W * (val / MAX_VAL)
    add_rect(slide, TRACK_L, top + ROW_H - 0.05, bar_fill_w, 0.04, fill=FOREST)
    # Label (left side)
    add_textbox(slide, label,
                l=TRACK_L+0.10, t=top+0.04, w=1.55, h=0.16,
                size=9, color=DEEP_FOREST, bold=True, valign='middle',
                word_wrap=False)
    # Value (right side)
    add_textbox(slide, value_text,
                l=TRACK_L + TRACK_W - 1.40, t=top+0.04, w=1.30, h=0.16,
                size=10, color=INK, bold=True, align='right', valign='middle',
                word_wrap=False)

prs.save(PPTX)
print(f"WROTE: {PPTX}")
