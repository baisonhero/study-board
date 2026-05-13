#!/usr/bin/env python3
"""
Build V10 of Claude Code 月次利用量レポート (A4 portrait, single-page dashboard).

V10 changes vs V9 (per user feedback):
  1) 14チームポジショニング図: チーム名ラベル全削除（凡例の「他チーム」も削除）。
     残すもの: 軸、自チーム=Amberのバブル＋「当チーム」マーカー、4象限メッセージ、
              示唆コールアウト、フットノート（簡潔化）。
  2) 「ミドル」全置換: 文脈に応じて「外注費用」「外注増員」など自然な日本語に。
       - 来期の投資規模: ミドル6.6人月相当     → 外注費用 6.6人月相当
       - ≈ ミドル人材 6.6人月分                → ≈ 外注費用 6.6人月分
       - ミドル増員 vs Claude活用拡大          → 外注増員 vs Claude活用拡大
       - 年度末予測 9,900 [千円] = ミドル6.6人月 → 年度末予測 9,900 [千円] = 外注費用 6.6人月
  3) バージョン/日付: V9 → V10, 2026-05-09 → 2026-05-11.

V9 changes vs V8 rev3 (preserved in V10):
  - 単位 [千円] 統一 (k/M 廃止).
  - 月次推移 band を LEFT 58% (二軸グラフ) / RIGHT 42% (プロジェクトメンバー
    コストTOP5) に分割.
  - メンバー別TOP3 を 3→2 panel (コストTOP3 削除, 生成行TOP3 + 効率TOP3 のみ).

Mock-data note: TOP1-3 cost preserved 1:1 from V7 (130/108/92 [千円]); V7
codenames ユーザーA/B/C are remapped to user-supplied real names
山田太郎 / 鈴木花子 / 佐藤健. TOP4 (田中美咲 78) and TOP5 (高橋翔 64) are mock
values provided by the user *for layout-confirmation only*. A footnote on-slide
explicitly says "※メンバー名・金額はレイアウト確認用のモックデータ".
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# matplotlib for the dual-axis monthly trend chart (rendered to PNG, embedded as picture)
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import MaxNLocator
from matplotlib import font_manager as fm

# Per-glyph font fallback (matplotlib >=3.6): set font.family to a LIST directly.
# DejaVu Sans handles Latin/digits/punctuation; Droid Sans Fallback covers CJK.
# Note: setting font.sans-serif (alone) does NOT enable per-glyph fallback —
# we have to pass the list to font.family itself.
plt.rcParams["font.family"] = [
    "DejaVu Sans",          # Latin & digits & punctuation
    "Droid Sans Fallback",  # CJK fallback (Linux sandbox)
    "Noto Sans CJK JP",     # preferred CJK if available
    "Hiragino Sans", "Yu Gothic", "Meiryo",
]
plt.rcParams["axes.unicode_minus"] = False

# ---------- House-style: Calm Forest + Amber ----------
DEEP_FOREST  = RGBColor(0x1F, 0x3A, 0x2E)  # primary
FOREST       = RGBColor(0x2E, 0x6F, 0x40)  # secondary
SAGE         = RGBColor(0xB5, 0xC9, 0xA8)  # tertiary (light fill)
SAGE_BG      = RGBColor(0xE8, 0xEF, 0xE3)  # very light card bg
AMBER        = RGBColor(0xF4, 0xB9, 0x42)  # accent
AMBER_DEEP   = RGBColor(0xC8, 0x8E, 0x1F)  # accent text
INK          = RGBColor(0x22, 0x28, 0x24)  # body
MUTED        = RGBColor(0x6B, 0x6B, 0x6B)  # caption
CARD_BORDER  = RGBColor(0xD7, 0xDC, 0xCF)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
RED          = RGBColor(0xC0, 0x39, 0x2B)  # for ↑ negative ratios
GREEN_OK     = RGBColor(0x2E, 0x6F, 0x40)  # for ↑ positive ratios

# bubble colors (2 colors only, per user revision)
COLOR_OWN    = AMBER                     # 当チーム — focal point
COLOR_OTHER  = RGBColor(0xA8, 0xB8, 0xA0)  # Sage Green — all 13 other teams (uniform)

JP_FONT  = "メイリオ"
EN_FONT  = "Segoe UI"

# ---------- helpers ----------
def set_run_font(run, latin=EN_FONT, ea=JP_FONT, size=12, color=INK, bold=False):
    run.font.name = latin
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn('a:ea')):
        rPr.remove(old)
    ea_elem = etree.SubElement(rPr, qn('a:ea'))
    ea_elem.set('typeface', ea)

def add_textbox(slide, text, l, t, w, h, size=12, color=INK, bold=False,
                align='left', valign='top', latin=EN_FONT, ea=JP_FONT,
                line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if valign == 'top':
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif valign == 'middle':
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == 'bottom':
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    if align == 'left':
        p.alignment = PP_ALIGN.LEFT
    elif align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    set_run_font(r, latin=latin, ea=ea, size=size, color=color, bold=bold)
    return tb

def add_rect(slide, l, t, w, h, fill=WHITE, line=None, line_width=None,
             corner=None):
    if corner:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(l), Inches(t), Inches(w), Inches(h))
        # set corner radius via XML adjust
        try:
            shp.adjustments[0] = corner
        except Exception:
            pass
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        if line_width:
            shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_oval(slide, l, t, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_line(slide, x1, y1, x2, y2, color=CARD_BORDER, width=0.75):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line

# ---------- slide setup (A4 portrait) ----------
prs = Presentation()
prs.slide_width  = Inches(8.27)
prs.slide_height = Inches(11.69)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# white background (default), draw nothing

# =====================================================================
# 1) HEADER (top)  T: 0.25 - 0.95
# =====================================================================
# left amber accent bar
add_rect(slide, 0.30, 0.30, 0.06, 0.50, fill=AMBER)

add_textbox(slide, "Claude Code 月次利用量レポート",
            l=0.42, t=0.20, w=6.30, h=0.45,
            size=22, color=DEEP_FOREST, bold=True)

# right-aligned version pill
add_rect(slide, 6.95, 0.24, 1.02, 0.34, fill=DEEP_FOREST, corner=0.30)
add_textbox(slide, "V10",
            l=6.95, t=0.27, w=1.02, h=0.28,
            size=14, color=AMBER, bold=True, align='center', valign='middle')

add_textbox(slide, "2026年3月度  /  報告者: Ogata  /  報告日: 2026-05-11  /  対象期間: 2026/03/01 - 03/31",
            l=0.42, t=0.65, w=7.50, h=0.20,
            size=9, color=MUTED)

# divider
add_line(slide, 0.30, 0.92, 7.97, 0.92, color=CARD_BORDER, width=1.0)

# =====================================================================
# 2) JUDGE box (left) + 年度末予測 (right)   T: 1.00 - 2.85
# =====================================================================
add_textbox(slide, "今月のジャッジ",
            l=0.30, t=1.00, w=4.70, h=0.22,
            size=14, color=DEEP_FOREST, bold=True)
add_textbox(slide, "数字から見えること — チームでの議論材料",
            l=0.30, t=1.20, w=4.70, h=0.18,
            size=9, color=MUTED)

judge_items = [
    ("✓", GREEN_OK,
     "効率性は改善傾向 (継続観測中)",
     "過去6ヶ月でコスト2.3倍に対し採用行数3.7倍"),
    ("▶", AMBER_DEEP,
     "来期の投資規模: 外注費用 6.6人月相当",
     "外注増員 vs Claude活用拡大、議論事項"),
    ("⚠", RED,
     "即対応: 認証基盤PJのモデル選定見直し",
     "局所的な過剰利用の兆候 (詳細は別ダッシュボード)"),
]
for i, (icon, icon_color, head, sub) in enumerate(judge_items):
    y = 1.42 + i * 0.47
    add_rect(slide, 0.30, y, 4.70, 0.43, fill=SAGE_BG, line=CARD_BORDER, line_width=0.5)
    add_rect(slide, 0.30, y, 0.05, 0.43, fill=icon_color)
    add_textbox(slide, icon, l=0.43, t=y+0.06, w=0.32, h=0.32,
                size=18, color=icon_color, bold=True, align='center', valign='middle')
    add_textbox(slide, head, l=0.82, t=y+0.03, w=4.10, h=0.21,
                size=12, color=DEEP_FOREST, bold=True)
    add_textbox(slide, sub, l=0.82, t=y+0.22, w=4.10, h=0.19,
                size=9, color=INK)

# right: 年度末着地予測
add_rect(slide, 5.10, 1.42, 2.87, 1.45, fill=DEEP_FOREST, line=None)
add_textbox(slide, "年度末着地予測",
            l=5.20, t=1.48, w=2.67, h=0.20,
            size=11, color=SAGE, bold=True)
add_textbox(slide, "9,900 [千円]",
            l=5.05, t=1.66, w=2.97, h=0.55,
            size=28, color=AMBER, bold=True, align='center')
add_line(slide, 5.40, 2.27, 7.67, 2.27, color=AMBER, width=0.5)
add_textbox(slide, "≈ 外注費用 6.6人月分",
            l=5.10, t=2.30, w=2.87, h=0.22,
            size=12, color=WHITE, bold=True, align='center')
add_textbox(slide, "1,500 [千円]/月換算 / 現行ペース×12ヶ月",
            l=5.10, t=2.51, w=2.87, h=0.18,
            size=9, color=SAGE, align='center')
add_textbox(slide, "投資判断: 人を雇うか / Claude活用拡大か",
            l=5.10, t=2.69, w=2.87, h=0.18,
            size=9, color=AMBER, align='center', bold=True)

# =====================================================================
# 3) KPI TILES   T: 3.00 - 3.78
# =====================================================================
kpi = [
    ("総コスト [千円]",   "822",     "+19.0%",  "vs 前月 691 [千円]", AMBER_DEEP),
    ("アクティブユーザー", "12名",    "+2名",    "vs 前月 10名",       GREEN_OK),
    ("1人あたり月額 [千円]", "68.5",  "−1.2%",   "vs 前月 69.1 [千円]", GREEN_OK),
    ("1人あたり採用行数",  "2,180行", "+18%",    "全社平均 1,520行",   GREEN_OK),
]
tile_w = 1.80
tile_gap = 0.16
start_x = 0.30
for i, (label, big, delta, sub, dc) in enumerate(kpi):
    x = start_x + i * (tile_w + tile_gap)
    add_rect(slide, x, 3.00, tile_w, 0.78, fill=WHITE, line=CARD_BORDER, line_width=0.5)
    add_rect(slide, x, 3.00, tile_w, 0.04, fill=DEEP_FOREST)
    add_textbox(slide, label, l=x+0.08, t=3.05, w=tile_w-0.16, h=0.17,
                size=9, color=MUTED)
    add_textbox(slide, big, l=x+0.08, t=3.20, w=tile_w-0.16, h=0.30,
                size=19, color=DEEP_FOREST, bold=True)
    add_textbox(slide, delta, l=x+0.08, t=3.50, w=0.99, h=0.16,
                size=9, color=dc, bold=True)
    add_textbox(slide, sub, l=x+0.08, t=3.62, w=tile_w-0.16, h=0.14,
                size=8, color=MUTED)

# =====================================================================
# 4) MONTHLY TREND CHART (LEFT 58%)  +  PROJECT-MEMBER TOP5 PANEL (RIGHT 42%)
#    Band: T = 3.92 - 5.50
# =====================================================================
# --- LEFT 58%: trend title + chart ---
TREND_L, TREND_W = 0.30, 4.40
TOP5_L,  TOP5_W  = 4.80, 3.17     # gap of 0.10 between panels

add_textbox(slide, "月次推移 — コスト × 採用コード行数 (過去6ヶ月)",
            l=TREND_L, t=3.92, w=TREND_W, h=0.22,
            size=10, color=DEEP_FOREST, bold=True)

# --- Render dual-axis chart with matplotlib (bar=cost on left, line=LOC on right) ---
cats = ['25/10', '25/11', '25/12', '26/01', '26/02', '26/03']
cost_k  = [357, 420, 510, 612, 691, 822]      # 千円
loc_k   = [7.0, 10.5, 12.0, 17.5, 22.0, 26.0]  # 千行

forest_hex = "#1F3A2E"
amber_hex  = "#F4B942"
muted_hex  = "#6B6B6B"
border_hex = "#D7DCCF"

# Match the slide's chart placeholder area: 4.40 x 1.33 inches (LEFT 58% panel)
chart_w_in, chart_h_in = 4.40, 1.33
fig, ax_left = plt.subplots(figsize=(chart_w_in, chart_h_in), dpi=220)
fig.patch.set_facecolor("white")

# Left Y axis: cost as bars
bars = ax_left.bar(cats, cost_k, color=forest_hex, width=0.55,
                   label="コスト [千円]", zorder=2)
ax_left.set_ylabel("コスト [千円]", color=forest_hex, fontsize=7.5, labelpad=2)
ax_left.tick_params(axis="y", labelcolor=forest_hex, labelsize=7,
                    length=2, pad=1)
ax_left.tick_params(axis="x", labelsize=7, colors=muted_hex,
                    length=2, pad=1)
ax_left.grid(True, axis="y", linestyle="--", linewidth=0.4,
             color=muted_hex, alpha=0.35, zorder=1)
ax_left.set_axisbelow(True)
for spine in ("top", "right"):
    ax_left.spines[spine].set_visible(False)
ax_left.spines["left"].set_color(forest_hex)
ax_left.spines["left"].set_linewidth(0.6)
ax_left.spines["bottom"].set_color(muted_hex)
ax_left.spines["bottom"].set_linewidth(0.6)
ax_left.yaxis.set_major_locator(MaxNLocator(nbins=5, integer=True))
ax_left.set_ylim(0, max(cost_k) * 1.18)

# Right Y axis: LOC as line + markers (Amber)
ax_right = ax_left.twinx()
line = ax_right.plot(cats, loc_k, color=amber_hex, marker="o",
                     markersize=5, markeredgecolor=amber_hex,
                     markerfacecolor=amber_hex, linewidth=1.8,
                     label="生成行数 [千行]", zorder=3)
ax_right.set_ylabel("生成行数 [千行]", color=amber_hex,
                    fontsize=7.5, labelpad=2)
ax_right.tick_params(axis="y", labelcolor=amber_hex, labelsize=7,
                     length=2, pad=1)
for spine in ("top", "left"):
    ax_right.spines[spine].set_visible(False)
ax_right.spines["right"].set_color(amber_hex)
ax_right.spines["right"].set_linewidth(0.6)
ax_right.set_ylim(0, max(loc_k) * 1.25)

# Combined legend in upper-left
handles = [bars, line[0]]
labels  = ["コスト [千円]", "生成行数 [千行]"]
ax_left.legend(handles, labels, loc="upper left",
               fontsize=7, frameon=False, handlelength=1.6,
               handletextpad=0.4, borderpad=0.0,
               labelcolor=[forest_hex, amber_hex])

# Tight margins so the rendered PNG fills the picture box
plt.subplots_adjust(left=0.07, right=0.93, top=0.93, bottom=0.16)

mpl_chart_png = "/tmp/v10_monthly_trend.png"
fig.savefig(mpl_chart_png, dpi=220, facecolor="white", edgecolor="none")
plt.close(fig)

# Embed the rendered PNG as a picture in the LEFT trend slot (a tad shorter
# so the in-band subtitle has room beneath it).
slide.shapes.add_picture(mpl_chart_png,
                         Inches(TREND_L), Inches(4.17),
                         width=Inches(TREND_W), height=Inches(1.18))

# Subtitle / "効率改善継続" note moved underneath the chart
add_textbox(slide, "コスト 2.3倍 vs 採用行数 3.7倍 = 効率改善継続",
            l=TREND_L, t=5.36, w=TREND_W, h=0.14,
            size=8, color=AMBER_DEEP, bold=True, align='right')

# =====================================================================
# 4b) RIGHT 42%: 自プロジェクトメンバー コストTOP5 panel  (NEW in V9)
# =====================================================================
add_textbox(slide, "プロジェクトメンバー コストTOP5",
            l=TOP5_L, t=3.92, w=TOP5_W-0.55, h=0.22,
            size=10, color=DEEP_FOREST, bold=True)
add_textbox(slide, "[千円]",
            l=TOP5_L + TOP5_W - 0.55, t=3.94, w=0.50, h=0.20,
            size=9, color=MUTED, bold=False, align='right', valign='middle')

# Panel card — matches the shortened trend chart height
add_rect(slide, TOP5_L, 4.17, TOP5_W, 1.18,
         fill=WHITE, line=CARD_BORDER, line_width=0.5)
add_rect(slide, TOP5_L, 4.17, TOP5_W, 0.04, fill=DEEP_FOREST)

# Mock data (per user, layout-confirmation only)
TOP5_MEMBERS = [
    # (rank, name, cost_kjpy)  — sorted by cost desc
    (1, "山田 太郎", 130),
    (2, "鈴木 花子", 108),
    (3, "佐藤 健",    92),
    (4, "田中 美咲",  78),
    (5, "高橋 翔",    64),
]
TOP5_MAX_COST = max(c for _, _, c in TOP5_MEMBERS)  # 130

# Row geometry inside the card (5 rows fit in ~ T=4.28..5.23)
row_top    = 4.28
row_h      = 0.18
inner_pad  = 0.08
rank_w     = 0.16
name_w     = 0.88
gap        = 0.04
# Bar area: from after name to before value text. Value cell wide enough to keep
# 3-digit numbers on a single line (unit shown only in the panel header).
value_w    = 0.45
bar_x      = TOP5_L + inner_pad + rank_w + gap + name_w + gap
bar_w_max  = TOP5_W - inner_pad - rank_w - gap - name_w - gap - value_w - gap - inner_pad

RANK_BG = [AMBER, SAGE, RGBColor(0xE0, 0xE0, 0xE0), SAGE_BG, SAGE_BG]

for i, (rank, name, cost) in enumerate(TOP5_MEMBERS):
    ry = row_top + i * row_h
    rx = TOP5_L + inner_pad
    # rank badge
    add_oval(slide, rx, ry + 0.01, rank_w, rank_w,
             fill=RANK_BG[i], line=None)
    add_textbox(slide, str(rank),
                l=rx, t=ry + 0.01, w=rank_w, h=rank_w,
                size=9, color=DEEP_FOREST, bold=True,
                align='center', valign='middle')
    # member name
    add_textbox(slide, name,
                l=rx + rank_w + gap, t=ry, w=name_w, h=row_h,
                size=10, color=INK, bold=(i == 0), valign='middle')
    # cost bar (proportional, sage with amber for #1)
    bw = bar_w_max * (cost / TOP5_MAX_COST)
    bar_color = AMBER if i == 0 else COLOR_OTHER
    add_rect(slide, bar_x, ry + 0.05, bw, row_h - 0.10,
             fill=bar_color, line=None)
    # cost value (right-aligned, number only — unit in panel header)
    add_textbox(slide, f"{cost}",
                l=TOP5_L + TOP5_W - inner_pad - value_w, t=ry,
                w=value_w, h=row_h,
                size=11, color=DEEP_FOREST, bold=True,
                align='right', valign='middle')

# Footnote — mock data disclosure (per user, condition: explicit on-slide notice)
add_textbox(slide,
            "※メンバー名・金額はレイアウト確認用のモックデータ／TOP6以降は別紙Excel参照",
            l=TOP5_L, t=5.36, w=TOP5_W, h=0.14,
            size=6, color=MUTED, align='left')

# =====================================================================
# 5) TEAM MAPPING (REDESIGNED, larger, full-width)  T: 5.65 - 8.50
# =====================================================================
# ---- panel ----
add_rect(slide, 0.30, 5.65, 7.67, 2.85, fill=WHITE, line=CARD_BORDER, line_width=0.5)
# title
add_textbox(slide,
            "全社14チームのポジショニング — 当チームは右上の理想ゾーン、下位3チームに導入支援余地",
            l=0.40, t=5.70, w=7.50, h=0.22,
            size=11, color=DEEP_FOREST, bold=True)
add_textbox(slide,
            "X軸: Claude利用量 (相対) →   /   Y軸: 効率 [採用行数 ÷ コスト] ↑   /   バブル径: 1人あたりコスト",
            l=0.40, t=5.92, w=7.50, h=0.16,
            size=8, color=MUTED)

# ---- plot box ----
PX, PY = 1.10, 6.15      # plot origin (top-left of plot area)
PW, PH = 5.20, 2.05      # plot dimensions (inches)
# axis lines + arrows (forest)
add_line(slide, PX, PY+PH, PX+PW, PY+PH, color=DEEP_FOREST, width=1.0)  # x
add_line(slide, PX, PY,    PX,     PY+PH, color=DEEP_FOREST, width=1.0)  # y
# 4-quadrant grid lines (light)
add_line(slide, PX+PW/2, PY, PX+PW/2, PY+PH, color=CARD_BORDER, width=0.5)
add_line(slide, PX, PY+PH/2, PX+PW, PY+PH/2, color=CARD_BORDER, width=0.5)

# axis labels
add_textbox(slide, "→ Claude利用量 (大)",
            l=PX+PW-1.50, t=PY+PH+0.05, w=1.50, h=0.18,
            size=8, color=DEEP_FOREST, bold=True, align='right')
add_textbox(slide, "(小)",
            l=PX-0.05, t=PY+PH+0.05, w=0.30, h=0.18,
            size=8, color=MUTED, align='left')
# Y axis label (vertical not supported simply; place rotated text via shape)
yl = slide.shapes.add_textbox(Inches(PX-0.95), Inches(PY+PH/2-0.40),
                              Inches(0.90), Inches(0.80))
yl.text_frame.word_wrap = True
yl.text_frame.margin_left = yl.text_frame.margin_right = Emu(0)
yl.text_frame.margin_top = yl.text_frame.margin_bottom = Emu(0)
ylp = yl.text_frame.paragraphs[0]
ylp.alignment = PP_ALIGN.RIGHT
ylr = ylp.add_run()
ylr.text = "↑ 効率\n(行/コスト)\n高"
set_run_font(ylr, size=8, color=DEEP_FOREST, bold=True)

# Quadrant captions (subtle, in the corners of plot area)
quad_specs = [
    # (text, sub, l_offset_frac, t_offset_frac, color)
    # All quadrant captions in muted gray so they don't compete with the focal Amber.
    ("活用優等生",        "高利用 × 高効率",   0.55, 0.04, MUTED),
    ("効率派・規模小",    "低利用 × 高効率",   0.04, 0.04, MUTED),
    ("要支援 (未活用)",   "低利用 × 低効率",   0.04, 0.78, MUTED),
    ("コスト先行・低効率", "高利用 × 低効率",  0.55, 0.78, MUTED),
]
for txt, sub, lf, tf, col in quad_specs:
    qx = PX + PW * lf
    qy = PY + PH * tf
    add_textbox(slide, txt, l=qx, t=qy, w=PW*0.42, h=0.18,
                size=9, color=col, bold=True)
    add_textbox(slide, sub, l=qx, t=qy+0.16, w=PW*0.42, h=0.16,
                size=7, color=MUTED)

# ---- bubble data (positions preserved from V7; team labels are placeholders
# because V7 has no team names; clearly noted in footnote) ----
# Format: (label, x_norm 0..1, y_norm 0..1 (1=top), size_in, is_own)
# All non-own teams use the same uniform sage color so the focal point (当チーム) pops.
TEAMS = [
    ("当チーム",  0.83, 0.95, 0.20, True),
    ("チームB",   0.86, 0.92, 0.20, False),
    ("チームC",   0.90, 0.78, 0.26, False),
    ("チームD",   0.80, 0.81, 0.21, False),
    ("チームE",   0.76, 0.84, 0.18, False),
    ("チームF",   0.70, 0.73, 0.19, False),
    ("チームG",   0.68, 0.66, 0.16, False),
    ("チームH",   0.58, 0.64, 0.18, False),
    ("チームI",   0.53, 0.60, 0.23, False),
    ("チームJ",   0.48, 0.56, 0.24, False),
    ("チームK",   0.42, 0.55, 0.26, False),
    ("チームL",   0.37, 0.48, 0.28, False),
    ("チームM",   0.32, 0.53, 0.31, False),
    ("チームN",   0.22, 0.39, 0.29, False),
]

# Helper to convert normalized (xn, yn_top1) into inches top-left for bubble
def bubble_pos(xn, yn, dia):
    cx = PX + xn * PW
    cy = PY + (1.0 - yn) * PH   # invert: yn=1 means top of plot
    return cx - dia/2, cy - dia/2

# draw bubbles — own team in Amber (focal), all others in uniform Sage.
# V10: per-bubble team-name labels removed entirely. Only the own team keeps
# its "当チーム" caption so its position remains identifiable.
# Draw "other" bubbles first so the own-team bubble sits on top visually.
for name, xn, yn, dia, is_own in sorted(TEAMS, key=lambda t: 1 if t[4] else 0):
    bx, by = bubble_pos(xn, yn, dia)
    fill = COLOR_OWN if is_own else COLOR_OTHER
    outline = AMBER_DEEP if is_own else FOREST
    add_oval(slide, bx, by, dia, dia, fill=fill, line=outline)
    # Label only the own team — other team names are intentionally suppressed.
    if is_own:
        label_x = bx + dia + 0.02
        label_y = by - 0.04
        if label_x + 0.60 > PX + PW + 0.10:
            label_x = bx - 0.60
        add_textbox(slide, "当チーム",
                    l=label_x, t=label_y, w=0.65, h=0.16,
                    size=9, color=AMBER_DEEP, bold=True)

# ---- legend (right side of plot) — V10: 当チーム のみ ----
LX = PX + PW + 0.30
LY = PY + 0.10
add_textbox(slide, "● 当チーム", l=LX, t=LY, w=1.30, h=0.22,
            size=10, color=AMBER_DEEP, bold=True)

# Key insight callout — shifted up since the legend now has 1 line instead of 2
add_rect(slide, LX, LY + 0.40, 1.30, 1.20, fill=SAGE_BG, line=AMBER, line_width=0.75)
add_textbox(slide, "示唆",
            l=LX+0.06, t=LY+0.43, w=1.20, h=0.18,
            size=9, color=AMBER_DEEP, bold=True)
add_textbox(slide,
            "当チームは右上に位置し\n高利用×高効率を実現。\n\n横展開支援で\n左下チームの底上げ\n余地あり。",
            l=LX+0.06, t=LY+0.64, w=1.20, h=0.94,
            size=7, color=INK, line_spacing=1.20)

# footnote — V10: team-name caveat dropped (labels are gone). Keep only the
# position-source note.
add_textbox(slide,
            "※ バブル位置はV7と一致する相対値。各バブルの個別データはダッシュボード参照。",
            l=0.40, t=8.32, w=7.50, h=0.16,
            size=7, color=MUTED)

# =====================================================================
# 6) TWO TOP3 PANELS  T: 8.55 - 9.85  (V9: コストTOP3 削除 → 生成行 + 効率 のみ)
# =====================================================================
# Header row
add_textbox(slide,
            "メンバー別 TOP3 — 当チーム内 (V7記載の3名を2視点で再ランキング)",
            l=0.30, t=8.55, w=7.67, h=0.20,
            size=12, color=DEEP_FOREST, bold=True)

# Two wider columns
panel_w = 3.78
panel_gap = 0.10
panel_left = 0.30

# Data: 3 named members corresponding to V7's ユーザーA/B/C
# Mapping (per user instruction):
#   ユーザーA → 山田 太郎  (130 [千円], 4,200行, 32.3 行/千円)
#   ユーザーB → 鈴木 花子  (108 [千円], 4,500行, 41.7 行/千円)
#   ユーザーC → 佐藤 健    ( 92 [千円], 2,800行, 30.4 行/千円)
USERS = {
    "山田 太郎": dict(cost_k=130, loc=4200),
    "鈴木 花子": dict(cost_k=108, loc=4500),
    "佐藤 健":   dict(cost_k=92,  loc=2800),
}
for u, d in USERS.items():
    d["eff"] = round(d["loc"] / d["cost_k"], 1)  # 行 / 千円

# rankings (descending) — only LOC and efficiency in V9
top_loc  = sorted(USERS.items(), key=lambda kv: -kv[1]["loc"])
top_eff  = sorted(USERS.items(), key=lambda kv: -kv[1]["eff"])

panels = [
    ("生成行TOP3",  "最もコード生成 (行)",   top_loc,
        lambda d: f"{d['loc']:,}行",
        lambda d: f"{d['cost_k']} [千円] / {d['eff']} 行/千円"),
    ("効率TOP3",    "コスパ (行/千円)",      top_eff,
        lambda d: f"{d['eff']}",
        lambda d: f"{d['cost_k']} [千円] → {d['loc']:,}行"),
]

panel_top = 8.78
panel_h = 1.07
RANK_BG = [AMBER, SAGE, RGBColor(0xE0, 0xE0, 0xE0)]  # 1st, 2nd, 3rd
for pi, (head, sub, ranking, big_fn, tail_fn) in enumerate(panels):
    px = panel_left + pi * (panel_w + panel_gap)
    add_rect(slide, px, panel_top, panel_w, panel_h,
             fill=WHITE, line=CARD_BORDER, line_width=0.5)
    # header strip
    add_rect(slide, px, panel_top, panel_w, 0.04, fill=DEEP_FOREST)
    add_textbox(slide, head, l=px+0.08, t=panel_top+0.06, w=panel_w-0.16, h=0.18,
                size=11, color=DEEP_FOREST, bold=True)
    add_textbox(slide, sub, l=px+0.08, t=panel_top+0.24, w=panel_w-0.16, h=0.14,
                size=8, color=MUTED)
    # rows
    row_top = panel_top + 0.40
    row_h = 0.22
    for ri, (uname, d) in enumerate(ranking):
        ry = row_top + ri * row_h
        # rank badge
        add_oval(slide, px+0.08, ry+0.02, 0.18, 0.18,
                 fill=RANK_BG[ri], line=None)
        add_textbox(slide, str(ri+1),
                    l=px+0.08, t=ry+0.02, w=0.18, h=0.18,
                    size=10, color=DEEP_FOREST, bold=True,
                    align='center', valign='middle')
        # name (already a real Japanese name in V9)
        add_textbox(slide, uname,
                    l=px+0.30, t=ry+0.01, w=1.40, h=0.18,
                    size=10, color=INK, bold=True, valign='middle')
        # big metric (right)
        add_textbox(slide, big_fn(d),
                    l=px+1.75, t=ry+0.01, w=panel_w-1.80, h=0.18,
                    size=12, color=AMBER_DEEP, bold=True,
                    align='right', valign='middle')
        # tail under
        # (omitted to keep panel clean — replaced with single summary line)
    # bottom mini-summary
    add_textbox(slide, tail_fn(ranking[0][1]) if ranking else "",
                l=px+0.08, t=panel_top+panel_h-0.18, w=panel_w-0.16, h=0.16,
                size=7, color=MUTED, align='right')

# =====================================================================
# 7) 当チーム実績 + 公開研究  T: 9.95 - 11.30
# =====================================================================
# left: 当チーム実績
LX1, LY1 = 0.30, 9.95
LW1, LH1 = 3.65, 1.35
add_rect(slide, LX1, LY1, LW1, LH1, fill=DEEP_FOREST, line=None)
add_textbox(slide, "当チーム実績 (3月)",
            l=LX1+0.10, t=LY1+0.06, w=LW1-0.20, h=0.18,
            size=11, color=SAGE, bold=True)
add_textbox(slide, "822 [千円] → 26,000行",
            l=LX1+0.10, t=LY1+0.24, w=LW1-0.20, h=0.30,
            size=18, color=AMBER, bold=True)
add_line(slide, LX1+0.10, LY1+0.55, LX1+LW1-0.10, LY1+0.55, color=AMBER, width=0.5)
add_textbox(slide, "1,500 [千円] 投資換算 (人月単価相当)",
            l=LX1+0.10, t=LY1+0.58, w=LW1-0.20, h=0.18,
            size=10, color=WHITE, bold=True)
add_textbox(slide, "約 47,400 行 / 人月",
            l=LX1+0.10, t=LY1+0.78, w=LW1-0.20, h=0.30,
            size=20, color=AMBER, bold=True)
add_textbox(slide, "業界中央値 600行/人月の約79倍 (※単純比較)",
            l=LX1+0.10, t=LY1+1.10, w=LW1-0.20, h=0.18,
            size=8, color=SAGE)

# right: 公開研究
RX1, RY1 = 4.10, 9.95
RW1, RH1 = 3.87, 1.35
add_rect(slide, RX1, RY1, RW1, RH1, fill=WHITE, line=CARD_BORDER, line_width=0.5)
add_textbox(slide, "【参考】 公開研究の月間コード生産量",
            l=RX1+0.10, t=RY1+0.06, w=RW1-0.20, h=0.18,
            size=11, color=DEEP_FOREST, bold=True)
add_textbox(slide, "※ 品質・複雑度・業務特性は考慮されない単純比較",
            l=RX1+0.10, t=RY1+0.24, w=RW1-0.20, h=0.14,
            size=8, color=MUTED)

bench = [
    ("保守的",   "325 〜 750 行/月",   "Capers Jones"),
    ("中央値的", "約 1,000 行/月",    "日本SI業界 (IPA)"),
    ("生産性高", "1,500 〜 1,600 行/月","個人実測 / Scrum"),
]
for i, (lvl, num, src) in enumerate(bench):
    by = RY1 + 0.43 + i * 0.30
    add_rect(slide, RX1+0.10, by, RW1-0.20, 0.27, fill=SAGE_BG,
             line=CARD_BORDER, line_width=0.3)
    add_rect(slide, RX1+0.10, by, 0.05, 0.27, fill=FOREST)
    add_textbox(slide, lvl, l=RX1+0.20, t=by+0.04, w=0.70, h=0.18,
                size=9, color=DEEP_FOREST, bold=True, valign='middle')
    add_textbox(slide, num, l=RX1+0.95, t=by+0.04, w=1.55, h=0.18,
                size=10, color=INK, bold=True, valign='middle')
    add_textbox(slide, src, l=RX1+2.50, t=by+0.04, w=RW1-2.65, h=0.18,
                size=8, color=MUTED, valign='middle', align='right')

# =====================================================================
# 8) BOTTOM CONCLUSION BAR  T: 11.30 - 11.55
# =====================================================================
add_rect(slide, 0.00, 11.30, 8.27, 0.39, fill=DEEP_FOREST)
# left amber accent
add_rect(slide, 0.00, 11.30, 0.06, 0.39, fill=AMBER)
add_textbox(slide,
            "効率は改善継続、年度末予測 9,900 [千円] = 外注費用 6.6人月相当。次手は『外注増員 vs Claude活用拡大』の経営判断",
            l=0.20, t=11.34, w=8.00, h=0.32,
            size=11, color=WHITE, align='left', valign='middle')

# footer (very thin, below)
# (page is full; place footnote inside the conclusion bar by omission)

# =====================================================================
# Speaker notes (sources)
# =====================================================================
notes = slide.notes_slide.notes_text_frame
notes.text = (
    "V10 — based on V9.\n"
    "V10 changes (per user 2026-05-11 feedback):\n"
    "  - TEAM MAP: per-bubble team-name labels removed. Only 当チーム is\n"
    "    captioned so its position remains identifiable. Legend reduced to\n"
    "    a single '当チーム' item (他チーム line removed). Quadrant captions,\n"
    "    axes, 示唆 callout, and 2-color (Amber/Sage) split preserved.\n"
    "  - WORDING: every 'ミドル' replaced with context-appropriate '外注':\n"
    "      'ミドル6.6人月'   → '外注費用 6.6人月'\n"
    "      'ミドル人材'      → '外注費用'\n"
    "      'ミドル増員'      → '外注増員'\n"
    "  - Version V9 → V10, date 2026-05-09 → 2026-05-11.\n"
    "\n"
    "--- prior V9 notes follow ---\n"
    "V9 — based on V8 rev3.\n"
    "V9 changes (per user 2026-05-09 feedback):\n"
    "  - UNIT: all currency normalized to [千円]. k/M notation removed.\n"
    "      ¥9.9M→9,900; ¥822k→822; ¥1.5M→1,500; ¥68.5k→68.5; etc.\n"
    "      効率TOP3 unit: 行/¥1k → 行/千円.\n"
    "  - LAYOUT (monthly trend band): 100% width chart → split into LEFT 58%\n"
    "      dual-axis trend chart + RIGHT 42% '自プロジェクトメンバー コストTOP5'\n"
    "      panel (NEW). Member-row TOP3 reduced 3→2 (生成行/効率 only; コスト\n"
    "      panel removed because TOP5 cost is now shown above).\n"
    "  - VERSION: V8 → V9; date 2026-05-08 → 2026-05-09.\n"
    "\n"
    "Mock-data disclosure (CRITICAL):\n"
    "  - TOP1-3 cost values (130/108/92 [千円]) preserved 1:1 from V7.\n"
    "  - V7 codenames remapped to real names per user instruction:\n"
    "      ユーザーA → 山田 太郎 (130, 4,200行, 32.3行/千円)\n"
    "      ユーザーB → 鈴木 花子 (108, 4,500行, 41.7行/千円)\n"
    "      ユーザーC → 佐藤 健    ( 92, 2,800行, 30.4行/千円)\n"
    "  - TOP4 (田中 美咲 78 [千円]) and TOP5 (高橋 翔 64 [千円]) are MOCK values\n"
    "    supplied by the user for layout-confirmation only. Disclosed on-slide\n"
    "    via footnote in the TOP5 panel.\n"
    "\n"
    "Data preserved 1:1 from V8 rev3: 月次推移6点, KPI 4枚, 9,900予測,\n"
    "  822→26,000行, 公開研究3水準, 14チーム配置 (2色, 4象限)."
)

out = "/sessions/lucid-festive-wozniak/mnt/_outputs/2026-05-11_Claude_monthly_report_V10.pptx"
prs.save(out)
print("WROTE:", out)
